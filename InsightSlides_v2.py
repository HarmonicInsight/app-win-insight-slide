# -*- coding: utf-8 -*-
"""
Insight Slides v2.0 - PowerPoint Text Extract & Update Tool
çµ±åˆç‰ˆ: æ—§UI + ã‚°ãƒªãƒƒãƒ‰ç·¨é›† + æ¯”è¼ƒæ©Ÿèƒ½ + ãƒ•ã‚£ãƒ«ã‚¿

by Harmonic Insight

ç‰¹å¾´:
- æŠ½å‡º/æ›´æ–°ãƒ¢ãƒ¼ãƒ‰åˆ‡æ›¿
- ã‚¤ãƒ³ãƒ©ã‚¤ãƒ³ã‚°ãƒªãƒƒãƒ‰ç·¨é›†
- PPTXæ¯”è¼ƒæ©Ÿèƒ½
- ãƒ•ã‚£ãƒ«ã‚¿æ©Ÿèƒ½
- çµ±ä¸€ãƒ©ã‚¤ã‚»ãƒ³ã‚¹å½¢å¼ (INS-SLIDE-{TIER}-XXXX-XXXX-CC)
- æŠ˜ã‚ŠãŸãŸã¿å¯èƒ½ãªã‚ªãƒ—ã‚·ãƒ§ãƒ³
"""
import tkinter as tk
from tkinter import ttk, filedialog, messagebox, scrolledtext, simpledialog
import pptx
import openpyxl
from openpyxl.styles import Font as XLFont, PatternFill
import os
import sys
import re
import csv
import json
import hashlib
import random
import string
import webbrowser
import traceback
from datetime import datetime, timedelta

# insight-common ãƒ©ã‚¤ã‚»ãƒ³ã‚¹ãƒ¢ã‚¸ãƒ¥ãƒ¼ãƒ«ã‚’ã‚¤ãƒ³ãƒãƒ¼ãƒˆ
import importlib.util
_license_module_path = os.path.join(os.path.dirname(__file__), 'insight-common', 'license', 'python', '__init__.py')
_spec = importlib.util.spec_from_file_location("insight_license", _license_module_path)
_license_module = importlib.util.module_from_spec(_spec)
_spec.loader.exec_module(_license_module)
ProductCode = _license_module.ProductCode
InsightLicenseTier = _license_module.LicenseTier
LicenseInfo = _license_module.LicenseInfo
LicenseValidator = _license_module.LicenseValidator
get_feature_limits = _license_module.get_feature_limits
TIER_NAMES = _license_module.TIER_NAMES
INSIGHT_TIERS = _license_module.TIERS
import threading
from pathlib import Path
from typing import Dict, Tuple, List, Optional
import shutil

# ============== App Info ==============
APP_VERSION = "2.0.0"
APP_NAME = "Insight Slides"

# ============== Config Paths ==============
CONFIG_DIR = Path.home() / ".insightslides"
CONFIG_FILE = CONFIG_DIR / "config.json"
LICENSE_FILE = CONFIG_DIR / "license.key"
ERROR_LOG_FILE = CONFIG_DIR / "error_log.txt"

# ============== Support Links ==============
SUPPORT_LINKS = {
    "faq": "https://example.com/insightslides/faq",
    "tutorial": "https://example.com/insightslides/tutorial",
    "purchase": "https://example.com/insightslides/purchase",
    "contact": "mailto:support@example.com",
}

# ============== Internationalization (i18n) ==============
LANGUAGES = {
    'en': {
        'app_subtitle': 'Extract â†’ Edit â†’ Update PowerPoint Text',
        'welcome_title': 'Welcome to Insight Slides!',
        'mode_extract': 'Extract Mode',
        'mode_update': 'Update Mode',
        'mode_extract_short': 'Extract Text',
        'mode_update_short': 'Overwrite',
        'panel_mode': 'Mode Selection',
        'panel_file': 'File Operations',
        'panel_input': 'Input (Single File)',
        'panel_output_file': 'Output (Single File)',
        'panel_settings': 'Settings',
        'panel_status': 'Status',
        'panel_output': 'Extracted Data',
        'panel_extract_options': 'Extract Options',
        'panel_update_options': 'Update Options',
        'panel_extract_run': 'Run Extract',
        'panel_update_run': 'Run Update',
        'panel_pro_features': 'Pro Features',
        'btn_load_pptx': 'Load PPTX',
        'btn_load_excel': 'Load Excel',
        'btn_load_json': 'Load JSON',
        'btn_single_file': 'Select File',
        'btn_from_excel': 'From Excel',
        'btn_from_json': 'From JSON',
        'btn_apply_pptx': 'Apply to PPTX',
        'btn_export_to_excel': 'Export to Excel',
        'btn_export_to_json': 'Export to JSON',
        'panel_batch': 'Folder Batch',
        'btn_batch_extract': 'Batch Extract',
        'btn_batch_update': 'Batch Update',
        'btn_batch_export_excel': 'Export to Folder (Excel)',
        'btn_batch_export_json': 'Export to Folder (JSON)',
        'btn_batch_import_excel': 'Import from Folder (Excel)',
        'btn_batch_import_json': 'Import from Folder (JSON)',
        'btn_diff_preview': 'Diff Preview',
        'btn_compare_pptx': 'Compare PPTX',
        'btn_cancel': 'Stop',
        'btn_clear': 'Clear Log',
        'btn_copy': 'Copy Log',
        'btn_license': 'License',
        'btn_activate': 'Activate',
        'btn_deactivate': 'Deactivate',
        'btn_purchase': 'Purchase',
        'btn_close': 'Close',
        'btn_start': 'Get Started',
        'btn_filter': 'Filter',
        'btn_clear_filter': 'Clear Filter',
        'setting_output_format': 'Output Format:',
        'setting_include_meta': 'Include file name & date',
        'setting_auto_backup': 'Auto backup before update',
        'chk_include_notes': 'Include Speaker Notes',
        'format_tab': 'Tab-separated',
        'format_csv': 'CSV',
        'format_excel': 'Excel',
        'status_waiting': 'Waiting...',
        'status_processing': 'Processing...',
        'status_complete': 'Complete',
        'status_cancelled': 'Cancelled',
        'status_error': 'Error',
        'msg_extract_desc': 'Extract text from PowerPoint files.',
        'msg_update_desc': 'Apply edited text back to PowerPoint.',
        'msg_update_limit': 'Update: First {0} slides only\nUpgrade to Standard for unlimited!',
        'msg_processing_file': 'Processing: {0}',
        'msg_saved': 'Saved: {0}',
        'msg_extracted': 'Extracted: {0} items from {1} slides',
        'msg_updated': 'Updated: {0} items, Skipped: {1}',
        'msg_no_pptx': 'No PPTX files found',
        'msg_no_data': 'No update data found',
        'msg_copied': 'Copied to clipboard',
        'license_title': 'License Management',
        'license_current': 'Current License',
        'license_enter_key': 'Enter License Key:',
        'license_activated': '{0} has been activated',
        'license_deactivated': 'License deactivated',
        'license_invalid': 'Invalid license key',
        'license_enter_prompt': 'Please enter a license key',
        'upgrade_title': 'Upgrade',
        'dialog_confirm': 'Confirm',
        'dialog_error': 'Error',
        'dialog_complete': 'Complete',
        'header_slide': 'Slide',
        'header_id': 'Object ID',
        'header_type': 'Type',
        'header_text': 'Text Content',
        'header_filename': 'Filename',
        'header_datetime': 'Extracted At',
        'diff_title': 'Diff Preview',
        'menu_help': 'Help',
        'menu_guide': 'User Guide',
        'menu_faq': 'FAQ',
        'menu_license': 'License Management',
        'menu_about': 'About',
        'lang_menu': 'Language',
        'font_size_menu': 'Font Size',
        'font_size_small': 'Small',
        'font_size_medium': 'Medium',
        'font_size_large': 'Large',
        'advanced_options': 'Advanced Options',
        'type_notes': 'Notes',
        'filter_placeholder': 'Filter text...',
        # UI elements
        'mode_section': 'Mode',
        'btn_compare': '2-File Compare',
        'show_detail': 'Show details',
        'welcome_guide_title': 'Edit PowerPoint Text',
        'guide_step1': 'Select a PPTX file from the left panel',
        'guide_step2': 'Text will be displayed in a list',
        'guide_step3': 'Double-click a cell to edit',
        'guide_step4': 'Click "Apply" to update PPTX',
        'btn_apply': 'Apply',
        'btn_export_excel': 'Excel Export',
        'btn_export_json': 'JSON Export',
        'filter_label': 'Filter:',
        'mode_desc_extract': 'Extract text from PPTX for editing',
        'mode_desc_update': 'Apply edited data to PPTX',
        # Grid toolbar
        'btn_clear_grid': 'Clear',
        'btn_replace_all': 'Replace All',
        'btn_undo': 'Undo',
        'btn_redo': 'Redo',
        # Replace dialog
        'replace_search': 'Search:',
        'replace_with': 'Replace:',
        'btn_replace': 'Replace',
        # Compare dialog
        'compare_title': 'Compare 2 PowerPoint files',
        'compare_file1': 'Original:',
        'compare_file2': 'New file:',
        'btn_browse': 'Browse',
        'compare_ignore_ws': 'Ignore whitespace',
        'btn_run_compare': 'Compare',
        # Compare result
        'btn_export_csv': 'CSV Export',
        'header_select': 'Select',
        'header_status': 'Status',
        'btn_select_original': 'All Original',
        'btn_select_new': 'All New',
        'btn_apply_selection': 'Apply Selection',
        # Log dialog
        'btn_copy_log': 'Copy',
        'btn_clear_log': 'Clear',
        # License dialog (auth)
        'license_auth_title': 'License Activation',
        'license_email': 'Email Address:',
        'license_key': 'License Key:',
        'license_wrong_product': 'This license key is not valid for Insight Slides',
        'license_perpetual': 'Perpetual',
        'license_expiry_warning': 'Your license will expire in {0} days ({1}). Please renew.',
        'license_expired': 'Your license has expired. Please renew to continue using all features.',
        'license_trial_link': 'Request Trial',
        'license_email_required': 'Please enter your email address',
        'license_status_active': 'Active',
        'license_status_expired': 'Expired',
        'license_valid_until': 'Valid until: {0}',
        'license_days_remaining': '({0} days remaining)',
        'license_feature_restricted': 'This feature requires a Pro license. Current: {0}',
        'license_batch_restricted': 'Batch processing requires a Pro license.',
        'license_json_restricted': 'JSON export requires a Pro license.',
        'license_continue_free': 'Continue as Free',
        # Status messages
        'status_slides_items': '{0} slides / {1} items',
        'status_complete_items': 'Complete: {0} items',
        'status_batch_complete': 'Batch extract complete: {0} items ({1})',
        'status_update_complete': 'Update complete: {0} items',
        'lang_changed': 'Language changed.',
        # Log messages
        'log_cancelled': 'Cancelled',
        'log_cancel_request': 'Cancellation requested...',
        'log_no_text': 'No text found',
        'log_error': 'Error: {0}',
        'log_found_files': 'Found: {0} files',
        'log_no_pptx_found': 'No PPTX files found',
        'log_invalid_header': 'Invalid header format',
        'log_no_update_data': 'No update data',
        'log_processing': 'Processing...',
        'dialog_select_folder': 'Select Folder (containing PPTX files)',
        'dialog_select_folder_update': 'Select Folder (*_extracted{0} + PPTX)',
        'dialog_select_pptx': 'Select PowerPoint to update',
        'dialog_processing_exit': 'Processing in progress. Exit anyway?',
        'dialog_confirm_title': 'Confirm',
        'result_updated': 'Updated: {0} items\nSkipped: {1} items',
        'result_replaced': '{0} items replaced',
        'result_applied': '{0} items applied',
        'result_csv_saved': 'CSV saved',
        'result_export_complete': 'Export complete: {0}',
    },
    'ja': {
        'app_subtitle': 'PowerPointãƒ†ã‚­ã‚¹ãƒˆã‚’æŠ½å‡º â†’ ç·¨é›† â†’ åæ˜ ',
        'welcome_title': 'Insight Slides ã¸ã‚ˆã†ã“ãï¼',
        'mode_extract': 'æŠ½å‡ºãƒ¢ãƒ¼ãƒ‰',
        'mode_update': 'æ›´æ–°ãƒ¢ãƒ¼ãƒ‰',
        'mode_extract_short': 'ãƒ†ã‚­ã‚¹ãƒˆæŠ½å‡º',
        'mode_update_short': 'ä¸Šæ›¸ãæ›´æ–°',
        'panel_mode': 'ãƒ¢ãƒ¼ãƒ‰é¸æŠ',
        'panel_file': 'ãƒ•ã‚¡ã‚¤ãƒ«æ“ä½œ',
        'panel_input': 'å…¥åŠ›ï¼ˆ1ãƒ•ã‚¡ã‚¤ãƒ«ï¼‰',
        'panel_output_file': 'å‡ºåŠ›ï¼ˆ1ãƒ•ã‚¡ã‚¤ãƒ«ï¼‰',
        'panel_settings': 'å‡¦ç†è¨­å®š',
        'panel_status': 'å‡¦ç†çŠ¶æ³',
        'panel_output': 'æŠ½å‡ºçµæœ',
        'panel_extract_options': 'æŠ½å‡ºã‚ªãƒ—ã‚·ãƒ§ãƒ³',
        'panel_update_options': 'æ›´æ–°ã‚ªãƒ—ã‚·ãƒ§ãƒ³',
        'panel_extract_run': 'æŠ½å‡ºå®Ÿè¡Œ',
        'panel_update_run': 'æ›´æ–°å®Ÿè¡Œ',
        'panel_pro_features': 'æ‹¡å¼µæ©Ÿèƒ½',
        'btn_load_pptx': 'PPTXèª­è¾¼',
        'btn_load_excel': 'Excelèª­è¾¼',
        'btn_load_json': 'JSONèª­è¾¼',
        'btn_single_file': 'ãƒ•ã‚¡ã‚¤ãƒ«é¸æŠ',
        'btn_from_excel': 'Excelã‹ã‚‰æ›´æ–°',
        'btn_from_json': 'JSONã‹ã‚‰æ›´æ–°',
        'btn_apply_pptx': 'PPTXã«åæ˜ ',
        'btn_export_to_excel': 'Excelå‡ºåŠ›',
        'btn_export_to_json': 'JSONå‡ºåŠ›',
        'panel_batch': 'ãƒ•ã‚©ãƒ«ãƒ€ä¸€æ‹¬',
        'btn_batch_extract': 'ä¸€æ‹¬æŠ½å‡º',
        'btn_batch_update': 'ä¸€æ‹¬æ›´æ–°',
        'btn_batch_export_excel': 'ãƒ•ã‚©ãƒ«ãƒ€ã«å‡ºåŠ› (Excel)',
        'btn_batch_export_json': 'ãƒ•ã‚©ãƒ«ãƒ€ã«å‡ºåŠ› (JSON)',
        'btn_batch_import_excel': 'ãƒ•ã‚©ãƒ«ãƒ€ã‹ã‚‰èª­è¾¼ (Excel)',
        'btn_batch_import_json': 'ãƒ•ã‚©ãƒ«ãƒ€ã‹ã‚‰èª­è¾¼ (JSON)',
        'btn_diff_preview': 'å·®åˆ†ãƒ—ãƒ¬ãƒ“ãƒ¥ãƒ¼',
        'btn_compare_pptx': 'PPTXæ¯”è¼ƒ',
        'btn_cancel': 'ä¸­æ­¢',
        'btn_clear': 'ãƒ­ã‚°ã‚¯ãƒªã‚¢',
        'btn_copy': 'ãƒ­ã‚°ã‚³ãƒ”ãƒ¼',
        'btn_license': 'ãƒ©ã‚¤ã‚»ãƒ³ã‚¹',
        'btn_activate': 'ã‚¢ã‚¯ãƒ†ã‚£ãƒ™ãƒ¼ãƒˆ',
        'btn_deactivate': 'ãƒ©ã‚¤ã‚»ãƒ³ã‚¹è§£é™¤',
        'btn_purchase': 'è³¼å…¥ãƒšãƒ¼ã‚¸',
        'btn_close': 'é–‰ã˜ã‚‹',
        'btn_start': 'å§‹ã‚ã‚‹',
        'btn_filter': 'ãƒ•ã‚£ãƒ«ã‚¿',
        'btn_clear_filter': 'ã‚¯ãƒªã‚¢',
        'setting_output_format': 'å‡ºåŠ›å½¢å¼:',
        'setting_include_meta': 'ãƒ•ã‚¡ã‚¤ãƒ«åãƒ»æ—¥æ™‚ã‚’å«ã‚ã‚‹',
        'setting_auto_backup': 'æ›´æ–°å‰ã«è‡ªå‹•ãƒãƒƒã‚¯ã‚¢ãƒƒãƒ—',
        'chk_include_notes': 'ã‚¹ãƒ”ãƒ¼ã‚«ãƒ¼ãƒãƒ¼ãƒˆå«ã‚€',
        'format_tab': 'ã‚¿ãƒ–åŒºåˆ‡ã‚Š',
        'format_csv': 'CSVå½¢å¼',
        'format_excel': 'Excelå½¢å¼',
        'status_waiting': 'å‡¦ç†å¾…æ©Ÿä¸­...',
        'status_processing': 'å‡¦ç†ä¸­...',
        'status_complete': 'å®Œäº†',
        'status_cancelled': 'ã‚­ãƒ£ãƒ³ã‚»ãƒ«ã•ã‚Œã¾ã—ãŸ',
        'status_error': 'ã‚¨ãƒ©ãƒ¼',
        'msg_extract_desc': 'PowerPointã‹ã‚‰ãƒ†ã‚­ã‚¹ãƒˆã‚’æŠ½å‡ºã—ã¾ã™ã€‚',
        'msg_update_desc': 'ç·¨é›†ã—ãŸãƒ•ã‚¡ã‚¤ãƒ«ã®å¤‰æ›´ã‚’PowerPointã«åæ˜ ã—ã¾ã™ã€‚',
        'msg_update_limit': 'æ›´æ–°æ©Ÿèƒ½: æœ€åˆã®{0}ã‚¹ãƒ©ã‚¤ãƒ‰ã®ã¿\nStandardç‰ˆã§ç„¡åˆ¶é™ã«ï¼',
        'msg_processing_file': 'å‡¦ç†ä¸­: {0}',
        'msg_saved': 'ä¿å­˜å®Œäº†: {0}',
        'msg_extracted': 'æŠ½å‡º: {0}ä»¶ / ã‚¹ãƒ©ã‚¤ãƒ‰: {1}æš',
        'msg_updated': 'æ›´æ–°: {0}ä»¶ / ã‚¹ã‚­ãƒƒãƒ—: {1}ä»¶',
        'msg_no_pptx': 'PPTXãƒ•ã‚¡ã‚¤ãƒ«ãŒè¦‹ã¤ã‹ã‚Šã¾ã›ã‚“',
        'msg_no_data': 'æ›´æ–°ãƒ‡ãƒ¼ã‚¿ãŒã‚ã‚Šã¾ã›ã‚“',
        'msg_copied': 'ã‚¯ãƒªãƒƒãƒ—ãƒœãƒ¼ãƒ‰ã«ã‚³ãƒ”ãƒ¼ã—ã¾ã—ãŸ',
        'license_title': 'ãƒ©ã‚¤ã‚»ãƒ³ã‚¹ç®¡ç†',
        'license_current': 'ç¾åœ¨ã®ãƒ©ã‚¤ã‚»ãƒ³ã‚¹',
        'license_enter_key': 'ãƒ©ã‚¤ã‚»ãƒ³ã‚¹ã‚­ãƒ¼:',
        'license_activated': '{0}ç‰ˆãŒã‚¢ã‚¯ãƒ†ã‚£ãƒ™ãƒ¼ãƒˆã•ã‚Œã¾ã—ãŸ',
        'license_deactivated': 'ãƒ©ã‚¤ã‚»ãƒ³ã‚¹ã‚’è§£é™¤ã—ã¾ã—ãŸ',
        'license_invalid': 'ç„¡åŠ¹ãªãƒ©ã‚¤ã‚»ãƒ³ã‚¹ã‚­ãƒ¼ã§ã™',
        'license_enter_prompt': 'ãƒ©ã‚¤ã‚»ãƒ³ã‚¹ã‚­ãƒ¼ã‚’å…¥åŠ›ã—ã¦ãã ã•ã„',
        'upgrade_title': 'ã‚¢ãƒƒãƒ—ã‚°ãƒ¬ãƒ¼ãƒ‰',
        'dialog_confirm': 'ç¢ºèª',
        'dialog_error': 'ã‚¨ãƒ©ãƒ¼',
        'dialog_complete': 'å®Œäº†',
        'header_slide': 'ã‚¹ãƒ©ã‚¤ãƒ‰ç•ªå·',
        'header_id': 'ã‚ªãƒ–ã‚¸ã‚§ã‚¯ãƒˆID',
        'header_type': 'ã‚¿ã‚¤ãƒ—',
        'header_text': 'ãƒ†ã‚­ã‚¹ãƒˆå†…å®¹',
        'header_filename': 'ãƒ•ã‚¡ã‚¤ãƒ«å',
        'header_datetime': 'æŠ½å‡ºæ—¥æ™‚',
        'diff_title': 'å·®åˆ†ãƒ—ãƒ¬ãƒ“ãƒ¥ãƒ¼',
        'menu_help': 'ãƒ˜ãƒ«ãƒ—',
        'menu_guide': 'ä½¿ã„æ–¹ã‚¬ã‚¤ãƒ‰',
        'menu_faq': 'ã‚ˆãã‚ã‚‹è³ªå•',
        'menu_license': 'ãƒ©ã‚¤ã‚»ãƒ³ã‚¹ç®¡ç†',
        'menu_about': 'ãƒãƒ¼ã‚¸ãƒ§ãƒ³æƒ…å ±',
        'lang_menu': 'è¨€èª / Language',
        'font_size_menu': 'æ–‡å­—ã‚µã‚¤ã‚º',
        'font_size_small': 'å°',
        'font_size_medium': 'ä¸­',
        'font_size_large': 'å¤§',
        'advanced_options': 'è©³ç´°ã‚ªãƒ—ã‚·ãƒ§ãƒ³',
        'type_notes': 'ãƒãƒ¼ãƒˆ',
        'filter_placeholder': 'ãƒ•ã‚£ãƒ«ã‚¿...',
        # UI elements
        'mode_section': 'æ“ä½œãƒ¢ãƒ¼ãƒ‰',
        'btn_compare': '2ãƒ•ã‚¡ã‚¤ãƒ«æ¯”è¼ƒ',
        'show_detail': 'è©³ç´°ã‚’è¡¨ç¤º',
        'welcome_guide_title': 'PowerPointãƒ†ã‚­ã‚¹ãƒˆã‚’ç·¨é›†',
        'guide_step1': 'å·¦ã®ãƒ‘ãƒãƒ«ã§PPTXãƒ•ã‚¡ã‚¤ãƒ«ã‚’é¸æŠ',
        'guide_step2': 'ãƒ†ã‚­ã‚¹ãƒˆãŒä¸€è¦§ã§è¡¨ç¤ºã•ã‚Œã¾ã™',
        'guide_step3': 'ã‚»ãƒ«ã‚’ãƒ€ãƒ–ãƒ«ã‚¯ãƒªãƒƒã‚¯ã—ã¦ç·¨é›†',
        'guide_step4': 'ã€Œæ›´æ–°ã‚’é©ç”¨ã€ã§PPTXã«åæ˜ ',
        'btn_apply': 'æ›´æ–°ã‚’é©ç”¨',
        'btn_export_excel': 'Excelã‚¨ã‚¯ã‚¹ãƒãƒ¼ãƒˆ',
        'btn_export_json': 'JSONã‚¨ã‚¯ã‚¹ãƒãƒ¼ãƒˆ',
        'filter_label': 'ãƒ•ã‚£ãƒ«ã‚¿:',
        'mode_desc_extract': 'PPTXã‹ã‚‰ãƒ†ã‚­ã‚¹ãƒˆã‚’æŠ½å‡ºã—ã¦ç·¨é›†',
        'mode_desc_update': 'ç·¨é›†ã—ãŸãƒ‡ãƒ¼ã‚¿ã‚’PPTXã«åæ˜ ',
        # Grid toolbar
        'btn_clear_grid': 'ã‚¯ãƒªã‚¢',
        'btn_replace_all': 'ä¸€æ‹¬ç½®æ›',
        'btn_undo': 'å…ƒã«æˆ»ã™',
        'btn_redo': 'ã‚„ã‚Šç›´ã—',
        # Replace dialog
        'replace_search': 'æ¤œç´¢:',
        'replace_with': 'ç½®æ›:',
        'btn_replace': 'ç½®æ›',
        # Compare dialog
        'compare_title': '2ã¤ã®PowerPointãƒ•ã‚¡ã‚¤ãƒ«ã‚’æ¯”è¼ƒ',
        'compare_file1': 'å…ƒãƒ•ã‚¡ã‚¤ãƒ«:',
        'compare_file2': 'æ–°ãƒ•ã‚¡ã‚¤ãƒ«:',
        'btn_browse': 'å‚ç…§',
        'compare_ignore_ws': 'ç©ºç™½ã®é•ã„ã‚’ç„¡è¦–',
        'btn_run_compare': 'æ¯”è¼ƒå®Ÿè¡Œ',
        # Compare result
        'btn_export_csv': 'CSVã‚¨ã‚¯ã‚¹ãƒãƒ¼ãƒˆ',
        'header_select': 'æ¡ç”¨',
        'header_status': 'çŠ¶æ…‹',
        'btn_select_original': 'å…¨ã¦å…ƒ',
        'btn_select_new': 'å…¨ã¦æ–°',
        'btn_apply_selection': 'é¸æŠã‚’åæ˜ ',
        # Log dialog
        'btn_copy_log': 'ã‚³ãƒ”ãƒ¼',
        'btn_clear_log': 'ã‚¯ãƒªã‚¢',
        # License dialog (auth)
        'license_auth_title': 'ãƒ©ã‚¤ã‚»ãƒ³ã‚¹èªè¨¼',
        'license_email': 'ãƒ¡ãƒ¼ãƒ«ã‚¢ãƒ‰ãƒ¬ã‚¹:',
        'license_key': 'ãƒ©ã‚¤ã‚»ãƒ³ã‚¹ã‚­ãƒ¼:',
        'license_wrong_product': 'ã“ã®ãƒ©ã‚¤ã‚»ãƒ³ã‚¹ã‚­ãƒ¼ã¯Insight Slidesã«ã¯é©ç”¨ã§ãã¾ã›ã‚“',
        'license_perpetual': 'æ°¸ç¶š',
        'license_expiry_warning': 'ãƒ©ã‚¤ã‚»ãƒ³ã‚¹ã®æœ‰åŠ¹æœŸé™ã¾ã§ã‚ã¨{0}æ—¥ã§ã™ï¼ˆ{1}ã¾ã§ï¼‰ã€‚æ›´æ–°ã‚’ã”æ¤œè¨ãã ã•ã„ã€‚',
        'license_expired': 'ãƒ©ã‚¤ã‚»ãƒ³ã‚¹ã®æœ‰åŠ¹æœŸé™ãŒåˆ‡ã‚Œã¾ã—ãŸã€‚ç¶™ç¶šã—ã¦ã”åˆ©ç”¨ã„ãŸã ãã«ã¯æ›´æ–°ãŒå¿…è¦ã§ã™ã€‚',
        'license_trial_link': 'ãƒˆãƒ©ã‚¤ã‚¢ãƒ«ç”³è«‹',
        'license_email_required': 'ãƒ¡ãƒ¼ãƒ«ã‚¢ãƒ‰ãƒ¬ã‚¹ã‚’å…¥åŠ›ã—ã¦ãã ã•ã„',
        'license_status_active': 'æœ‰åŠ¹',
        'license_status_expired': 'æœŸé™åˆ‡ã‚Œ',
        'license_valid_until': 'æœ‰åŠ¹æœŸé™: {0}',
        'license_days_remaining': 'ï¼ˆæ®‹ã‚Š{0}æ—¥ï¼‰',
        'license_feature_restricted': 'ã“ã®æ©Ÿèƒ½ã¯Proãƒ©ã‚¤ã‚»ãƒ³ã‚¹ãŒå¿…è¦ã§ã™ã€‚ç¾åœ¨: {0}',
        'license_batch_restricted': 'ãƒ•ã‚©ãƒ«ãƒ€ä¸€æ‹¬å‡¦ç†ã¯Proãƒ©ã‚¤ã‚»ãƒ³ã‚¹ãŒå¿…è¦ã§ã™ã€‚',
        'license_json_restricted': 'JSONå‡ºåŠ›ã¯Proãƒ©ã‚¤ã‚»ãƒ³ã‚¹ãŒå¿…è¦ã§ã™ã€‚',
        'license_continue_free': 'Freeç‰ˆã§ç¶šè¡Œ',
        # Status messages
        'status_slides_items': '{0}ã‚¹ãƒ©ã‚¤ãƒ‰ / {1}é …ç›®',
        'status_complete_items': 'å®Œäº†: {0}ä»¶',
        'status_batch_complete': 'ãƒãƒƒãƒæŠ½å‡ºå®Œäº†: {0}ä»¶ ({1})',
        'status_update_complete': 'æ›´æ–°å®Œäº†: {0}ä»¶',
        'lang_changed': 'è¨€èªã‚’å¤‰æ›´ã—ã¾ã—ãŸã€‚',
        # Log messages
        'log_cancelled': 'ã‚­ãƒ£ãƒ³ã‚»ãƒ«ã•ã‚Œã¾ã—ãŸ',
        'log_cancel_request': 'ã‚­ãƒ£ãƒ³ã‚»ãƒ«ã‚’ãƒªã‚¯ã‚¨ã‚¹ãƒˆ...',
        'log_no_text': 'ãƒ†ã‚­ã‚¹ãƒˆãŒè¦‹ã¤ã‹ã‚Šã¾ã›ã‚“ã§ã—ãŸ',
        'log_error': 'ã‚¨ãƒ©ãƒ¼: {0}',
        'log_found_files': 'ç™ºè¦‹: {0}ä»¶',
        'log_no_pptx_found': 'PPTXãƒ•ã‚¡ã‚¤ãƒ«ãŒè¦‹ã¤ã‹ã‚Šã¾ã›ã‚“',
        'log_invalid_header': 'ãƒ˜ãƒƒãƒ€ãƒ¼å½¢å¼ãŒä¸æ­£ã§ã™',
        'log_no_update_data': 'æ›´æ–°ãƒ‡ãƒ¼ã‚¿ãªã—',
        'log_processing': 'å‡¦ç†ä¸­...',
        'dialog_select_folder': 'ãƒ•ã‚©ãƒ«ãƒ€ã‚’é¸æŠ (PPTXãƒ•ã‚¡ã‚¤ãƒ«ã‚’å«ã‚€)',
        'dialog_select_folder_update': 'ãƒ•ã‚©ãƒ«ãƒ€ã‚’é¸æŠ (*_æŠ½å‡º{0} + PPTX)',
        'dialog_select_pptx': 'æ›´æ–°ã™ã‚‹PowerPointã‚’é¸æŠ',
        'dialog_processing_exit': 'å‡¦ç†ä¸­ã§ã™ã€‚çµ‚äº†ã—ã¾ã™ã‹ï¼Ÿ',
        'dialog_confirm_title': 'ç¢ºèª',
        'result_updated': 'æ›´æ–°: {0}ä»¶\nã‚¹ã‚­ãƒƒãƒ—: {1}ä»¶',
        'result_replaced': '{0} ä»¶ã‚’ç½®æ›ã—ã¾ã—ãŸ',
        'result_applied': '{0} ä»¶ã‚’åæ˜ ã—ã¾ã—ãŸ',
        'result_csv_saved': 'CSVã‚’ä¿å­˜ã—ã¾ã—ãŸ',
        'result_export_complete': 'ã‚¨ã‚¯ã‚¹ãƒãƒ¼ãƒˆå®Œäº†: {0}',
    },
}

_current_lang = 'ja'

def t(key: str, *args) -> str:
    text = LANGUAGES.get(_current_lang, LANGUAGES['ja']).get(key, key)
    if args:
        return text.format(*args)
    return text

def set_language(lang: str):
    global _current_lang
    if lang in LANGUAGES:
        _current_lang = lang

def get_language() -> str:
    return _current_lang


# ============== ãƒ©ã‚¤ã‚»ãƒ³ã‚¹è¨­å®šï¼ˆinsight-common çµ±åˆï¼‰ ==============
PRODUCT_CODE = ProductCode.SLIDE
EXPIRY_WARNING_DAYS = 30  # æœŸé™åˆ‡ã‚Œè­¦å‘Šã®æ—¥æ•°

# ãƒ­ãƒ¼ã‚«ãƒ«ãƒ†ã‚£ã‚¢å®šç¾©ï¼ˆFREEè¿½åŠ ï¼‰
class LicenseTier:
    FREE = "FREE"
    TRIAL = "TRIAL"
    STD = "STD"
    PRO = "PRO"
    ENT = "ENT"

# ãƒ†ã‚£ã‚¢åˆ¥è¨­å®šï¼ˆInsightSlideå›ºæœ‰ï¼‰
# json: 1ãƒ•ã‚¡ã‚¤ãƒ«JSONå…¥å‡ºåŠ›, batch: ãƒ•ã‚©ãƒ«ãƒ€ä¸€æ‹¬å‡¦ç†, compare: 2ãƒ•ã‚¡ã‚¤ãƒ«æ¯”è¼ƒ
TIERS = {
    LicenseTier.FREE: {'name': 'Free', 'name_ja': 'ç„¡æ–™ç‰ˆ', 'badge': 'Free', 'update_limit': 3, 'batch': False, 'json': False, 'compare': False},
    LicenseTier.TRIAL: {'name': 'Trial', 'name_ja': 'ãƒˆãƒ©ã‚¤ã‚¢ãƒ«', 'badge': 'Trial', 'update_limit': None, 'batch': True, 'json': True, 'compare': True},
    LicenseTier.STD: {'name': 'Standard', 'name_ja': 'ã‚¹ã‚¿ãƒ³ãƒ€ãƒ¼ãƒ‰', 'badge': 'Standard', 'update_limit': None, 'batch': False, 'json': False, 'compare': True},
    LicenseTier.PRO: {'name': 'Professional', 'name_ja': 'ãƒ—ãƒ­ãƒ•ã‚§ãƒƒã‚·ãƒ§ãƒŠãƒ«', 'badge': 'Pro', 'update_limit': None, 'batch': True, 'json': True, 'compare': True},
    LicenseTier.ENT: {'name': 'Enterprise', 'name_ja': 'ã‚¨ãƒ³ã‚¿ãƒ¼ãƒ—ãƒ©ã‚¤ã‚º', 'badge': 'Enterprise', 'update_limit': None, 'batch': True, 'json': True, 'compare': True},
}


class LicenseManager:
    """insight-common çµ±åˆãƒ©ã‚¤ã‚»ãƒ³ã‚¹ãƒãƒãƒ¼ã‚¸ãƒ£ãƒ¼"""

    def __init__(self):
        CONFIG_DIR.mkdir(parents=True, exist_ok=True)
        self.validator = LicenseValidator()
        self.license_info: Dict = {}
        self.insight_info: Optional[LicenseInfo] = None
        self._load_license()

    def _load_license(self):
        """ä¿å­˜ã•ã‚ŒãŸãƒ©ã‚¤ã‚»ãƒ³ã‚¹æƒ…å ±ã‚’èª­ã¿è¾¼ã‚€"""
        self.license_info = {'type': LicenseTier.FREE, 'key': '', 'email': '', 'expires': None}

        if LICENSE_FILE.exists():
            try:
                with open(LICENSE_FILE, 'r', encoding='utf-8') as f:
                    data = json.load(f)

                if data.get('key'):
                    # æœ‰åŠ¹æœŸé™ã‚’å¾©å…ƒ
                    expires_at = None
                    if data.get('expires'):
                        try:
                            expires_at = datetime.fromisoformat(data['expires'])
                        except:
                            pass

                    # insight-common ã§æ¤œè¨¼
                    self.insight_info = self.validator.validate(data['key'], expires_at)

                    if self.insight_info.is_valid:
                        # è£½å“ãƒã‚§ãƒƒã‚¯
                        if self.validator.is_product_covered(self.insight_info, PRODUCT_CODE):
                            tier = self._map_insight_tier(self.insight_info.tier)
                            self.license_info = {
                                'type': tier,
                                'key': data['key'],
                                'email': data.get('email', ''),
                                'expires': data.get('expires')
                            }
                            return

            except Exception as e:
                print(f"License load error: {e}")

    def _save_license(self):
        """ãƒ©ã‚¤ã‚»ãƒ³ã‚¹æƒ…å ±ã‚’ä¿å­˜"""
        with open(LICENSE_FILE, 'w', encoding='utf-8') as f:
            json.dump(self.license_info, f, ensure_ascii=False, indent=2)

    def _map_insight_tier(self, tier: Optional[InsightLicenseTier]) -> str:
        """insight-common ã®ãƒ†ã‚£ã‚¢ã‚’ãƒ­ãƒ¼ã‚«ãƒ«ãƒ†ã‚£ã‚¢ã«ãƒãƒƒãƒ—"""
        if not tier:
            return LicenseTier.FREE
        mapping = {
            InsightLicenseTier.TRIAL: LicenseTier.TRIAL,
            InsightLicenseTier.STD: LicenseTier.STD,
            InsightLicenseTier.PRO: LicenseTier.PRO,
            InsightLicenseTier.ENT: LicenseTier.ENT,
        }
        return mapping.get(tier, LicenseTier.FREE)

    def activate(self, email: str, key: str) -> Tuple[bool, str]:
        """ãƒ©ã‚¤ã‚»ãƒ³ã‚¹ã‚’ã‚¢ã‚¯ãƒ†ã‚£ãƒ™ãƒ¼ãƒˆ"""
        if not email or not key:
            return False, t('license_enter_prompt')

        # insight-common ã§æ¤œè¨¼
        self.insight_info = self.validator.validate(key.strip())

        if not self.insight_info.is_valid:
            error_msg = self.insight_info.error or t('license_invalid')
            return False, error_msg

        # è£½å“ãƒã‚§ãƒƒã‚¯
        if not self.validator.is_product_covered(self.insight_info, PRODUCT_CODE):
            return False, t('license_wrong_product')

        # æœ‰åŠ¹æœŸé™ã‚’è¨ˆç®—ï¼ˆåˆå›ã‚¢ã‚¯ãƒ†ã‚£ãƒ™ãƒ¼ã‚·ãƒ§ãƒ³æ™‚ï¼‰
        expires_str = None
        if self.insight_info.tier and self.insight_info.tier != InsightLicenseTier.ENT:
            tier_config = INSIGHT_TIERS.get(self.insight_info.tier, {})
            duration_months = tier_config.get('duration_months')
            duration_days = tier_config.get('duration_days')

            if duration_days:
                expires = datetime.now() + timedelta(days=duration_days)
                expires_str = expires.isoformat()
            elif duration_months:
                now = datetime.now()
                new_month = now.month + duration_months
                new_year = now.year + (new_month - 1) // 12
                new_month = (new_month - 1) % 12 + 1
                expires = datetime(new_year, new_month, min(now.day, 28))
                expires_str = expires.isoformat()

        tier = self._map_insight_tier(self.insight_info.tier)
        self.license_info = {
            'type': tier,
            'key': key.strip().upper(),
            'email': email.strip(),
            'expires': expires_str
        }
        self._save_license()

        tier_info = TIERS.get(tier, TIERS[LicenseTier.FREE])
        name = tier_info['name_ja'] if get_language() == 'ja' else tier_info['name']
        return True, t('license_activated', name)

    def deactivate(self):
        """ãƒ©ã‚¤ã‚»ãƒ³ã‚¹ã‚’è§£é™¤"""
        self.license_info = {'type': LicenseTier.FREE, 'key': '', 'email': '', 'expires': None}
        self.insight_info = None
        if LICENSE_FILE.exists():
            LICENSE_FILE.unlink()

    def get_tier(self) -> str:
        return self.license_info.get('type', LicenseTier.FREE)

    def get_tier_info(self) -> Dict:
        return TIERS.get(self.get_tier(), TIERS[LicenseTier.FREE])

    def get_update_limit(self) -> Optional[int]:
        return self.get_tier_info().get('update_limit')

    def can_batch(self) -> bool:
        """ãƒ•ã‚©ãƒ«ãƒ€ä¸€æ‹¬å‡¦ç†ãŒå¯èƒ½ã‹ï¼ˆPRO/Trial/ENTï¼‰"""
        return self.get_tier_info().get('batch', False)

    def can_json(self) -> bool:
        """1ãƒ•ã‚¡ã‚¤ãƒ«JSONå…¥å‡ºåŠ›ãŒå¯èƒ½ã‹ï¼ˆPRO/Trial/ENTï¼‰"""
        return self.get_tier_info().get('json', False)

    def can_compare(self) -> bool:
        """2ãƒ•ã‚¡ã‚¤ãƒ«æ¯”è¼ƒãŒå¯èƒ½ã‹ï¼ˆSTDä»¥ä¸Šï¼‰"""
        return self.get_tier_info().get('compare', False)

    def is_pro(self) -> bool:
        """å¾Œæ–¹äº’æ›æ€§ã®ãŸã‚ç¶­æŒï¼ˆPROæ©Ÿèƒ½ = batch + jsonï¼‰"""
        return self.can_batch() and self.can_json()

    def is_activated(self) -> bool:
        """ãƒ©ã‚¤ã‚»ãƒ³ã‚¹ãŒã‚¢ã‚¯ãƒ†ã‚£ãƒ™ãƒ¼ãƒˆã•ã‚Œã¦ã„ã‚‹ã‹"""
        return self.get_tier() != LicenseTier.FREE

    def get_days_until_expiry(self) -> Optional[int]:
        """æœ‰åŠ¹æœŸé™ã¾ã§ã®æ—¥æ•°ã‚’å–å¾—ï¼ˆæœŸé™ãªã—ã®å ´åˆã¯Noneï¼‰"""
        expires_str = self.license_info.get('expires')
        if not expires_str:
            return None
        try:
            expires = datetime.fromisoformat(expires_str)
            delta = expires - datetime.now()
            return delta.days
        except:
            return None

    def should_show_expiry_warning(self) -> bool:
        """æœŸé™åˆ‡ã‚Œè­¦å‘Šã‚’è¡¨ç¤ºã™ã¹ãã‹"""
        days = self.get_days_until_expiry()
        if days is None:
            return False
        return 0 < days <= EXPIRY_WARNING_DAYS

    def get_expiry_date_str(self) -> str:
        """æœ‰åŠ¹æœŸé™ã®è¡¨ç¤ºæ–‡å­—åˆ—"""
        expires_str = self.license_info.get('expires')
        if not expires_str:
            return t('license_perpetual') if self.get_tier() == LicenseTier.ENT else '-'
        try:
            expires = datetime.fromisoformat(expires_str)
            return expires.strftime('%Y/%m/%d')
        except:
            return '-'


# ============== ãƒ¢ãƒ€ãƒ³ãƒ‡ã‚¶ã‚¤ãƒ³ã‚·ã‚¹ãƒ†ãƒ  ==============
# B2B SaaSå“è³ª - Notion/Linear/Figmaé¢¨

# ã‚«ãƒ©ãƒ¼ãƒ‘ãƒ¬ãƒƒãƒˆï¼ˆæ´—ç·´ã•ã‚ŒãŸãƒ‹ãƒ¥ãƒ¼ãƒˆãƒ©ãƒ« + è½ã¡ç€ã„ãŸãƒ–ãƒ«ãƒ¼ï¼‰
COLOR_PALETTE = {
    # èƒŒæ™¯
    "bg_primary": "#FFFFFF",       # ãƒ¡ã‚¤ãƒ³èƒŒæ™¯
    "bg_secondary": "#F8FAFC",     # ã‚»ã‚«ãƒ³ãƒ€ãƒªèƒŒæ™¯ï¼ˆã‚«ãƒ¼ãƒ‰å†…ï¼‰
    "bg_elevated": "#F1F5F9",      # å¼·èª¿èƒŒæ™¯ï¼ˆãƒ›ãƒãƒ¼ç­‰ï¼‰
    "bg_sidebar": "#FAFBFC",       # ã‚µã‚¤ãƒ‰ãƒãƒ¼èƒŒæ™¯
    "bg_card": "#FFFFFF",          # ã‚«ãƒ¼ãƒ‰èƒŒæ™¯
    "bg_input": "#FFFFFF",         # å…¥åŠ›ãƒ•ã‚£ãƒ¼ãƒ«ãƒ‰èƒŒæ™¯

    # ãƒ†ã‚­ã‚¹ãƒˆï¼ˆ4æ®µéšã®éšå±¤ï¼‰
    "text_primary": "#1F2937",     # ãƒ¡ã‚¤ãƒ³ãƒ†ã‚­ã‚¹ãƒˆï¼ˆè¦‹å‡ºã—ï¼‰
    "text_secondary": "#374151",   # æœ¬æ–‡ãƒ†ã‚­ã‚¹ãƒˆ
    "text_tertiary": "#6B7280",    # è£œåŠ©ãƒ†ã‚­ã‚¹ãƒˆ
    "text_muted": "#9CA3AF",       # è–„ã„ãƒ†ã‚­ã‚¹ãƒˆï¼ˆæ³¨é‡ˆï¼‰
    "text_placeholder": "#D1D5DB", # ãƒ—ãƒ¬ãƒ¼ã‚¹ãƒ›ãƒ«ãƒ€ãƒ¼

    # ãƒ–ãƒ©ãƒ³ãƒ‰ã‚«ãƒ©ãƒ¼ï¼ˆè½ã¡ç€ã„ãŸãƒ–ãƒ«ãƒ¼ç³»ï¼‰
    "brand_primary": "#2563EB",    # ãƒ—ãƒ©ã‚¤ãƒãƒªãƒ–ãƒ«ãƒ¼
    "brand_hover": "#1D4ED8",      # ãƒ›ãƒãƒ¼æ™‚ï¼ˆæ¿ƒã„ï¼‰
    "brand_light": "#DBEAFE",      # è–„ã„ãƒ–ãƒ«ãƒ¼ï¼ˆé¸æŠèƒŒæ™¯ï¼‰
    "brand_muted": "#93C5FD",      # ãƒŸãƒ¥ãƒ¼ãƒˆãƒ–ãƒ«ãƒ¼

    # ã‚»ã‚«ãƒ³ãƒ€ãƒªã‚¢ã‚¯ã‚·ãƒ§ãƒ³
    "secondary_default": "#F3F4F6",  # ã‚»ã‚«ãƒ³ãƒ€ãƒªãƒœã‚¿ãƒ³èƒŒæ™¯
    "secondary_hover": "#E5E7EB",    # ã‚»ã‚«ãƒ³ãƒ€ãƒªãƒ›ãƒãƒ¼
    "secondary_border": "#D1D5DB",   # ã‚»ã‚«ãƒ³ãƒ€ãƒªãƒœãƒ¼ãƒ€ãƒ¼

    # æ©Ÿèƒ½åˆ¥ã‚«ãƒ©ãƒ¼
    "action_update": "#059669",    # æ›´æ–°ï¼ˆã‚°ãƒªãƒ¼ãƒ³ï¼‰
    "action_compare": "#7C3AED",   # æ¯”è¼ƒï¼ˆãƒ‘ãƒ¼ãƒ—ãƒ«ï¼‰
    "action_danger": "#DC2626",    # å±é™ºï¼ˆèµ¤ãƒ»æ§ãˆã‚ï¼‰

    # ã‚¹ãƒ†ãƒ¼ã‚¿ã‚¹
    "success": "#10B981",
    "success_light": "#D1FAE5",
    "warning": "#F59E0B",
    "warning_light": "#FEF3C7",
    "error": "#EF4444",
    "error_light": "#FEE2E2",
    "info": "#3B82F6",
    "info_light": "#DBEAFE",

    # ãƒœãƒ¼ãƒ€ãƒ¼ãƒ»åŒºåˆ‡ã‚Š
    "border_light": "#E5E7EB",     # è–„ã„ãƒœãƒ¼ãƒ€ãƒ¼
    "border_default": "#D1D5DB",   # æ¨™æº–ãƒœãƒ¼ãƒ€ãƒ¼
    "border_dark": "#9CA3AF",      # æ¿ƒã„ãƒœãƒ¼ãƒ€ãƒ¼
    "divider": "#F3F4F6",          # ã‚»ã‚¯ã‚·ãƒ§ãƒ³åŒºåˆ‡ã‚Š

    # å·®åˆ†è¡¨ç¤º
    "diff_changed": "#FEF3C7",
    "diff_added": "#D1FAE5",
    "diff_removed": "#FEE2E2",
}

# ãƒ•ã‚©ãƒ³ãƒˆè¨­å®šï¼ˆæ—¥æœ¬èªå¯¾å¿œï¼‰
FONT_FAMILY_SANS = "Meiryo UI"       # ã‚¯ãƒªãƒ¼ãƒ³ãªæ—¥æœ¬èªãƒ•ã‚©ãƒ³ãƒˆ
FONT_FAMILY_MONO = "MS Gothic"       # æ—¥æœ¬èªå¯¾å¿œç­‰å¹…ãƒ•ã‚©ãƒ³ãƒˆ

def get_fonts(size_preset: str = 'medium') -> dict:
    base = {'small': 10, 'medium': 11, 'large': 13}.get(size_preset, 11)
    return {
        # è¦‹å‡ºã—ç³»ï¼ˆSemiboldï¼‰
        "display": (FONT_FAMILY_SANS, base + 8, "bold"),      # ã‚¢ãƒ—ãƒªã‚¿ã‚¤ãƒˆãƒ«
        "title": (FONT_FAMILY_SANS, base + 4, "bold"),        # ç”»é¢ã‚¿ã‚¤ãƒˆãƒ«
        "heading": (FONT_FAMILY_SANS, base + 2, "bold"),      # ã‚»ã‚¯ã‚·ãƒ§ãƒ³è¦‹å‡ºã—

        # æœ¬æ–‡ç³»
        "body": (FONT_FAMILY_SANS, base, "normal"),           # æœ¬æ–‡
        "body_medium": (FONT_FAMILY_SANS, base, "bold"),      # æœ¬æ–‡ï¼ˆå¼·èª¿ï¼‰
        "body_bold": (FONT_FAMILY_SANS, base, "bold"),        # ãƒœã‚¿ãƒ³ãƒ©ãƒ™ãƒ«

        # è£œåŠ©ç³»
        "caption": (FONT_FAMILY_SANS, base - 1, "normal"),    # ã‚­ãƒ£ãƒ—ã‚·ãƒ§ãƒ³
        "small": (FONT_FAMILY_SANS, base - 2, "normal"),      # æ³¨é‡ˆ
        "tiny": (FONT_FAMILY_SANS, base - 3, "normal"),       # æ¥µå°

        # ãƒ­ã‚°ãƒ»ãƒ‡ãƒ¼ã‚¿è¡¨ç¤ºç”¨ï¼ˆæ—¥æœ¬èªå¯¾å¿œï¼‰
        "mono": (FONT_FAMILY_SANS, base, "normal"),
        "mono_small": (FONT_FAMILY_SANS, base - 1, "normal"),
    }

FONTS = get_fonts('medium')

# ã‚¹ãƒšãƒ¼ã‚·ãƒ³ã‚°ã‚·ã‚¹ãƒ†ãƒ ï¼ˆ8pxãƒ™ãƒ¼ã‚¹ï¼‰
SPACING = {
    "none": 0,
    "xs": 4,
    "sm": 8,
    "md": 12,
    "lg": 16,
    "xl": 24,
    "2xl": 32,
    "3xl": 48,
}

# è§’ä¸¸
RADIUS = {
    "none": 0,
    "sm": 4,
    "default": 6,
    "md": 8,
    "lg": 12,
    "full": 9999,
}


class ConfigManager:
    DEFAULT = {
        'language': 'ja', 'output_format': 'excel', 'include_metadata': True,
        'auto_backup': True, 'last_directory': '', 'font_size': 'medium',
        'advanced_expanded': False,
    }

    def __init__(self):
        global FONTS
        CONFIG_DIR.mkdir(parents=True, exist_ok=True)
        self.config = self._load()
        set_language(self.config.get('language', 'ja'))
        FONTS = get_fonts(self.config.get('font_size', 'medium'))

    def _load(self) -> Dict:
        if CONFIG_FILE.exists():
            try:
                with open(CONFIG_FILE, 'r', encoding='utf-8') as f:
                    return {**self.DEFAULT, **json.load(f)}
            except:
                pass
        return self.DEFAULT.copy()

    def save(self):
        with open(CONFIG_FILE, 'w', encoding='utf-8') as f:
            json.dump(self.config, f, ensure_ascii=False, indent=2)

    def get(self, key: str, default=None):
        return self.config.get(key, default)

    def set(self, key: str, value):
        self.config[key] = value
        self.save()


def save_error_log(error: Exception, context: str = ""):
    try:
        CONFIG_DIR.mkdir(parents=True, exist_ok=True)
        with open(ERROR_LOG_FILE, 'a', encoding='utf-8') as f:
            f.write(f"\n{'='*60}\n{datetime.now()}\n{context}\n{error}\n{traceback.format_exc()}\n")
    except:
        pass


# ============== ã‚°ãƒªãƒƒãƒ‰UI (Undo/Redoå¯¾å¿œ) ==============
class UndoManager:
    def __init__(self, max_history: int = 50):
        self.undo_stack: List[Dict] = []
        self.redo_stack: List[Dict] = []
        self.max_history = max_history

    def push(self, action: Dict):
        self.undo_stack.append(action)
        self.redo_stack.clear()
        if len(self.undo_stack) > self.max_history:
            self.undo_stack.pop(0)

    def undo(self) -> Optional[Dict]:
        if not self.undo_stack:
            return None
        action = self.undo_stack.pop()
        self.redo_stack.append(action)
        return action

    def redo(self) -> Optional[Dict]:
        if not self.redo_stack:
            return None
        action = self.redo_stack.pop()
        self.undo_stack.append(action)
        return action

    def clear(self):
        self.undo_stack.clear()
        self.redo_stack.clear()


class EditableGrid(ttk.Frame):
    """ã‚¤ãƒ³ãƒ©ã‚¤ãƒ³ç·¨é›†å¯¾å¿œã‚°ãƒªãƒƒãƒ‰ + ãƒ•ã‚£ãƒ«ã‚¿æ©Ÿèƒ½"""

    def __init__(self, parent, on_change=None, **kwargs):
        super().__init__(parent, **kwargs)

        self.on_change = on_change
        self.undo_manager = UndoManager()
        self._edit_widget = None
        self._editing_item = None
        self._editing_column = None
        self._all_data: List[Dict] = []
        self._filter_text = ""

        self._create_widgets()
        self._setup_bindings()

    def _create_widgets(self):
        # ãƒ„ãƒ¼ãƒ«ãƒãƒ¼
        toolbar = ttk.Frame(self)
        toolbar.pack(fill="x", pady=(0, 5))

        # ãƒ•ã‚£ãƒ«ã‚¿
        ttk.Label(toolbar, text=t('filter_label')).pack(side="left", padx=(0, 5))
        self.filter_var = tk.StringVar()
        self.filter_entry = ttk.Entry(toolbar, textvariable=self.filter_var, width=20)
        self.filter_entry.pack(side="left", padx=(0, 5))
        self.filter_var.trace_add("write", lambda *args: self._apply_filter())

        ttk.Button(toolbar, text=t('btn_clear_grid'), command=self._clear_filter).pack(side="left")

        # ã‚¹ãƒšãƒ¼ã‚µãƒ¼
        ttk.Frame(toolbar).pack(side="left", fill="x", expand=True)

        # ä¸€æ‹¬ç½®æ›ãƒœã‚¿ãƒ³
        ttk.Button(toolbar, text=t('btn_replace_all'), command=self._show_replace_dialog).pack(side="left", padx=2)

        # Undo/Redo
        self.undo_btn = ttk.Button(toolbar, text=t('btn_undo'), command=self._do_undo)
        self.undo_btn.pack(side="left", padx=2)
        self.redo_btn = ttk.Button(toolbar, text=t('btn_redo'), command=self._do_redo)
        self.redo_btn.pack(side="left", padx=2)

        # Treeview
        tree_frame = ttk.Frame(self)
        tree_frame.pack(fill="both", expand=True)

        columns = ("slide", "id", "type", "text")
        self.tree = ttk.Treeview(tree_frame, columns=columns, show="headings")

        self.tree.heading("slide", text=t('header_slide'))
        self.tree.heading("id", text=t('header_id'))
        self.tree.heading("type", text=t('header_type'))
        self.tree.heading("text", text=t('header_text'))

        self.tree.column("slide", width=80, minwidth=60, anchor="center", stretch=False)
        self.tree.column("id", width=100, minwidth=80, stretch=False)
        self.tree.column("type", width=80, minwidth=60, stretch=False)
        self.tree.column("text", width=800, minwidth=300, stretch=True)

        vsb = ttk.Scrollbar(tree_frame, orient="vertical", command=self.tree.yview)
        hsb = ttk.Scrollbar(tree_frame, orient="horizontal", command=self.tree.xview)
        self.tree.configure(yscrollcommand=vsb.set, xscrollcommand=hsb.set)

        self.tree.grid(row=0, column=0, sticky="nsew")
        vsb.grid(row=0, column=1, sticky="ns")
        hsb.grid(row=1, column=0, sticky="ew")

        tree_frame.grid_rowconfigure(0, weight=1)
        tree_frame.grid_columnconfigure(0, weight=1)

        self.tree.tag_configure("modified", background=COLOR_PALETTE["diff_changed"])
        self.tree.tag_configure("filtered", background=COLOR_PALETTE["diff_added"])

    def _setup_bindings(self):
        self.tree.bind("<Double-1>", self._on_double_click)
        self.tree.bind("<Control-z>", lambda e: self._do_undo())
        self.tree.bind("<Control-y>", lambda e: self._do_redo())

    def _on_double_click(self, event):
        region = self.tree.identify_region(event.x, event.y)
        if region != "cell":
            return

        column = self.tree.identify_column(event.x)
        item = self.tree.identify_row(event.y)

        if not item or column != "#4":  # textã‚«ãƒ©ãƒ ã®ã¿ç·¨é›†å¯èƒ½
            return

        self._start_edit(item, "text")

    def _start_edit(self, item: str, column: str):
        bbox = self.tree.bbox(item, column)
        if not bbox:
            return

        current_value = self.tree.set(item, column)
        self._editing_item = item
        self._editing_column = column

        self._edit_widget = tk.Entry(self.tree, font=FONTS["body"])
        self._edit_widget.insert(0, current_value)
        self._edit_widget.select_range(0, tk.END)
        self._edit_widget.place(x=bbox[0], y=bbox[1], width=bbox[2], height=bbox[3])
        self._edit_widget.focus_set()

        self._edit_widget.bind("<Return>", lambda e: self._finish_edit())
        self._edit_widget.bind("<Escape>", lambda e: self._cancel_edit())
        self._edit_widget.bind("<FocusOut>", lambda e: self._finish_edit())

    def _finish_edit(self):
        if not self._edit_widget or not self._editing_item:
            return

        new_value = self._edit_widget.get()
        old_value = self.tree.set(self._editing_item, self._editing_column)

        if new_value != old_value:
            self.tree.set(self._editing_item, self._editing_column, new_value)
            self.tree.item(self._editing_item, tags=("modified",))

            self.undo_manager.push({
                "type": "edit", "item": self._editing_item,
                "column": self._editing_column, "old": old_value, "new": new_value
            })

            # å…ƒãƒ‡ãƒ¼ã‚¿ã‚‚æ›´æ–°
            idx = self.tree.index(self._editing_item)
            if idx < len(self._all_data):
                self._all_data[idx]["text"] = new_value

            if self.on_change:
                self.on_change(self._editing_item, self._editing_column, new_value)

        self._cancel_edit()

    def _cancel_edit(self):
        if self._edit_widget:
            self._edit_widget.destroy()
            self._edit_widget = None
        self._editing_item = None
        self._editing_column = None

    def _do_undo(self):
        action = self.undo_manager.undo()
        if action:
            self.tree.set(action["item"], action["column"], action["old"])

    def _do_redo(self):
        action = self.undo_manager.redo()
        if action:
            self.tree.set(action["item"], action["column"], action["new"])
            self.tree.item(action["item"], tags=("modified",))

    def _apply_filter(self):
        self._filter_text = self.filter_var.get().lower()
        self._refresh_display()

    def _clear_filter(self):
        self.filter_var.set("")
        self._filter_text = ""
        self._refresh_display()

    def _refresh_display(self):
        for item in self.tree.get_children():
            self.tree.delete(item)

        for row in self._all_data:
            if self._filter_text:
                if self._filter_text not in str(row.get("text", "")).lower():
                    continue

            self.tree.insert("", "end", values=(
                row.get("slide", ""),
                row.get("id", ""),
                row.get("type", ""),
                row.get("text", "")
            ))

    def _show_replace_dialog(self):
        dialog = tk.Toplevel(self)
        dialog.title(t('btn_replace_all'))
        dialog.geometry("400x150")
        dialog.transient(self)
        dialog.grab_set()

        ttk.Label(dialog, text=t('replace_search')).grid(row=0, column=0, padx=10, pady=10, sticky="w")
        find_var = tk.StringVar()
        ttk.Entry(dialog, textvariable=find_var, width=40).grid(row=0, column=1, padx=10, pady=10)

        ttk.Label(dialog, text=t('replace_with')).grid(row=1, column=0, padx=10, pady=10, sticky="w")
        replace_var = tk.StringVar()
        ttk.Entry(dialog, textvariable=replace_var, width=40).grid(row=1, column=1, padx=10, pady=10)

        def do_replace():
            find_text = find_var.get()
            replace_text = replace_var.get()
            if not find_text:
                return

            count = 0
            for item in self.tree.get_children():
                old_text = self.tree.set(item, "text")
                if find_text in old_text:
                    new_text = old_text.replace(find_text, replace_text)
                    self.tree.set(item, "text", new_text)
                    self.tree.item(item, tags=("modified",))
                    count += 1

            dialog.destroy()
            messagebox.showinfo(t('dialog_complete'), t('result_replaced', count))

        ttk.Button(dialog, text=t('btn_replace'), command=do_replace).grid(row=2, column=1, pady=10, sticky="e")

    def load_data(self, data: List[Dict]):
        self._all_data = data.copy()
        self.undo_manager.clear()
        self._refresh_display()

    def clear(self):
        self._all_data = []
        for item in self.tree.get_children():
            self.tree.delete(item)
        self.undo_manager.clear()

    def get_data(self) -> List[Dict]:
        result = []
        for item in self.tree.get_children():
            result.append({
                "slide": self.tree.set(item, "slide"),
                "id": self.tree.set(item, "id"),
                "type": self.tree.set(item, "type"),
                "text": self.tree.set(item, "text"),
            })
        return result


# ============== æ¯”è¼ƒæ©Ÿèƒ½ ==============
class CompareDialog:
    def __init__(self, parent, callback):
        self.callback = callback
        self.dialog = tk.Toplevel(parent)
        self.dialog.title("PPTXæ¯”è¼ƒ")
        self.dialog.geometry("600x280")
        self.dialog.transient(parent)
        self.dialog.grab_set()

        self._create_widgets()

    def _create_widgets(self):
        frame = ttk.Frame(self.dialog, padding=20)
        frame.pack(fill='both', expand=True)

        ttk.Label(frame, text=t('compare_title'), font=FONTS["heading"]).pack(anchor='w', pady=(0, 15))

        # ãƒ•ã‚¡ã‚¤ãƒ«1
        f1 = ttk.Frame(frame)
        f1.pack(fill='x', pady=5)
        ttk.Label(f1, text=t('compare_file1'), width=12).pack(side='left')
        self.file1_var = tk.StringVar()
        ttk.Entry(f1, textvariable=self.file1_var, width=45).pack(side='left', padx=5)
        ttk.Button(f1, text=t('btn_browse'), command=lambda: self._browse(self.file1_var)).pack(side='left')

        # ãƒ•ã‚¡ã‚¤ãƒ«2
        f2 = ttk.Frame(frame)
        f2.pack(fill='x', pady=5)
        ttk.Label(f2, text=t('compare_file2'), width=12).pack(side='left')
        self.file2_var = tk.StringVar()
        ttk.Entry(f2, textvariable=self.file2_var, width=45).pack(side='left', padx=5)
        ttk.Button(f2, text=t('btn_browse'), command=lambda: self._browse(self.file2_var)).pack(side='left')

        # ã‚ªãƒ—ã‚·ãƒ§ãƒ³
        opt = ttk.Frame(frame)
        opt.pack(fill='x', pady=15)
        self.ignore_ws = tk.BooleanVar(value=True)
        ttk.Checkbutton(opt, text=t('compare_ignore_ws'), variable=self.ignore_ws).pack(side='left')

        # ãƒœã‚¿ãƒ³
        btn = ttk.Frame(frame)
        btn.pack(fill='x', pady=10)
        ttk.Button(btn, text=t('btn_cancel'), command=self.dialog.destroy).pack(side='left')
        tk.Button(btn, text=t('btn_run_compare'), bg=COLOR_PALETTE["brand_primary"], fg="#FFFFFF",
                  command=self._execute).pack(side='left', padx=10)

    def _browse(self, var):
        path = filedialog.askopenfilename(filetypes=[("PowerPoint", "*.pptx")])
        if path:
            var.set(path)

    def _execute(self):
        f1, f2 = self.file1_var.get(), self.file2_var.get()
        if not f1 or not f2:
            messagebox.showwarning("è­¦å‘Š", "2ã¤ã®ãƒ•ã‚¡ã‚¤ãƒ«ã‚’é¸æŠã—ã¦ãã ã•ã„")
            return
        self.callback(f1, f2, self.ignore_ws.get())
        self.dialog.destroy()


class CompareResultWindow:
    def __init__(self, parent, file1_name, file2_name, diff_data, stats, on_apply=None):
        self.window = tk.Toplevel(parent)
        self.window.title(f"æ¯”è¼ƒçµæœ: {file1_name} â†” {file2_name}")
        self.window.geometry("1100x700")

        self.diff_data = diff_data
        self.on_apply = on_apply
        self.selections = {}

        for i, row in enumerate(diff_data):
            if row["status"] == "å¤‰æ›´":
                self.selections[i] = None
            elif row["status"] == "è¿½åŠ ":
                self.selections[i] = "after"
            elif row["status"] == "å‰Šé™¤":
                self.selections[i] = "before"
            else:
                self.selections[i] = "same"

        self._create_widgets(stats, file1_name, file2_name)

    def _create_widgets(self, stats, f1, f2):
        # çµ±è¨ˆ
        top = ttk.Frame(self.window, padding=10)
        top.pack(fill='x')
        ttk.Label(top, text=f"ğŸ“Š {stats['same']} | {stats['changed']} | {stats['added']} | {stats['removed']}",
                  font=FONTS["heading"]).pack(side='left')

        ttk.Button(top, text=t('btn_export_csv'), command=self._export_csv).pack(side='right')

        # ã‚°ãƒªãƒƒãƒ‰
        grid_frame = ttk.Frame(self.window, padding=10)
        grid_frame.pack(fill='both', expand=True)

        cols = ("select", "slide", "id", "status", "before", "after")
        self.tree = ttk.Treeview(grid_frame, columns=cols, show="headings")

        self.tree.heading("select", text=t('header_select'))
        self.tree.heading("slide", text=t('header_slide'))
        self.tree.heading("id", text="ID")
        self.tree.heading("status", text=t('header_status'))
        self.tree.heading("before", text=f"å…ƒ: {f1}")
        self.tree.heading("after", text=f"æ–°: {f2}")

        self.tree.column("select", width=60, anchor="center")
        self.tree.column("slide", width=60, anchor="center")
        self.tree.column("id", width=80)
        self.tree.column("status", width=60, anchor="center")
        self.tree.column("before", width=350)
        self.tree.column("after", width=350)

        vsb = ttk.Scrollbar(grid_frame, orient="vertical", command=self.tree.yview)
        self.tree.configure(yscrollcommand=vsb.set)
        vsb.pack(side='right', fill='y')
        self.tree.pack(fill='both', expand=True)

        self.tree.tag_configure("same", background=COLOR_PALETTE["bg_secondary"])
        self.tree.tag_configure("changed", background=COLOR_PALETTE["diff_changed"])
        self.tree.tag_configure("added", background=COLOR_PALETTE["diff_added"])
        self.tree.tag_configure("removed", background=COLOR_PALETTE["diff_removed"])

        self.tree.bind("<Button-1>", self._on_click)
        self.item_map = {}
        self._refresh()

        # ãƒœã‚¿ãƒ³
        bottom = ttk.Frame(self.window, padding=10)
        bottom.pack(fill='x')
        ttk.Button(bottom, text=t('btn_select_original'), command=lambda: self._select_all("before")).pack(side='left', padx=2)
        ttk.Button(bottom, text=t('btn_select_new'), command=lambda: self._select_all("after")).pack(side='left', padx=2)
        tk.Button(bottom, text=t('btn_apply_selection'), font=(FONT_FAMILY_SANS, 10),
                  bg=COLOR_PALETTE["action_update"], fg="#FFFFFF", relief="flat",
                  activebackground="#047857", padx=SPACING["lg"], pady=SPACING["sm"],
                  cursor="hand2", command=self._apply).pack(side='right', padx=5)
        ttk.Button(bottom, text=t('btn_close'), command=self.window.destroy).pack(side='right')

    def _refresh(self):
        for item in self.tree.get_children():
            self.tree.delete(item)
        self.item_map = {}

        for i, row in enumerate(self.diff_data):
            sel = self.selections.get(i)
            sel_text = {"before": "â—€ å…ƒ", "after": "æ–° â–¶", "same": "â”€"}.get(sel, "")
            tag = {"ä¸€è‡´": "same", "å¤‰æ›´": "changed", "è¿½åŠ ": "added", "å‰Šé™¤": "removed"}.get(row["status"], "same")

            before = (row.get("before") or "").replace("\n", " â†µ ")[:50]
            after = (row.get("after") or "").replace("\n", " â†µ ")[:50]

            item_id = self.tree.insert("", "end", values=(
                sel_text, row["slide"], row.get("id", ""), row["status"], before, after
            ), tags=(tag,))
            self.item_map[item_id] = i

    def _on_click(self, event):
        item = self.tree.identify_row(event.y)
        col = self.tree.identify_column(event.x)
        if not item:
            return

        idx = self.item_map.get(item)
        if idx is None or self.diff_data[idx]["status"] == "ä¸€è‡´":
            return

        current = self.selections.get(idx)
        if col == "#5":
            self.selections[idx] = "before"
        elif col == "#6":
            self.selections[idx] = "after"
        else:
            self.selections[idx] = "after" if current == "before" else "before"
        self._refresh()

    def _select_all(self, choice):
        for i, row in enumerate(self.diff_data):
            if row["status"] != "ä¸€è‡´":
                self.selections[i] = choice
        self._refresh()

    def _apply(self):
        selected = []
        for i, row in enumerate(self.diff_data):
            sel = self.selections.get(i)
            if sel in ("before", "after"):
                text = row["before"] if sel == "before" else row["after"]
                selected.append({"slide": row["slide"], "id": row.get("id"), "text": text})

        if not selected:
            messagebox.showwarning("è­¦å‘Š", "åæ˜ ã™ã‚‹é …ç›®ãŒã‚ã‚Šã¾ã›ã‚“")
            return

        if self.on_apply:
            self.on_apply(selected)
            messagebox.showinfo(t('dialog_complete'), t('result_applied', len(selected)))
            self.window.destroy()

    def _export_csv(self):
        path = filedialog.asksaveasfilename(defaultextension=".csv", filetypes=[("CSV", "*.csv")])
        if not path:
            return

        with open(path, 'w', newline='', encoding='utf-8-sig') as f:
            w = csv.writer(f)
            w.writerow(["ã‚¹ãƒ©ã‚¤ãƒ‰", "ID", "çŠ¶æ…‹", "å…ƒ", "æ–°"])
            for row in self.diff_data:
                w.writerow([row["slide"], row.get("id", ""), row["status"], row.get("before", ""), row.get("after", "")])
        messagebox.showinfo(t('dialog_complete'), t('result_csv_saved'))


# ============== ãƒ¡ã‚¤ãƒ³ã‚¢ãƒ—ãƒªã‚±ãƒ¼ã‚·ãƒ§ãƒ³ ==============
class InsightSlidesApp:
    def __init__(self, root):
        self.root = root
        self.license_manager = LicenseManager()
        self.config_manager = ConfigManager()
        self.processing = False
        self.cancel_requested = False
        self.presentation = None
        self.log_buffer = []
        self.extracted_data = []  # ã‚°ãƒªãƒƒãƒ‰ç”¨
        self.include_notes_var = tk.BooleanVar(value=False)
        self.auto_backup_var = tk.BooleanVar(value=self.config_manager.get('auto_backup', True))

        self._setup_window()
        self._apply_styles()
        self._create_menu()
        self._create_layout()
        self.root.protocol("WM_DELETE_WINDOW", self._on_closing)

        # èµ·å‹•æ™‚ãƒ©ã‚¤ã‚»ãƒ³ã‚¹ãƒã‚§ãƒƒã‚¯ï¼ˆUIãŒè¡¨ç¤ºã•ã‚ŒãŸå¾Œã«å®Ÿè¡Œï¼‰
        self.root.after(100, self._check_license_on_startup)

    def _setup_window(self):
        tier = self.license_manager.get_tier_info()
        self.root.title(f"{APP_NAME} v{APP_VERSION} - {tier['name']}")
        self.root.geometry("1300x900")
        self.root.minsize(1100, 700)
        self.root.configure(bg=COLOR_PALETTE["bg_primary"])

    def _apply_styles(self):
        """ã‚·ãƒ³ãƒ—ãƒ«ã§çµ±ä¸€æ„Ÿã®ã‚ã‚‹ã‚¹ã‚¿ã‚¤ãƒ«"""
        self.style = ttk.Style()
        self.style.theme_use('clam')

        # çµ±ä¸€èƒŒæ™¯è‰²ï¼ˆå…¨ä½“ã§ä¸€è²«æ€§ã‚’æŒãŸã›ã‚‹ï¼‰
        BG = COLOR_PALETTE["bg_primary"]  # #FFFFFF
        BG_LIGHT = COLOR_PALETTE["bg_secondary"]  # #F8FAFC
        TEXT = COLOR_PALETTE["text_primary"]  # #1F2937
        TEXT_SUB = COLOR_PALETTE["text_tertiary"]  # #6B7280
        BORDER = COLOR_PALETTE["border_light"]  # #E5E7EB

        # === ãƒ•ãƒ¬ãƒ¼ãƒ  ===
        self.style.configure('Main.TFrame', background=BG)
        self.style.configure('Card.TFrame', background=BG)
        self.style.configure('Sidebar.TFrame', background=BG)
        self.style.configure('TFrame', background=BG)

        # === ãƒ©ãƒ™ãƒ«ãƒ•ãƒ¬ãƒ¼ãƒ  ===
        self.style.configure('TLabelframe', background=BG, bordercolor=BORDER)
        self.style.configure('TLabelframe.Label', background=BG, foreground=TEXT,
                            font=(FONT_FAMILY_SANS, 11, "bold"))

        # === ãƒ©ãƒ™ãƒ«ï¼ˆå…¨ã¦åŒã˜èƒŒæ™¯ï¼‰ ===
        self.style.configure('TLabel', background=BG, foreground=TEXT,
                            font=(FONT_FAMILY_SANS, 10))
        self.style.configure('Muted.TLabel', background=BG, foreground=TEXT_SUB,
                            font=(FONT_FAMILY_SANS, 9))
        self.style.configure('Caption.TLabel', background=BG, foreground=TEXT_SUB,
                            font=(FONT_FAMILY_SANS, 9))

        # === ãƒœã‚¿ãƒ³ ===
        self.style.configure('TButton', background=BG_LIGHT, foreground=TEXT,
                            bordercolor=BORDER, padding=(12, 6),
                            font=(FONT_FAMILY_SANS, 10))
        self.style.map('TButton',
                      background=[('active', COLOR_PALETTE["bg_elevated"])])

        # === ãƒã‚§ãƒƒã‚¯ãƒœãƒƒã‚¯ã‚¹ ===
        self.style.configure('TCheckbutton', background=BG, foreground=TEXT,
                            font=(FONT_FAMILY_SANS, 10))
        self.style.map('TCheckbutton', background=[('active', BG)])

        # === ã‚³ãƒ³ãƒœãƒœãƒƒã‚¯ã‚¹ ===
        self.style.configure('TCombobox', fieldbackground=BG, background=BG,
                            foreground=TEXT, bordercolor=BORDER,
                            padding=(4, 2), font=(FONT_FAMILY_SANS, 10))

        # === ã‚¨ãƒ³ãƒˆãƒª ===
        self.style.configure('TEntry', fieldbackground=BG, foreground=TEXT,
                            bordercolor=BORDER, padding=(4, 2))

        # === Notebookï¼ˆã‚¿ãƒ–ï¼‰ ===
        self.style.configure('TNotebook', background=BG, bordercolor=BORDER)
        self.style.configure('TNotebook.Tab', background=BG_LIGHT, foreground=TEXT_SUB,
                            padding=(16, 8), font=(FONT_FAMILY_SANS, 10))
        self.style.map('TNotebook.Tab',
                      background=[('selected', BG)],
                      foreground=[('selected', TEXT)])

        # === ãƒ—ãƒ­ã‚°ãƒ¬ã‚¹ãƒãƒ¼ ===
        self.style.configure('TProgressbar', background=COLOR_PALETTE["brand_primary"],
                            troughcolor=BG_LIGHT, bordercolor=BORDER)

        # === Treeview ===
        self.style.configure('Treeview', background=BG, foreground=TEXT,
                            fieldbackground=BG, bordercolor=BORDER,
                            rowheight=28, font=(FONT_FAMILY_SANS, 10))
        self.style.configure('Treeview.Heading', background=BG_LIGHT, foreground=TEXT,
                            bordercolor=BORDER, font=(FONT_FAMILY_SANS, 10, "bold"))
        self.style.map('Treeview',
                      background=[('selected', COLOR_PALETTE["brand_light"])],
                      foreground=[('selected', COLOR_PALETTE["brand_primary"])])

        # === ã‚¹ã‚¯ãƒ­ãƒ¼ãƒ«ãƒãƒ¼ ===
        self.style.configure('Vertical.TScrollbar', background=BG_LIGHT,
                            troughcolor=BG, bordercolor=BG)
        self.style.configure('Horizontal.TScrollbar', background=BG_LIGHT,
                            troughcolor=BG, bordercolor=BG)

    def _create_menu(self):
        # ãƒ¡ãƒ‹ãƒ¥ãƒ¼ã‚¹ã‚¿ã‚¤ãƒ«è¨­å®š
        menu_font = (FONT_FAMILY_SANS, 9)
        menu_style = {
            'font': menu_font,
            'bg': COLOR_PALETTE["bg_primary"],
            'fg': COLOR_PALETTE["text_primary"],
            'activebackground': COLOR_PALETTE["brand_primary"],
            'activeforeground': '#FFFFFF',
            'relief': 'flat',
            'bd': 0,
        }

        menubar = tk.Menu(self.root, **menu_style)
        self.root.config(menu=menubar)

        help_menu = tk.Menu(menubar, tearoff=0, **menu_style)
        menubar.add_cascade(label=t('menu_help'), menu=help_menu)
        help_menu.add_command(label=t('menu_guide'), command=lambda: webbrowser.open(SUPPORT_LINKS["tutorial"]))
        help_menu.add_command(label=t('menu_faq'), command=lambda: webbrowser.open(SUPPORT_LINKS["faq"]))
        help_menu.add_separator()
        help_menu.add_command(label=t('menu_license'), command=self._show_license_dialog)
        help_menu.add_separator()

        lang_menu = tk.Menu(help_menu, tearoff=0, **menu_style)
        help_menu.add_cascade(label=t('lang_menu'), menu=lang_menu)
        lang_menu.add_command(label="English", command=lambda: self._change_language('en'))
        lang_menu.add_command(label="æ—¥æœ¬èª", command=lambda: self._change_language('ja'))

        help_menu.add_separator()
        help_menu.add_command(label=t('menu_about'), command=self._show_about)

    def _create_layout(self):
        if hasattr(self, 'main_container') and self.main_container:
            self.main_container.destroy()

        self.main_container = ttk.Frame(self.root, style='Main.TFrame')
        self.main_container.pack(fill='both', expand=True, padx=SPACING["xl"], pady=SPACING["xl"])
        self.main_container.grid_columnconfigure(0, weight=1)
        self.main_container.grid_rowconfigure(1, weight=1)

        self._create_header(self.main_container)

        content = ttk.Frame(self.main_container, style='Main.TFrame')
        content.grid(row=1, column=0, sticky='nsew')
        content.grid_columnconfigure(1, weight=1)
        content.grid_rowconfigure(0, weight=1)

        self._create_controls(content)
        self._create_output(content)

    def _create_header(self, parent):
        """æ´—ç·´ã•ã‚ŒãŸãƒ˜ãƒƒãƒ€ãƒ¼ - ä½™ç™½ã§åŒºåˆ‡ã‚Šã€ã‚·ãƒ³ãƒ—ãƒ«ã«"""
        header = tk.Frame(parent, bg=COLOR_PALETTE["bg_primary"])
        header.grid(row=0, column=0, sticky='ew', pady=(0, SPACING["lg"]))

        # å·¦: ã‚¿ã‚¤ãƒˆãƒ« + ã‚µãƒ–ã‚¿ã‚¤ãƒˆãƒ«
        left = tk.Frame(header, bg=COLOR_PALETTE["bg_primary"])
        left.pack(side='left')

        # ã‚¢ãƒ—ãƒªåï¼ˆã‚·ãƒ³ãƒ—ãƒ«ã«ï¼‰
        tk.Label(left, text="Insight Slides", font=FONTS["display"],
                 fg=COLOR_PALETTE["text_primary"], bg=COLOR_PALETTE["bg_primary"]).pack(side='left')

        # ãƒ©ã‚¤ã‚»ãƒ³ã‚¹ãƒãƒƒã‚¸ï¼ˆã‚ã‚Œã°ï¼‰
        tier = self.license_manager.get_tier_info()
        if tier['name'] != 'Free':
            badge = tk.Label(left, text=f" {tier['name']} ", font=FONTS["small"],
                            fg=COLOR_PALETTE["brand_primary"], bg=COLOR_PALETTE["brand_light"],
                            padx=6, pady=2)
            badge.pack(side='left', padx=(SPACING["sm"], 0))

        # ã‚µãƒ–ã‚¿ã‚¤ãƒˆãƒ«
        tk.Label(left, text=t('app_subtitle'), font=FONTS["caption"],
                 fg=COLOR_PALETTE["text_muted"], bg=COLOR_PALETTE["bg_primary"]).pack(side='left', padx=(SPACING["lg"], 0))

        # å³: ãƒãƒ¼ã‚¸ãƒ§ãƒ³
        right = tk.Frame(header, bg=COLOR_PALETTE["bg_primary"])
        right.pack(side='right')

        tk.Label(right, text=f"v{APP_VERSION}", font=FONTS["small"],
                 fg=COLOR_PALETTE["text_muted"], bg=COLOR_PALETTE["bg_primary"]).pack(side='right')

    def _create_controls(self, parent):
        """å·¦ã‚µã‚¤ãƒ‰ãƒãƒ¼ - 3ã‚»ã‚¯ã‚·ãƒ§ãƒ³æ§‹æˆï¼ˆå…¥åŠ›/å‡ºåŠ›/ãƒ•ã‚©ãƒ«ãƒ€ä¸€æ‹¬ï¼‰"""
        frame = ttk.Frame(parent, style='Sidebar.TFrame')
        frame.grid(row=0, column=0, sticky='nsew', padx=(0, SPACING["xl"]))
        frame.grid_rowconfigure(4, weight=1)

        btn_font = (FONT_FAMILY_SANS, 10)
        can_json = self.license_manager.can_json()
        can_batch = self.license_manager.can_batch()
        can_compare = self.license_manager.can_compare()

        # ============ å…¥åŠ›ï¼ˆ1ãƒ•ã‚¡ã‚¤ãƒ«ï¼‰ã‚»ã‚¯ã‚·ãƒ§ãƒ³ ============
        input_card = ttk.LabelFrame(frame, text=t('panel_input'), padding=SPACING["md"])
        input_card.grid(row=0, column=0, sticky='ew', pady=(0, SPACING["md"]))
        input_card.grid_columnconfigure(0, weight=1)

        # PPTXèª­è¾¼ãƒœã‚¿ãƒ³ï¼ˆãƒ—ãƒ©ã‚¤ãƒãƒªï¼‰
        tk.Button(input_card, text=t('btn_load_pptx'), font=btn_font,
                  bg=COLOR_PALETTE["brand_primary"], fg="#FFFFFF", relief="flat",
                  activebackground=COLOR_PALETTE["brand_hover"],
                  padx=SPACING["lg"], pady=SPACING["sm"],
                  cursor="hand2", command=self._extract_single).grid(row=0, column=0, sticky='ew', pady=(0, SPACING["xs"]))

        # Excelèª­è¾¼ãƒœã‚¿ãƒ³
        tk.Button(input_card, text=t('btn_load_excel'), font=btn_font,
                  bg=COLOR_PALETTE["secondary_default"], fg=COLOR_PALETTE["text_secondary"], relief="flat",
                  activebackground=COLOR_PALETTE["secondary_hover"],
                  padx=SPACING["md"], pady=SPACING["sm"],
                  cursor="hand2", command=self._load_excel_to_grid).grid(row=1, column=0, sticky='ew', pady=(0, SPACING["xs"]))

        # JSONèª­è¾¼ãƒœã‚¿ãƒ³ï¼ˆProï¼‰
        if can_json:
            tk.Button(input_card, text=t('btn_load_json'), font=btn_font,
                      bg=COLOR_PALETTE["secondary_default"], fg=COLOR_PALETTE["text_secondary"], relief="flat",
                      activebackground=COLOR_PALETTE["secondary_hover"],
                      padx=SPACING["md"], pady=SPACING["sm"],
                      cursor="hand2", command=self._load_json_to_grid).grid(row=2, column=0, sticky='ew')
        else:
            tk.Label(input_card, text=f"{t('btn_load_json')} (Pro)", font=btn_font,
                     fg=COLOR_PALETTE["text_muted"], bg=COLOR_PALETTE["bg_primary"]).grid(row=2, column=0, sticky='w')

        # ============ å‡ºåŠ›ï¼ˆ1ãƒ•ã‚¡ã‚¤ãƒ«ï¼‰ã‚»ã‚¯ã‚·ãƒ§ãƒ³ ============
        output_card = ttk.LabelFrame(frame, text=t('panel_output_file'), padding=SPACING["md"])
        output_card.grid(row=1, column=0, sticky='ew', pady=(0, SPACING["md"]))
        output_card.grid_columnconfigure(0, weight=1)

        # ã‚¹ãƒ©ã‚¤ãƒ‰åˆ¶é™è­¦å‘Š
        limit = self.license_manager.get_update_limit()
        if limit:
            warn_frame = tk.Frame(output_card, bg=COLOR_PALETTE["warning_light"], padx=SPACING["sm"], pady=SPACING["xs"])
            warn_frame.grid(row=0, column=0, sticky='ew', pady=(0, SPACING["sm"]))
            tk.Label(warn_frame, text=t('msg_update_limit', limit), font=FONTS["small"],
                    fg=COLOR_PALETTE["warning"], bg=COLOR_PALETTE["warning_light"]).pack(anchor='w')

        # PPTXã«åæ˜ ãƒœã‚¿ãƒ³ï¼ˆãƒ—ãƒ©ã‚¤ãƒãƒªï¼‰
        tk.Button(output_card, text=t('btn_apply_pptx'), font=btn_font,
                  bg=COLOR_PALETTE["action_update"], fg="#FFFFFF", relief="flat",
                  activebackground="#047857",
                  padx=SPACING["lg"], pady=SPACING["sm"],
                  cursor="hand2", command=self._apply_grid_to_pptx).grid(row=1, column=0, sticky='ew', pady=(0, SPACING["xs"]))

        # Excelå‡ºåŠ›ãƒœã‚¿ãƒ³
        tk.Button(output_card, text=t('btn_export_to_excel'), font=btn_font,
                  bg=COLOR_PALETTE["secondary_default"], fg=COLOR_PALETTE["text_secondary"], relief="flat",
                  activebackground=COLOR_PALETTE["secondary_hover"],
                  padx=SPACING["md"], pady=SPACING["sm"],
                  cursor="hand2", command=self._export_grid_excel).grid(row=2, column=0, sticky='ew', pady=(0, SPACING["xs"]))

        # JSONå‡ºåŠ›ãƒœã‚¿ãƒ³ï¼ˆProï¼‰
        if can_json:
            tk.Button(output_card, text=t('btn_export_to_json'), font=btn_font,
                      bg=COLOR_PALETTE["secondary_default"], fg=COLOR_PALETTE["text_secondary"], relief="flat",
                      activebackground=COLOR_PALETTE["secondary_hover"],
                      padx=SPACING["md"], pady=SPACING["sm"],
                      cursor="hand2", command=self._export_grid_json).grid(row=3, column=0, sticky='ew')
        else:
            tk.Label(output_card, text=f"{t('btn_export_to_json')} (Pro)", font=btn_font,
                     fg=COLOR_PALETTE["text_muted"], bg=COLOR_PALETTE["bg_primary"]).grid(row=3, column=0, sticky='w')

        # ============ ãƒ•ã‚©ãƒ«ãƒ€ä¸€æ‹¬ã‚»ã‚¯ã‚·ãƒ§ãƒ³ ============
        batch_card = ttk.LabelFrame(frame, text=t('panel_batch'), padding=SPACING["md"])
        batch_card.grid(row=2, column=0, sticky='ew', pady=(0, SPACING["md"]))
        batch_card.grid_columnconfigure(0, weight=1)

        if can_batch:
            # Proãƒãƒƒã‚¸
            tk.Label(batch_card, text="Pro", font=(FONT_FAMILY_SANS, 8, 'bold'),
                     fg=COLOR_PALETTE["brand_primary"], bg=COLOR_PALETTE["bg_primary"]).grid(row=0, column=0, sticky='e')

            # ä¸€æ‹¬æŠ½å‡ºãƒœã‚¿ãƒ³
            tk.Button(batch_card, text=t('btn_batch_extract'), font=btn_font,
                      bg=COLOR_PALETTE["secondary_default"], fg=COLOR_PALETTE["text_secondary"], relief="flat",
                      activebackground=COLOR_PALETTE["secondary_hover"],
                      padx=SPACING["md"], pady=SPACING["sm"],
                      cursor="hand2", command=self._batch_extract_dialog).grid(row=1, column=0, sticky='ew', pady=(0, SPACING["xs"]))

            # ä¸€æ‹¬æ›´æ–°ãƒœã‚¿ãƒ³
            tk.Button(batch_card, text=t('btn_batch_update'), font=btn_font,
                      bg=COLOR_PALETTE["secondary_default"], fg=COLOR_PALETTE["text_secondary"], relief="flat",
                      activebackground=COLOR_PALETTE["secondary_hover"],
                      padx=SPACING["md"], pady=SPACING["sm"],
                      cursor="hand2", command=self._batch_update_dialog).grid(row=2, column=0, sticky='ew')
        else:
            tk.Label(batch_card, text=f"{t('btn_batch_extract')} (Pro)", font=btn_font,
                     fg=COLOR_PALETTE["text_muted"], bg=COLOR_PALETTE["bg_primary"]).grid(row=0, column=0, sticky='w', pady=(0, SPACING["xs"]))
            tk.Label(batch_card, text=f"{t('btn_batch_update')} (Pro)", font=btn_font,
                     fg=COLOR_PALETTE["text_muted"], bg=COLOR_PALETTE["bg_primary"]).grid(row=1, column=0, sticky='w')

        # ============ 2ãƒ•ã‚¡ã‚¤ãƒ«æ¯”è¼ƒãƒœã‚¿ãƒ³ ============
        compare_text = t('btn_compare') if can_compare else f"{t('btn_compare')} (STD)"
        tk.Button(frame, text=compare_text, font=btn_font,
                  bg=COLOR_PALETTE["secondary_default"] if can_compare else COLOR_PALETTE["bg_secondary"],
                  fg=COLOR_PALETTE["text_secondary"] if can_compare else COLOR_PALETTE["text_muted"],
                  activebackground=COLOR_PALETTE["secondary_hover"] if can_compare else COLOR_PALETTE["bg_secondary"],
                  relief="flat", padx=SPACING["md"], pady=SPACING["sm"],
                  cursor="hand2" if can_compare else "arrow",
                  command=self._show_compare_dialog if can_compare else None,
                  state='normal' if can_compare else 'disabled').grid(row=3, column=0, sticky='ew', pady=(0, SPACING["md"]))

        # ã‚¹ãƒ†ãƒ¼ã‚¿ã‚¹ï¼†ãƒŸãƒ‹ãƒ­ã‚°
        status_frame = ttk.Frame(frame, style='Main.TFrame')
        status_frame.grid(row=4, column=0, sticky='sew')

        # ãƒ—ãƒ­ã‚°ãƒ¬ã‚¹ãƒãƒ¼ï¼ˆå‡¦ç†ä¸­ã®ã¿è¡¨ç¤ºï¼‰
        self.progress = ttk.Progressbar(status_frame, mode='indeterminate')
        self.progress.pack(fill='x', pady=(0, SPACING["sm"]))

        # ãƒŸãƒ‹ãƒ­ã‚°ï¼ˆ1-2è¡Œã€ã‚¯ãƒªãƒƒã‚¯ã§è©³ç´°è¡¨ç¤ºï¼‰
        log_frame = tk.Frame(status_frame, bg=COLOR_PALETTE["bg_secondary"], cursor="hand2")
        log_frame.pack(fill='x')
        log_frame.bind("<Button-1>", lambda e: self._show_log_detail())

        self.mini_log_label = tk.Label(log_frame, text=t('status_waiting'),
                                       font=(FONT_FAMILY_SANS, 9), fg=COLOR_PALETTE["text_tertiary"],
                                       bg=COLOR_PALETTE["bg_secondary"], anchor='w', padx=SPACING["sm"], pady=SPACING["xs"])
        self.mini_log_label.pack(fill='x')
        self.mini_log_label.bind("<Button-1>", lambda e: self._show_log_detail())

        # è©³ç´°ãƒªãƒ³ã‚¯
        detail_link = tk.Label(log_frame, text=t('show_detail'),
                               font=(FONT_FAMILY_SANS, 8), fg=COLOR_PALETTE["brand_primary"],
                               bg=COLOR_PALETTE["bg_secondary"], cursor="hand2")
        detail_link.pack(anchor='e', padx=SPACING["sm"], pady=(0, SPACING["xs"]))
        detail_link.bind("<Button-1>", lambda e: self._show_log_detail())

        # ã‚­ãƒ£ãƒ³ã‚»ãƒ«ãƒœã‚¿ãƒ³ï¼ˆå‡¦ç†ä¸­ã®ã¿ã‚¢ã‚¯ãƒ†ã‚£ãƒ–ï¼‰
        btn_frame = ttk.Frame(status_frame)
        btn_frame.pack(fill='x', pady=(SPACING["sm"], 0))
        self.cancel_btn = ttk.Button(btn_frame, text=t('btn_cancel'), command=self._cancel, state='disabled')
        self.cancel_btn.pack(side='left')

    def _create_output(self, parent):
        """å³å´ãƒ¡ã‚¤ãƒ³ã‚³ãƒ³ãƒ†ãƒ³ãƒ„ - ç·¨é›†å°‚ç”¨ã‚¨ãƒªã‚¢"""
        # ãƒ¡ã‚¤ãƒ³ã‚«ãƒ¼ãƒ‰ï¼ˆã‚¿ã‚¤ãƒˆãƒ«ãªã— - æ§‹é€ ã§å½¹å‰²ã‚’ç¤ºã™ï¼‰
        card = ttk.Frame(parent, style='Card.TFrame')
        card.grid(row=0, column=1, sticky='nsew')
        card.grid_columnconfigure(0, weight=1)
        card.grid_rowconfigure(1, weight=1)

        # ãƒ•ã‚¡ã‚¤ãƒ«æƒ…å ±ãƒ˜ãƒƒãƒ€ãƒ¼ï¼ˆã‚³ãƒ³ãƒ‘ã‚¯ãƒˆï¼‰
        file_info_frame = tk.Frame(card, bg=COLOR_PALETTE["bg_primary"], pady=SPACING["sm"])
        file_info_frame.grid(row=0, column=0, sticky='ew', padx=SPACING["md"])

        self.file_name_label = tk.Label(file_info_frame, text="",
                                        font=(FONT_FAMILY_SANS, 10), bg=COLOR_PALETTE["bg_primary"],
                                        fg=COLOR_PALETTE["text_secondary"])
        self.file_name_label.pack(side='left')

        self.file_info_detail = tk.Label(file_info_frame, text="",
                                         font=FONTS["caption"], bg=COLOR_PALETTE["bg_primary"],
                                         fg=COLOR_PALETTE["text_tertiary"])
        self.file_info_detail.pack(side='right')

        # ãƒ¡ã‚¤ãƒ³ç·¨é›†ã‚¨ãƒªã‚¢ï¼ˆã‚°ãƒªãƒƒãƒ‰ï¼‰
        edit_area = ttk.Frame(card, style='Main.TFrame')
        edit_area.grid(row=1, column=0, sticky='nsew', padx=SPACING["md"])
        edit_area.grid_columnconfigure(0, weight=1)
        edit_area.grid_rowconfigure(0, weight=1)

        # ã‚¦ã‚§ãƒ«ã‚«ãƒ ã‚¬ã‚¤ãƒ‰ï¼ˆåˆæœŸçŠ¶æ…‹ï¼‰
        self.welcome_frame = tk.Frame(edit_area, bg=COLOR_PALETTE["bg_primary"])
        self.welcome_frame.grid(row=0, column=0, sticky='nsew')
        self.welcome_frame.grid_columnconfigure(0, weight=1)
        self.welcome_frame.grid_rowconfigure(0, weight=1)
        self._create_welcome_guide()

        # ã‚°ãƒªãƒƒãƒ‰ãƒ“ãƒ¥ãƒ¼ï¼ˆãƒ‡ãƒ¼ã‚¿èª­è¾¼å¾Œã«è¡¨ç¤ºï¼‰
        self.grid_container = ttk.Frame(edit_area, style='Main.TFrame')
        self.grid_view = EditableGrid(self.grid_container, on_change=self._on_grid_change)
        self.grid_view.grid(row=0, column=0, sticky='nsew')
        self.grid_container.grid_columnconfigure(0, weight=1)
        self.grid_container.grid_rowconfigure(0, weight=1)

        # ã‚¢ã‚¯ã‚·ãƒ§ãƒ³ãƒãƒ¼ï¼ˆä¸‹éƒ¨å›ºå®šï¼‰
        action_bar = tk.Frame(card, bg=COLOR_PALETTE["bg_primary"], pady=SPACING["md"])
        action_bar.grid(row=2, column=0, sticky='ew', padx=SPACING["md"])

        # ãƒ—ãƒ©ã‚¤ãƒãƒªã‚¢ã‚¯ã‚·ãƒ§ãƒ³
        self.apply_btn = tk.Button(action_bar, text=t('btn_apply'), font=(FONT_FAMILY_SANS, 10),
                  bg=COLOR_PALETTE["action_update"], fg="#FFFFFF", relief="flat",
                  padx=SPACING["lg"], pady=SPACING["sm"],
                  activebackground="#047857",
                  cursor="hand2", command=self._apply_grid_to_pptx, state='disabled')
        self.apply_btn.pack(side='right')

        # ã‚¨ã‚¯ã‚¹ãƒãƒ¼ãƒˆãƒœã‚¿ãƒ³ï¼ˆExcelï¼‰
        self.export_excel_btn = tk.Button(action_bar, text=t('btn_export_excel'), font=(FONT_FAMILY_SANS, 10),
                  bg=COLOR_PALETTE["secondary_default"], fg=COLOR_PALETTE["text_secondary"], relief="flat",
                  padx=SPACING["md"], pady=SPACING["sm"],
                  activebackground=COLOR_PALETTE["secondary_hover"],
                  cursor="hand2", command=self._export_grid_excel, state='disabled')
        self.export_excel_btn.pack(side='right', padx=(0, SPACING["sm"]))

        # ã‚¨ã‚¯ã‚¹ãƒãƒ¼ãƒˆãƒœã‚¿ãƒ³ï¼ˆJSONï¼‰- Proç‰ˆä»¥ä¸Š
        if self.license_manager.can_json():
            self.export_json_btn = tk.Button(action_bar, text=t('btn_export_json'), font=(FONT_FAMILY_SANS, 10),
                      bg=COLOR_PALETTE["secondary_default"], fg=COLOR_PALETTE["text_secondary"], relief="flat",
                      padx=SPACING["md"], pady=SPACING["sm"],
                      activebackground=COLOR_PALETTE["secondary_hover"],
                      cursor="hand2", command=self._export_grid_json, state='disabled')
            self.export_json_btn.pack(side='right', padx=(0, SPACING["sm"]))
        else:
            self.export_json_btn = tk.Button(action_bar, text=f"{t('btn_export_json')} (Pro)", font=(FONT_FAMILY_SANS, 10),
                      bg=COLOR_PALETTE["bg_secondary"], fg=COLOR_PALETTE["text_muted"], relief="flat",
                      padx=SPACING["md"], pady=SPACING["sm"], state='disabled')
            self.export_json_btn.pack(side='right', padx=(0, SPACING["sm"]))

    def _create_welcome_guide(self):
        """åˆæœŸçŠ¶æ…‹ã®ã‚¦ã‚§ãƒ«ã‚«ãƒ ã‚¬ã‚¤ãƒ‰"""
        center_frame = tk.Frame(self.welcome_frame, bg=COLOR_PALETTE["bg_primary"])
        center_frame.place(relx=0.5, rely=0.45, anchor='center')

        # ã‚¿ã‚¤ãƒˆãƒ«
        tk.Label(center_frame, text=t('welcome_guide_title'),
                 font=(FONT_FAMILY_SANS, 16, "bold"), fg=COLOR_PALETTE["text_primary"],
                 bg=COLOR_PALETTE["bg_primary"]).pack(pady=(0, SPACING["lg"]))

        # æ‰‹é †
        steps = [
            ("1", t('guide_step1')),
            ("2", t('guide_step2')),
            ("3", t('guide_step3')),
            ("4", t('guide_step4')),
        ]

        for num, text in steps:
            step_frame = tk.Frame(center_frame, bg=COLOR_PALETTE["bg_primary"])
            step_frame.pack(anchor='w', pady=SPACING["xs"])

            tk.Label(step_frame, text=num, font=(FONT_FAMILY_SANS, 10, "bold"),
                     fg=COLOR_PALETTE["brand_primary"], bg=COLOR_PALETTE["bg_primary"],
                     width=2).pack(side='left')
            tk.Label(step_frame, text=text, font=(FONT_FAMILY_SANS, 10),
                     fg=COLOR_PALETTE["text_secondary"], bg=COLOR_PALETTE["bg_primary"]).pack(side='left')

    def _show_edit_area(self):
        """ã‚¦ã‚§ãƒ«ã‚«ãƒ ã‚¬ã‚¤ãƒ‰ã‚’éš ã—ã¦ã‚°ãƒªãƒƒãƒ‰ã‚’è¡¨ç¤º"""
        self.welcome_frame.grid_remove()
        self.grid_container.grid(row=0, column=0, sticky='nsew')
        self.apply_btn.configure(state='normal')
        self.export_excel_btn.configure(state='normal')
        self.export_json_btn.configure(state='normal')

    def _show_welcome_area(self):
        """ã‚°ãƒªãƒƒãƒ‰ã‚’éš ã—ã¦ã‚¦ã‚§ãƒ«ã‚«ãƒ ã‚¬ã‚¤ãƒ‰ã‚’è¡¨ç¤º"""
        self.grid_container.grid_remove()
        self.welcome_frame.grid(row=0, column=0, sticky='nsew')
        self.apply_btn.configure(state='disabled')
        self.export_excel_btn.configure(state='disabled')
        self.export_json_btn.configure(state='disabled')
        self.file_name_label.configure(text="")
        self.file_info_detail.configure(text="")

    def _update_file_info(self, filename: str, item_count: int = 0, slide_count: int = 0):
        """ãƒ•ã‚¡ã‚¤ãƒ«æƒ…å ±ãƒ˜ãƒƒãƒ€ãƒ¼ã‚’æ›´æ–°"""
        self.file_name_label.configure(text=filename)
        if item_count > 0:
            self.file_info_detail.configure(text=t('status_slides_items', slide_count, item_count))
        else:
            self.file_info_detail.configure(text="")

    def _show_welcome(self):
        """åˆæœŸã‚¦ã‚§ãƒ«ã‚«ãƒ è¡¨ç¤ºï¼ˆãƒŸãƒ‹ãƒ­ã‚°ã®ã¿æ›´æ–°ï¼‰"""
        tier = self.license_manager.get_tier_info()
        self._update_mini_log(f"{APP_NAME} v{APP_VERSION} ({tier['name']}) - æº–å‚™å®Œäº†")

    # === Output helpers ===
    def _update_output(self, text, clear=False):
        """ãƒ­ã‚°ãƒãƒƒãƒ•ã‚¡ã«è¿½åŠ ã—ã€ãƒŸãƒ‹ãƒ­ã‚°ã‚’æ›´æ–°"""
        if clear:
            self.log_buffer = []
        self.log_buffer.append(text)
        # ãƒŸãƒ‹ãƒ­ã‚°ã«ã¯æœ€æ–°ã®1è¡Œã®ã¿è¡¨ç¤º
        last_line = text.strip().split('\n')[-1] if text.strip() else ""
        self._update_mini_log(last_line)

    def _update_output_safe(self, text, clear=False):
        self.root.after(0, lambda: self._update_output(text, clear))

    def _update_mini_log(self, text):
        """ãƒŸãƒ‹ãƒ­ã‚°ãƒ©ãƒ™ãƒ«ã‚’æ›´æ–°ï¼ˆæœ€æ–°ãƒ¡ãƒƒã‚»ãƒ¼ã‚¸ã®ã¿ï¼‰"""
        # é•·ã™ãã‚‹ãƒ†ã‚­ã‚¹ãƒˆã¯çœç•¥
        max_len = 50
        display_text = text[:max_len] + "..." if len(text) > max_len else text
        self.mini_log_label.configure(text=display_text)

    def _update_mini_log_safe(self, text):
        self.root.after(0, lambda: self._update_mini_log(text))

    def _update_status(self, text, color=None):
        """ã‚¹ãƒ†ãƒ¼ã‚¿ã‚¹æ›´æ–°ï¼ˆãƒŸãƒ‹ãƒ­ã‚°ã«çµ±åˆï¼‰"""
        self._update_mini_log(text)

    def _update_status_safe(self, text, color=None):
        self.root.after(0, lambda: self._update_status(text, color))

    def _log(self, msg, level="info"):
        timestamp = datetime.now().strftime("%H:%M:%S")
        prefix = {"error": "Ã— ", "warning": "! ", "success": "âœ“ "}.get(level, "")
        full_msg = f"[{timestamp}] {prefix}{msg}"
        self._update_output_safe(f"{full_msg}\n")
        # ã‚¨ãƒ©ãƒ¼æ™‚ã¯è‰²ã‚’å¤‰ãˆã‚‹
        if level == "error":
            self.mini_log_label.configure(fg=COLOR_PALETTE["error"])
        elif level == "success":
            self.mini_log_label.configure(fg=COLOR_PALETTE["success"])
        else:
            self.mini_log_label.configure(fg=COLOR_PALETTE["text_tertiary"])

    def _start_progress(self):
        self.progress.start(10)
        self.processing = True
        self.cancel_requested = False
        self.root.after(0, lambda: self.cancel_btn.configure(state='normal'))

    def _stop_progress(self):
        self.progress.stop()
        self.processing = False
        self.root.after(0, lambda: self.cancel_btn.configure(state='disabled'))

    def _cancel(self):
        if self.processing:
            self.cancel_requested = True
            self._log(t('log_cancel_request'), "warning")

    def _show_log_detail(self):
        """ãƒ­ã‚°è©³ç´°ãƒ¢ãƒ¼ãƒ€ãƒ«ã‚’è¡¨ç¤º"""
        dialog = tk.Toplevel(self.root)
        dialog.title("å‡¦ç†ãƒ­ã‚°")
        dialog.geometry("600x400")
        dialog.transient(self.root)
        dialog.grab_set()

        # ãƒ­ã‚°ãƒ†ã‚­ã‚¹ãƒˆã‚¨ãƒªã‚¢
        text_frame = ttk.Frame(dialog)
        text_frame.pack(fill='both', expand=True, padx=SPACING["md"], pady=SPACING["md"])

        log_text = scrolledtext.ScrolledText(text_frame, wrap=tk.WORD,
                                              font=FONTS["mono"],
                                              bg=COLOR_PALETTE["bg_primary"],
                                              fg=COLOR_PALETTE["text_secondary"],
                                              relief="flat", bd=1)
        log_text.pack(fill='both', expand=True)

        # ãƒ­ã‚°å†…å®¹ã‚’è¡¨ç¤º
        log_content = "".join(self.log_buffer) if self.log_buffer else "ãƒ­ã‚°ã¯ã‚ã‚Šã¾ã›ã‚“"
        log_text.insert('1.0', log_content)
        log_text.configure(state=tk.DISABLED)
        log_text.see(tk.END)

        # ãƒœã‚¿ãƒ³ãƒ•ãƒ¬ãƒ¼ãƒ 
        btn_frame = tk.Frame(dialog, bg=COLOR_PALETTE["bg_primary"])
        btn_frame.pack(fill='x', padx=SPACING["md"], pady=(0, SPACING["md"]))

        def copy_log():
            content = "".join(self.log_buffer)
            if content:
                self.root.clipboard_clear()
                self.root.clipboard_append(content)
                messagebox.showinfo("ã‚³ãƒ”ãƒ¼å®Œäº†", "ãƒ­ã‚°ã‚’ã‚¯ãƒªãƒƒãƒ—ãƒœãƒ¼ãƒ‰ã«ã‚³ãƒ”ãƒ¼ã—ã¾ã—ãŸ")

        def clear_log():
            self.log_buffer = []
            log_text.configure(state=tk.NORMAL)
            log_text.delete('1.0', tk.END)
            log_text.insert('1.0', "ãƒ­ã‚°ã‚’ã‚¯ãƒªã‚¢ã—ã¾ã—ãŸ")
            log_text.configure(state=tk.DISABLED)
            self._update_mini_log("æº–å‚™å®Œäº†")

        tk.Button(btn_frame, text=t('btn_copy_log'), font=(FONT_FAMILY_SANS, 9),
                  bg=COLOR_PALETTE["secondary_default"], fg=COLOR_PALETTE["text_secondary"],
                  relief="flat", padx=SPACING["md"], pady=SPACING["xs"],
                  command=copy_log).pack(side='left', padx=(0, SPACING["sm"]))

        tk.Button(btn_frame, text=t('btn_clear_log'), font=(FONT_FAMILY_SANS, 9),
                  bg=COLOR_PALETTE["secondary_default"], fg=COLOR_PALETTE["text_secondary"],
                  relief="flat", padx=SPACING["md"], pady=SPACING["xs"],
                  command=clear_log).pack(side='left')

        tk.Button(btn_frame, text=t('btn_close'), font=(FONT_FAMILY_SANS, 9),
                  bg=COLOR_PALETTE["brand_primary"], fg="#FFFFFF",
                  relief="flat", padx=SPACING["md"], pady=SPACING["xs"],
                  command=dialog.destroy).pack(side='right')

    def _change_language(self, lang):
        if lang != get_language():
            self.config_manager.set('language', lang)
            set_language(lang)
            self._create_menu()  # ãƒ¡ãƒ‹ãƒ¥ãƒ¼ã‚’å†ä½œæˆ
            self._create_layout()
            self._setup_window()  # ã‚¦ã‚£ãƒ³ãƒ‰ã‚¦ã‚¿ã‚¤ãƒˆãƒ«ã‚’æ›´æ–°
            messagebox.showinfo(t('dialog_complete'), t('lang_changed'))

    # === Utility ===
    def clean_text(self, text):
        if text is None:
            return ""
        text = re.sub(r'[\x00-\x08\x0B-\x0C\x0E-\x1F]', '', text)
        text = text.replace('\r\n', '\n').replace('\r', '\n').replace('\v', '\n')
        return text

    def _normalize_for_compare(self, text):
        if text is None:
            return ""
        text = text.replace('\r\n', '\n').replace('\r', '\n').replace('\v', '\n')
        text = text.replace('\u00A0', ' ').replace('\u3000', ' ')
        return text.strip()

    def _texts_are_equal(self, old_text, new_text):
        return self._normalize_for_compare(old_text) == self._normalize_for_compare(new_text)

    def get_shape_type(self, shape):
        try:
            if shape.is_placeholder:
                types_ja = {1: "ã‚¿ã‚¤ãƒˆãƒ«", 2: "æœ¬æ–‡", 3: "å›³è¡¨", 4: "æ—¥ä»˜", 5: "ã‚¹ãƒ©ã‚¤ãƒ‰ç•ªå·"}
                return types_ja.get(shape.placeholder_format.type, "ãã®ä»–")
            elif hasattr(shape, "has_table") and shape.has_table:
                return "è¡¨"
            elif shape.shape_type == 1:
                return "ãƒ†ã‚­ã‚¹ãƒˆãƒœãƒƒã‚¯ã‚¹"
            return "ãã®ä»–"
        except:
            return "ä¸æ˜"

    def _create_backup(self, path: str):
        if not self.license_manager.is_pro() or not self.auto_backup_var.get():
            return
        try:
            backup_dir = Path(path).parent / "backup"
            backup_dir.mkdir(exist_ok=True)
            backup_name = f"{Path(path).stem}_{datetime.now().strftime('%Y%m%d_%H%M%S')}{Path(path).suffix}"
            shutil.copy2(path, backup_dir / backup_name)
            self._log(f"ãƒãƒƒã‚¯ã‚¢ãƒƒãƒ—ä½œæˆ: {backup_name}")
        except Exception as e:
            self._log(f"ãƒãƒƒã‚¯ã‚¢ãƒƒãƒ—å¤±æ•—: {e}", "warning")

    # === Extract ===
    def extract_from_ppt(self, path: str, include_notes: bool = False) -> Tuple[List, Dict]:
        try:
            prs = pptx.Presentation(path)
            data = []
            meta = {'file_name': os.path.basename(path), 'slide_count': len(prs.slides)}

            for slide_num, slide in enumerate(prs.slides, 1):
                if self.cancel_requested:
                    break
                for shape in slide.shapes:
                    try:
                        sid = str(shape.shape_id)
                        stype = self.get_shape_type(shape)

                        if hasattr(shape, "text") and shape.text.strip():
                            data.append({
                                "slide": slide_num, "id": sid, "type": stype, "text": self.clean_text(shape.text)
                            })

                        if hasattr(shape, "has_table") and shape.has_table:
                            for r, row in enumerate(shape.table.rows):
                                for c, cell in enumerate(row.cells):
                                    if cell.text.strip():
                                        data.append({
                                            "slide": slide_num, "id": f"{sid}_t{r}_{c}",
                                            "type": f"è¡¨({r+1},{c+1})", "text": self.clean_text(cell.text)
                                        })
                    except:
                        pass

                if include_notes:
                    try:
                        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                            notes_text = slide.notes_slide.notes_text_frame.text.strip()
                            if notes_text:
                                data.append({
                                    "slide": slide_num, "id": "notes", "type": t('type_notes'),
                                    "text": self.clean_text(notes_text)
                                })
                    except:
                        pass

            return data, meta
        except Exception as e:
            save_error_log(e, f"extract_from_ppt: {path}")
            self._log(f"èª­ã¿è¾¼ã¿ã‚¨ãƒ©ãƒ¼: {e}", "error")
            return [], {}

    def save_to_file(self, data: List[Dict], path: str, fmt: str = "excel") -> bool:
        try:
            if fmt == "excel":
                wb = openpyxl.Workbook()
                ws = wb.active
                ws.append([t('header_slide'), t('header_id'), t('header_type'), t('header_text')])
                for row in data:
                    ws.append([row["slide"], row["id"], row["type"], row["text"]])
                wb.save(path)
            elif fmt == "json":
                with open(path, 'w', encoding='utf-8') as f:
                    json.dump(data, f, ensure_ascii=False, indent=2)
            else:
                with open(path, 'w', encoding='utf-8', newline='') as f:
                    w = csv.writer(f, delimiter='\t')
                    w.writerow([t('header_slide'), t('header_id'), t('header_type'), t('header_text')])
                    for row in data:
                        w.writerow([row["slide"], row["id"], row["type"], row["text"]])
            return True
        except Exception as e:
            save_error_log(e, f"save_to_file: {path}")
            self._log(f"ä¿å­˜ã‚¨ãƒ©ãƒ¼: {e}", "error")
            return False

    def _extract_single(self):
        if self.processing:
            return
        path = filedialog.askopenfilename(title="PowerPointãƒ•ã‚¡ã‚¤ãƒ«ã‚’é¸æŠ", filetypes=[("PowerPoint", "*.pptx")])
        if not path:
            return
        self.config_manager.set('last_directory', os.path.dirname(path))

        include_notes = self.include_notes_var.get() if self.license_manager.is_pro() else False

        def run():
            try:
                self._start_progress()
                self._update_status_safe(t('log_processing'))
                self._update_output_safe(f"\nğŸ“„ å‡¦ç†é–‹å§‹: {os.path.basename(path)}\n", clear=True)

                data, meta = self.extract_from_ppt(path, include_notes)
                if self.cancel_requested:
                    return self._log(t('log_cancelled'), "warning")

                if data:
                    # ãƒ•ã‚¡ã‚¤ãƒ«æƒ…å ±ã‚’æ›´æ–°
                    filename = os.path.basename(path)
                    slide_count = meta.get('slide_count', 0)
                    self.root.after(0, lambda: self._update_file_info(filename, len(data), slide_count))

                    # ã‚°ãƒªãƒƒãƒ‰ã«ãƒ­ãƒ¼ãƒ‰
                    self.extracted_data = data
                    self.root.after(0, lambda: self.grid_view.load_data(data))
                    self.root.after(0, lambda: self._show_edit_area())

                    # ãƒ•ã‚¡ã‚¤ãƒ«ä¿å­˜ï¼ˆãƒ‡ãƒ•ã‚©ãƒ«ãƒˆã¯Excelï¼‰
                    out = os.path.splitext(path)[0] + "_æŠ½å‡º.xlsx"
                    if self.save_to_file(data, out, "excel"):
                        self._log(f"âœ… {t('status_complete_items', len(data))} â†’ {os.path.basename(out)}", "success")
                        self._update_status_safe(t('status_complete_items', len(data)))
                else:
                    self._log(t('log_no_text'), "warning")
            except Exception as e:
                save_error_log(e, "_extract_single")
                self._log(t('log_error', e), "error")
            finally:
                self._stop_progress()

        threading.Thread(target=run, daemon=True).start()

    def _extract_batch(self, format: str = "excel"):
        """ãƒ•ã‚©ãƒ«ãƒ€ä¸€æ‹¬æŠ½å‡º (excel/json)"""
        if self.processing:
            return
        folder = filedialog.askdirectory(title=t('dialog_select_folder'))
        if not folder:
            return

        include_notes = self.include_notes_var.get() if self.license_manager.is_pro() else False
        ext = ".xlsx" if format == "excel" else ".json"

        def run():
            try:
                self._start_progress()
                self._update_output_safe(f"\nğŸ“ ãƒ•ã‚©ãƒ«ãƒ€ä¸€æ‹¬å‡ºåŠ› ({format.upper()}): {folder}\n", clear=True)

                files = [f for f in Path(folder).glob("*.pptx") if not f.name.startswith("~$")]
                if not files:
                    return self._log(t('log_no_pptx_found'), "warning")

                self._log(t('log_found_files', len(files)))
                total = 0

                for i, f in enumerate(files, 1):
                    if self.cancel_requested:
                        break
                    self._log(f"[{i}/{len(files)}] {f.name}")
                    data, meta = self.extract_from_ppt(str(f), include_notes)
                    if data:
                        out = str(f.with_suffix('')) + f"_æŠ½å‡º{ext}"
                        self.save_to_file(data, out, format)
                        total += len(data)

                self._log(f"âœ… {t('status_batch_complete', total, format.upper())}", "success")
            except Exception as e:
                self._log(t('log_error', e), "error")
            finally:
                self._stop_progress()

        threading.Thread(target=run, daemon=True).start()

    # === Update ===
    def _load_updates(self, path: str, source: str) -> Dict:
        updates = {}
        try:
            if source == "excel":
                wb = openpyxl.load_workbook(path)
                ws = wb.active
                headers = [c.value for c in ws[1]]
                try:
                    si = headers.index("ã‚¹ãƒ©ã‚¤ãƒ‰ç•ªå·") if "ã‚¹ãƒ©ã‚¤ãƒ‰ç•ªå·" in headers else headers.index("slide")
                    oi = headers.index("ã‚ªãƒ–ã‚¸ã‚§ã‚¯ãƒˆID") if "ã‚ªãƒ–ã‚¸ã‚§ã‚¯ãƒˆID" in headers else headers.index("id")
                    ti = headers.index("ãƒ†ã‚­ã‚¹ãƒˆå†…å®¹") if "ãƒ†ã‚­ã‚¹ãƒˆå†…å®¹" in headers else headers.index("text")
                except:
                    self._log(t('log_invalid_header'), "error")
                    return {}
                for row in list(ws.rows)[1:]:
                    try:
                        sn = int(row[si].value) if row[si].value else None
                        oid = str(row[oi].value) if row[oi].value else None
                        txt = str(row[ti].value) if row[ti].value else ""
                        if txt == "None":
                            txt = ""
                        if sn and oid:
                            updates[(sn, oid)] = txt
                    except:
                        pass
            elif source == "json":
                with open(path, 'r', encoding='utf-8') as f:
                    data = json.load(f)
                for item in data:
                    sn = item.get('ã‚¹ãƒ©ã‚¤ãƒ‰ç•ªå·') or item.get('slide')
                    oid = item.get('ã‚ªãƒ–ã‚¸ã‚§ã‚¯ãƒˆID') or item.get('id')
                    txt = item.get('ãƒ†ã‚­ã‚¹ãƒˆå†…å®¹') or item.get('text', '')
                    if sn and oid:
                        updates[(int(sn), str(oid))] = str(txt)
        except Exception as e:
            self._log(f"èª­ã¿è¾¼ã¿ã‚¨ãƒ©ãƒ¼: {e}", "error")
        return updates

    def _update_ppt(self, ppt_path: str, updates: Dict, preview: bool = False) -> Tuple[int, int, List]:
        limit = self.license_manager.get_update_limit()
        self.presentation = pptx.Presentation(ppt_path)
        updated, skipped = 0, 0
        changes = []

        for slide_idx, slide in enumerate(self.presentation.slides, 1):
            if self.cancel_requested:
                break
            if limit and slide_idx > limit:
                skipped += len([k for k in updates if k[0] == slide_idx])
                continue

            for shape in slide.shapes:
                try:
                    sid = str(shape.shape_id)
                    key = (slide_idx, sid)

                    if key in updates and hasattr(shape, "text"):
                        new_txt = updates[key]
                        old_txt = shape.text
                        if not self._texts_are_equal(old_txt, new_txt):
                            changes.append({'slide': slide_idx, 'id': sid, 'old': old_txt[:50], 'new': new_txt[:50]})
                            if not preview:
                                shape.text = self._normalize_for_compare(new_txt)
                                updated += 1

                    if hasattr(shape, "has_table") and shape.has_table:
                        for r, row in enumerate(shape.table.rows):
                            for c, cell in enumerate(row.cells):
                                ckey = (slide_idx, f"{sid}_t{r}_{c}")
                                if ckey in updates:
                                    new_txt = updates[ckey]
                                    old_txt = cell.text
                                    if not self._texts_are_equal(old_txt, new_txt):
                                        changes.append({'slide': slide_idx, 'id': ckey[1], 'old': old_txt[:30], 'new': new_txt[:30]})
                                        if not preview:
                                            cell.text = self._normalize_for_compare(new_txt)
                                            updated += 1
                except Exception as e:
                    skipped += 1

        return updated, skipped, changes

    def _run_update(self, source: str):
        if self.processing:
            return

        limit = self.license_manager.get_update_limit()
        if limit:
            if not messagebox.askyesno("ç¢ºèª", f"Freeç‰ˆã§ã¯æœ€åˆã®{limit}ã‚¹ãƒ©ã‚¤ãƒ‰ã®ã¿æ›´æ–°ã•ã‚Œã¾ã™ã€‚ç¶šè¡Œã—ã¾ã™ã‹ï¼Ÿ"):
                return

        ftypes = [("Excel", "*.xlsx")] if source == "excel" else [("JSON", "*.json")]
        data_path = filedialog.askopenfilename(title="ç·¨é›†æ¸ˆã¿ãƒ•ã‚¡ã‚¤ãƒ«ã‚’é¸æŠ", filetypes=ftypes)
        if not data_path:
            return
        ppt_path = filedialog.askopenfilename(title="PowerPointãƒ•ã‚¡ã‚¤ãƒ«ã‚’é¸æŠ", filetypes=[("PowerPoint", "*.pptx")])
        if not ppt_path:
            return

        def run():
            try:
                self._start_progress()
                self._update_output_safe(f"\nğŸ“¥ æ›´æ–°å‡¦ç†é–‹å§‹\n", clear=True)
                self._create_backup(ppt_path)

                updates = self._load_updates(data_path, source)
                if not updates:
                    return self._log(t('log_no_update_data'), "warning")

                self._log(f"èª­ã¿è¾¼ã¿: {len(updates)}ä»¶")
                updated, skipped, _ = self._update_ppt(ppt_path, updates)

                if self.cancel_requested:
                    return

                def save():
                    out = filedialog.asksaveasfilename(defaultextension=".pptx", filetypes=[("PowerPoint", "*.pptx")],
                                                       initialfile=os.path.splitext(os.path.basename(ppt_path))[0] + "_æ›´æ–°æ¸ˆã¿.pptx")
                    if out:
                        self.presentation.save(out)
                        self._log(f"âœ… ä¿å­˜å®Œäº†: {os.path.basename(out)}", "success")
                        messagebox.showinfo(t('dialog_complete'), t('result_updated', updated, skipped))

                self.root.after(0, save)
            except Exception as e:
                self._log(t('log_error', e), "error")
            finally:
                self._stop_progress()

        threading.Thread(target=run, daemon=True).start()

    def _update_excel(self):
        self._run_update("excel")

    def _update_json(self):
        self._run_update("json")

    def _update_batch(self, format: str = "excel"):
        """ãƒ•ã‚©ãƒ«ãƒ€å†…ã®Excel/JSONãƒ•ã‚¡ã‚¤ãƒ«ã¨PPTXã‚’ä¸€æ‹¬æ›´æ–°"""
        if self.processing:
            return

        ext = ".xlsx" if format == "excel" else ".json"
        folder = filedialog.askdirectory(title=t('dialog_select_folder_update', ext))
        if not folder:
            return

        def run():
            try:
                self._start_progress()
                self._update_output_safe(f"\nğŸ“ ãƒ•ã‚©ãƒ«ãƒ€ä¸€æ‹¬èª­è¾¼ ({format.upper()}): {folder}\n", clear=True)

                folder_path = Path(folder)

                # æŒ‡å®šå½¢å¼ã®ãƒ•ã‚¡ã‚¤ãƒ«ã‚’æ¤œç´¢
                data_files = list(folder_path.glob(f"*_æŠ½å‡º{ext}"))

                if not data_files:
                    return self._log(f"æŠ½å‡ºãƒ•ã‚¡ã‚¤ãƒ« (*_æŠ½å‡º{ext}) ãŒè¦‹ã¤ã‹ã‚Šã¾ã›ã‚“", "warning")

                self._log(t('log_found_files', len(data_files)))
                updated_count = 0
                error_count = 0

                for i, data_file in enumerate(data_files, 1):
                    if self.cancel_requested:
                        break

                    # å¯¾å¿œã™ã‚‹PPTXãƒ•ã‚¡ã‚¤ãƒ«ã‚’æ¤œç´¢
                    base_name = data_file.stem.replace("_æŠ½å‡º", "")
                    pptx_path = folder_path / f"{base_name}.pptx"

                    if not pptx_path.exists():
                        self._log(f"[{i}/{len(data_files)}] {data_file.name}: PPTXãªã— (ã‚¹ã‚­ãƒƒãƒ—)", "warning")
                        continue

                    self._log(f"[{i}/{len(data_files)}] {pptx_path.name}")

                    try:
                        updates = self._load_updates(str(data_file), format)

                        if not updates:
                            self._log(f"  â†’ æ›´æ–°ãƒ‡ãƒ¼ã‚¿ãªã—", "warning")
                            continue

                        self._create_backup(str(pptx_path))
                        updated, skipped, _ = self._update_ppt(str(pptx_path), updates)

                        # æ›´æ–°æ¸ˆã¿ãƒ•ã‚¡ã‚¤ãƒ«ã‚’ä¿å­˜
                        out_path = folder_path / f"{base_name}_æ›´æ–°æ¸ˆã¿.pptx"
                        self.presentation.save(str(out_path))
                        self._log(f"  â†’ {updated}ä»¶æ›´æ–°, ä¿å­˜: {out_path.name}")
                        updated_count += 1

                    except Exception as e:
                        self._log(f"  â†’ ã‚¨ãƒ©ãƒ¼: {e}", "error")
                        error_count += 1

                self._log(f"\nâœ… ãƒãƒƒãƒèª­è¾¼å®Œäº† ({format.upper()}): {updated_count}ä»¶æˆåŠŸ, {error_count}ä»¶ã‚¨ãƒ©ãƒ¼", "success")

            except Exception as e:
                self._log(t('log_error', e), "error")
            finally:
                self._stop_progress()

        threading.Thread(target=run, daemon=True).start()

    def _run_preview(self):
        data_path = filedialog.askopenfilename(title="ç·¨é›†æ¸ˆã¿ãƒ•ã‚¡ã‚¤ãƒ«ã‚’é¸æŠ", filetypes=[("Excel/TXT", "*.xlsx *.txt")])
        if not data_path:
            return
        ppt_path = filedialog.askopenfilename(title="PowerPointãƒ•ã‚¡ã‚¤ãƒ«ã‚’é¸æŠ", filetypes=[("PowerPoint", "*.pptx")])
        if not ppt_path:
            return

        def run():
            try:
                self._start_progress()
                self._update_output_safe(f"\nğŸ‘ å·®åˆ†ãƒ—ãƒ¬ãƒ“ãƒ¥ãƒ¼\n", clear=True)
                source = "excel" if data_path.endswith('.xlsx') else "json"
                updates = self._load_updates(data_path, source)
                if not updates:
                    return self._log(t('log_no_update_data'), "warning")

                _, _, changes = self._update_ppt(ppt_path, updates, preview=True)
                if changes:
                    self._log(f"\nå¤‰æ›´ç®‡æ‰€: {len(changes)}ä»¶")
                    for i, c in enumerate(changes[:20], 1):
                        self._update_output_safe(f"[{i}] ã‚¹ãƒ©ã‚¤ãƒ‰{c['slide']} ID:{c['id']}\n  æ—§: {c['old']}\n  æ–°: {c['new']}\n\n")
                else:
                    self._log("å¤‰æ›´ç®‡æ‰€ãªã—")
            except Exception as e:
                self._log(t('log_error', e), "error")
            finally:
                self._stop_progress()

        threading.Thread(target=run, daemon=True).start()

    # === æ¯”è¼ƒæ©Ÿèƒ½ ===
    def _show_compare_dialog(self):
        CompareDialog(self.root, self._run_compare)

    def _run_compare(self, file1: str, file2: str, ignore_ws: bool):
        def run():
            try:
                self._start_progress()
                self._update_output_safe(f"\nğŸ”€ æ¯”è¼ƒå‡¦ç†ä¸­...\n", clear=True)

                data1, _ = self.extract_from_ppt(file1)
                data2, _ = self.extract_from_ppt(file2)

                # ãƒãƒƒãƒ”ãƒ³ã‚°
                map1 = {(d["slide"], d["id"]): d["text"] for d in data1}
                map2 = {(d["slide"], d["id"]): d["text"] for d in data2}

                all_keys = set(map1.keys()) | set(map2.keys())
                diff_data = []
                stats = {"same": 0, "changed": 0, "added": 0, "removed": 0}

                for key in sorted(all_keys):
                    t1 = map1.get(key)
                    t2 = map2.get(key)

                    if t1 and t2:
                        if self._texts_are_equal(t1, t2):
                            status = "ä¸€è‡´"
                            stats["same"] += 1
                        else:
                            status = "å¤‰æ›´"
                            stats["changed"] += 1
                    elif t1:
                        status = "å‰Šé™¤"
                        stats["removed"] += 1
                    else:
                        status = "è¿½åŠ "
                        stats["added"] += 1

                    diff_data.append({
                        "slide": key[0], "id": key[1], "status": status,
                        "before": t1 or "", "after": t2 or ""
                    })

                self._log(f"æ¯”è¼ƒå®Œäº†: ä¸€è‡´{stats['same']} å¤‰æ›´{stats['changed']} è¿½åŠ {stats['added']} å‰Šé™¤{stats['removed']}")

                self.root.after(0, lambda: CompareResultWindow(
                    self.root, os.path.basename(file1), os.path.basename(file2),
                    diff_data, stats, on_apply=self._apply_compare_result
                ))
            except Exception as e:
                self._log(t('log_error', e), "error")
            finally:
                self._stop_progress()

        threading.Thread(target=run, daemon=True).start()

    def _apply_compare_result(self, selected_data: List[Dict]):
        # æ¯”è¼ƒçµæœã‚’ã‚°ãƒªãƒƒãƒ‰ã«åæ˜ 
        grid_data = []
        for item in selected_data:
            grid_data.append({
                "slide": item["slide"], "id": item.get("id", ""),
                "type": "", "text": item["text"]
            })
        self.grid_view.load_data(grid_data)
        self._show_edit_area()

    # === ã‚°ãƒªãƒƒãƒ‰æ“ä½œ ===
    def _on_grid_change(self, item, column, value):
        pass  # å¤‰æ›´æ™‚ã®è¿½åŠ å‡¦ç†ãŒã‚ã‚Œã°

    def _apply_grid_to_pptx(self):
        if not self.grid_view.get_data():
            messagebox.showwarning("è­¦å‘Š", "ã‚°ãƒªãƒƒãƒ‰ã«ãƒ‡ãƒ¼ã‚¿ãŒã‚ã‚Šã¾ã›ã‚“")
            return

        ppt_path = filedialog.askopenfilename(title="æ›´æ–°ã™ã‚‹PowerPointã‚’é¸æŠ", filetypes=[("PowerPoint", "*.pptx")])
        if not ppt_path:
            return

        # ã‚°ãƒªãƒƒãƒ‰ãƒ‡ãƒ¼ã‚¿ã‹ã‚‰æ›´æ–°è¾æ›¸ä½œæˆ
        grid_data = self.grid_view.get_data()
        updates = {}
        for row in grid_data:
            try:
                sn = int(row["slide"])
                oid = row["id"]
                txt = row["text"]
                if sn and oid:
                    updates[(sn, oid)] = txt
            except:
                pass

        if not updates:
            messagebox.showwarning("è­¦å‘Š", "æœ‰åŠ¹ãªæ›´æ–°ãƒ‡ãƒ¼ã‚¿ãŒã‚ã‚Šã¾ã›ã‚“")
            return

        def run():
            try:
                self._start_progress()
                self._log(f"ã‚°ãƒªãƒƒãƒ‰ã‹ã‚‰æ›´æ–°: {len(updates)}ä»¶")
                self._create_backup(ppt_path)

                updated, skipped, _ = self._update_ppt(ppt_path, updates)

                def save():
                    out = filedialog.asksaveasfilename(defaultextension=".pptx", filetypes=[("PowerPoint", "*.pptx")],
                                                       initialfile=os.path.splitext(os.path.basename(ppt_path))[0] + "_æ›´æ–°æ¸ˆã¿.pptx")
                    if out:
                        self.presentation.save(out)
                        self._log(f"âœ… ä¿å­˜å®Œäº†: {out}", "success")
                        messagebox.showinfo(t('dialog_complete'), t('status_update_complete', updated))

                self.root.after(0, save)
            except Exception as e:
                self._log(t('log_error', e), "error")
            finally:
                self._stop_progress()

        threading.Thread(target=run, daemon=True).start()

    def _export_grid_excel(self):
        data = self.grid_view.get_data()
        if not data:
            messagebox.showwarning("è­¦å‘Š", "ã‚¨ã‚¯ã‚¹ãƒãƒ¼ãƒˆã™ã‚‹ãƒ‡ãƒ¼ã‚¿ãŒã‚ã‚Šã¾ã›ã‚“")
            return

        path = filedialog.asksaveasfilename(defaultextension=".xlsx", filetypes=[("Excel", "*.xlsx")])
        if not path:
            return

        if self.save_to_file(data, path, "excel"):
            messagebox.showinfo(t('dialog_complete'), t('result_export_complete', path))

    def _export_grid_json(self):
        data = self.grid_view.get_data()
        if not data:
            messagebox.showwarning("è­¦å‘Š", "ã‚¨ã‚¯ã‚¹ãƒãƒ¼ãƒˆã™ã‚‹ãƒ‡ãƒ¼ã‚¿ãŒã‚ã‚Šã¾ã›ã‚“")
            return

        path = filedialog.asksaveasfilename(defaultextension=".json", filetypes=[("JSON", "*.json")])
        if not path:
            return

        if self.save_to_file(data, path, "json"):
            messagebox.showinfo(t('dialog_complete'), t('result_export_complete', path))

    def _load_excel_to_grid(self):
        """Excelãƒ•ã‚¡ã‚¤ãƒ«ã‚’ã‚°ãƒªãƒƒãƒ‰ã«ãƒ­ãƒ¼ãƒ‰ã™ã‚‹"""
        path = filedialog.askopenfilename(
            title=t('btn_load_excel'),
            filetypes=[("Excel", "*.xlsx")]
        )
        if not path:
            return

        try:
            wb = openpyxl.load_workbook(path)
            ws = wb.active
            headers = [c.value for c in ws[1]]

            # ãƒ˜ãƒƒãƒ€ãƒ¼åˆ—ã‚’ç‰¹å®š
            try:
                si = headers.index("ã‚¹ãƒ©ã‚¤ãƒ‰ç•ªå·") if "ã‚¹ãƒ©ã‚¤ãƒ‰ç•ªå·" in headers else headers.index("slide")
                oi = headers.index("ã‚ªãƒ–ã‚¸ã‚§ã‚¯ãƒˆID") if "ã‚ªãƒ–ã‚¸ã‚§ã‚¯ãƒˆID" in headers else headers.index("id")
                ti = headers.index("ãƒ†ã‚­ã‚¹ãƒˆå†…å®¹") if "ãƒ†ã‚­ã‚¹ãƒˆå†…å®¹" in headers else headers.index("text")
                type_i = headers.index("ã‚¿ã‚¤ãƒ—") if "ã‚¿ã‚¤ãƒ—" in headers else (headers.index("type") if "type" in headers else None)
            except ValueError:
                messagebox.showerror(t('dialog_error'), t('log_invalid_header'))
                return

            data = []
            for row in list(ws.rows)[1:]:
                try:
                    slide = str(row[si].value) if row[si].value else ""
                    oid = str(row[oi].value) if row[oi].value else ""
                    txt = str(row[ti].value) if row[ti].value else ""
                    obj_type = str(row[type_i].value) if type_i is not None and row[type_i].value else ""
                    if txt == "None":
                        txt = ""
                    if slide and oid:
                        data.append({
                            "slide": slide,
                            "id": oid,
                            "type": obj_type,
                            "text": txt
                        })
                except Exception:
                    pass

            if data:
                self.extracted_data = data
                self.grid_view.load_data(data)
                self._show_edit_area()
                self._update_file_info(os.path.basename(path), len(data))
                self._log(f"âœ… {t('status_complete_items', len(data))}", "success")
            else:
                messagebox.showwarning(t('dialog_error'), t('log_no_text'))

        except Exception as e:
            save_error_log(e, "_load_excel_to_grid")
            messagebox.showerror(t('dialog_error'), str(e))

    def _load_json_to_grid(self):
        """JSONãƒ•ã‚¡ã‚¤ãƒ«ã‚’ã‚°ãƒªãƒƒãƒ‰ã«ãƒ­ãƒ¼ãƒ‰ã™ã‚‹"""
        if not self.license_manager.can_json():
            messagebox.showinfo(t('dialog_error'), "JSONæ©Ÿèƒ½ã¯Proç‰ˆä»¥ä¸ŠãŒå¿…è¦ã§ã™")
            return

        path = filedialog.askopenfilename(
            title=t('btn_load_json'),
            filetypes=[("JSON", "*.json")]
        )
        if not path:
            return

        try:
            with open(path, 'r', encoding='utf-8') as f:
                json_data = json.load(f)

            data = []
            for item in json_data:
                slide = str(item.get('ã‚¹ãƒ©ã‚¤ãƒ‰ç•ªå·') or item.get('slide', ''))
                oid = str(item.get('ã‚ªãƒ–ã‚¸ã‚§ã‚¯ãƒˆID') or item.get('id', ''))
                txt = str(item.get('ãƒ†ã‚­ã‚¹ãƒˆå†…å®¹') or item.get('text', ''))
                obj_type = str(item.get('ã‚¿ã‚¤ãƒ—') or item.get('type', ''))
                if slide and oid:
                    data.append({
                        "slide": slide,
                        "id": oid,
                        "type": obj_type,
                        "text": txt
                    })

            if data:
                self.extracted_data = data
                self.grid_view.load_data(data)
                self._show_edit_area()
                self._update_file_info(os.path.basename(path), len(data))
                self._log(f"âœ… {t('status_complete_items', len(data))}", "success")
            else:
                messagebox.showwarning(t('dialog_error'), t('log_no_text'))

        except Exception as e:
            save_error_log(e, "_load_json_to_grid")
            messagebox.showerror(t('dialog_error'), str(e))

    def _batch_extract_dialog(self):
        """ãƒ•ã‚©ãƒ«ãƒ€ä¸€æ‹¬æŠ½å‡ºï¼ˆExcelå½¢å¼ï¼‰"""
        self._extract_batch("excel")

    def _batch_update_dialog(self):
        """ãƒ•ã‚©ãƒ«ãƒ€ä¸€æ‹¬æ›´æ–°ï¼ˆExcelå½¢å¼ï¼‰"""
        self._update_batch("excel")

    # === Dialogs ===
    def _check_license_on_startup(self):
        """èµ·å‹•æ™‚ã®ãƒ©ã‚¤ã‚»ãƒ³ã‚¹ãƒã‚§ãƒƒã‚¯"""
        # ãƒ©ã‚¤ã‚»ãƒ³ã‚¹ãŒæœªã‚¢ã‚¯ãƒ†ã‚£ãƒ™ãƒ¼ãƒˆã®å ´åˆã€èªè¨¼ãƒ€ã‚¤ã‚¢ãƒ­ã‚°ã‚’è¡¨ç¤º
        if not self.license_manager.is_activated():
            self._show_license_dialog(startup_check=True)
            return

        # æœ‰åŠ¹æœŸé™åˆ‡ã‚Œãƒã‚§ãƒƒã‚¯
        days = self.license_manager.get_days_until_expiry()
        if days is not None and days <= 0:
            # æœŸé™åˆ‡ã‚Œ - Freeç‰ˆã«ãƒ€ã‚¦ãƒ³ã‚°ãƒ¬ãƒ¼ãƒ‰
            messagebox.showwarning(
                t('dialog_error'),
                t('license_expired')
            )
            self.license_manager.deactivate()
            self._create_layout()
            self._show_license_dialog(startup_check=True)
            return

        # æœŸé™åˆ‡ã‚Œè­¦å‘Šï¼ˆ30æ—¥ä»¥å†…ï¼‰
        if self.license_manager.should_show_expiry_warning():
            expiry_str = self.license_manager.get_expiry_date_str()
            messagebox.showinfo(
                t('license_title'),
                t('license_expiry_warning', days, expiry_str)
            )

    def _show_license_dialog(self, startup_check: bool = False):
        """ãƒ©ã‚¤ã‚»ãƒ³ã‚¹èªè¨¼ãƒ€ã‚¤ã‚¢ãƒ­ã‚°ã‚’è¡¨ç¤º

        Args:
            startup_check: èµ·å‹•æ™‚ãƒã‚§ãƒƒã‚¯ã®å ´åˆTrueï¼ˆã‚­ãƒ£ãƒ³ã‚»ãƒ«ä¸å¯ï¼‰
        """
        dialog = tk.Toplevel(self.root)
        dialog.title(t('license_auth_title'))
        dialog.geometry("550x520")
        dialog.minsize(550, 520)
        dialog.transient(self.root)
        dialog.grab_set()

        if startup_check:
            dialog.protocol("WM_DELETE_WINDOW", lambda: None)  # é–‰ã˜ã‚‹ãƒœã‚¿ãƒ³ç„¡åŠ¹

        frame = ttk.Frame(dialog, padding=20)
        frame.pack(fill='both', expand=True)

        # ãƒ˜ãƒƒãƒ€ãƒ¼
        header_frame = ttk.Frame(frame)
        header_frame.pack(fill='x', pady=(0, 15))
        ttk.Label(header_frame, text=APP_NAME, font=FONTS["heading"]).pack(side='left')

        tier = self.license_manager.get_tier_info()
        tier_name = tier['name_ja'] if get_language() == 'ja' else tier['name']

        # ç¾åœ¨ã®ã‚¹ãƒ†ãƒ¼ã‚¿ã‚¹è¡¨ç¤º
        status_frame = ttk.Frame(frame, style="Card.TFrame")
        status_frame.pack(fill='x', pady=(0, 15))

        ttk.Label(status_frame, text=t('license_current'), font=FONTS["body_bold"]).pack(anchor='w', pady=(10, 5), padx=10)

        status_inner = ttk.Frame(status_frame)
        status_inner.pack(fill='x', padx=10, pady=(0, 10))

        badge_text = f"{tier['badge']}"
        ttk.Label(status_inner, text=badge_text, font=FONTS["body_bold"],
                  foreground=COLOR_PALETTE["brand_primary"]).pack(side='left')

        if self.license_manager.is_activated():
            # æœ‰åŠ¹æœŸé™è¡¨ç¤º
            expiry_str = self.license_manager.get_expiry_date_str()
            days = self.license_manager.get_days_until_expiry()

            if days is not None:
                if days > 0:
                    status_text = f"  |  {t('license_valid_until', expiry_str)} {t('license_days_remaining', days)}"
                    status_color = COLOR_PALETTE["success"] if days > 30 else COLOR_PALETTE["warning"]
                else:
                    status_text = f"  |  {t('license_status_expired')}"
                    status_color = COLOR_PALETTE["error"]
            else:
                status_text = f"  |  {t('license_perpetual')}"
                status_color = COLOR_PALETTE["success"]

            ttk.Label(status_inner, text=status_text, font=FONTS["small"],
                      foreground=status_color).pack(side='left')

        ttk.Separator(frame, orient='horizontal').pack(fill='x', pady=10)

        # ãƒ•ã‚©ãƒ¼ãƒãƒƒãƒˆèª¬æ˜
        ttk.Label(frame, text="å½¢å¼: INS-SLIDE-{TIER}-XXXX-XXXX-CC", font=FONTS["small"],
                  foreground=COLOR_PALETTE["text_muted"]).pack(anchor='w', pady=(0, 10))

        # ãƒ¡ãƒ¼ãƒ«ã‚¢ãƒ‰ãƒ¬ã‚¹å…¥åŠ›
        ttk.Label(frame, text=t('license_email'), font=FONTS["body"]).pack(anchor='w')
        email_var = tk.StringVar(value=self.license_manager.license_info.get('email', ''))
        email_entry = ttk.Entry(frame, textvariable=email_var, width=45, font=FONTS["body"])
        email_entry.pack(fill='x', pady=(5, 10))

        # ãƒ©ã‚¤ã‚»ãƒ³ã‚¹ã‚­ãƒ¼å…¥åŠ›
        ttk.Label(frame, text=t('license_key'), font=FONTS["body"]).pack(anchor='w')
        key_var = tk.StringVar(value=self.license_manager.license_info.get('key', ''))
        key_entry = ttk.Entry(frame, textvariable=key_var, width=45, font=FONTS["body"])
        key_entry.pack(fill='x', pady=(5, 10))

        # ã‚¨ãƒ©ãƒ¼ãƒ¡ãƒƒã‚»ãƒ¼ã‚¸è¡¨ç¤ºç”¨
        error_var = tk.StringVar()
        error_label = ttk.Label(frame, textvariable=error_var, font=FONTS["small"],
                                foreground=COLOR_PALETTE["error"])
        error_label.pack(anchor='w', pady=(0, 10))

        def activate():
            email = email_var.get().strip()
            key = key_var.get().strip()

            if not email:
                error_var.set(t('license_email_required'))
                return
            if not key:
                error_var.set(t('license_enter_prompt'))
                return

            ok, msg = self.license_manager.activate(email, key)
            if ok:
                messagebox.showinfo(t('dialog_complete'), msg)
                dialog.destroy()
                self._create_layout()
            else:
                error_var.set(msg)

        def deactivate():
            self.license_manager.deactivate()
            messagebox.showinfo(t('dialog_complete'), t('license_deactivated'))
            dialog.destroy()
            self._create_layout()

        def skip_free():
            """Freeç‰ˆã¨ã—ã¦ç¶šè¡Œ"""
            dialog.destroy()

        # ãƒœã‚¿ãƒ³ãƒ•ãƒ¬ãƒ¼ãƒ 
        btn_frame = ttk.Frame(frame)
        btn_frame.pack(fill='x', pady=(10, 0))

        if self.license_manager.is_activated():
            ttk.Button(btn_frame, text=t('btn_deactivate'), command=deactivate).pack(side='left')

        ttk.Button(btn_frame, text=t('btn_activate'), command=activate, style="Accent.TButton").pack(side='left', padx=5)

        if not startup_check:
            ttk.Button(btn_frame, text=t('btn_close'), command=dialog.destroy).pack(side='right')
        else:
            # èµ·å‹•æ™‚ã¯Freeç‰ˆã¨ã—ã¦ç¶šè¡Œå¯èƒ½
            ttk.Button(btn_frame, text=t('license_continue_free'), command=skip_free).pack(side='right')

        # ãƒªãƒ³ã‚¯ãƒ•ãƒ¬ãƒ¼ãƒ 
        link_frame = ttk.Frame(frame)
        link_frame.pack(fill='x', pady=(20, 0))

        def open_trial():
            webbrowser.open(SUPPORT_LINKS.get('contact', ''))

        def open_purchase():
            webbrowser.open(SUPPORT_LINKS.get('purchase', ''))

        trial_link = ttk.Label(link_frame, text=t('license_trial_link'), font=FONTS["small"],
                               foreground=COLOR_PALETTE["brand_primary"], cursor="hand2")
        trial_link.pack(side='left')
        trial_link.bind("<Button-1>", lambda e: open_trial())

        ttk.Label(link_frame, text="  |  ", font=FONTS["small"],
                  foreground=COLOR_PALETTE["text_muted"]).pack(side='left')

        purchase_link = ttk.Label(link_frame, text=t('btn_purchase'), font=FONTS["small"],
                                  foreground=COLOR_PALETTE["brand_primary"], cursor="hand2")
        purchase_link.pack(side='left')
        purchase_link.bind("<Button-1>", lambda e: open_purchase())

    def _show_about(self):
        tier = self.license_manager.get_tier_info()
        messagebox.showinfo(t('menu_about'),
            f"{APP_NAME} v{APP_VERSION}\n\n"
            f"ãƒ©ã‚¤ã‚»ãƒ³ã‚¹: {tier['name']}\n\n"
            f"çµ±ä¸€ãƒ©ã‚¤ã‚»ãƒ³ã‚¹å½¢å¼:\n"
            f"INS-SLIDE-{{TIER}}-XXXX-XXXX-CC\n\n"
            f"by Harmonic Insight\nÂ© 2025"
        )

    def _on_closing(self):
        if self.processing:
            if not messagebox.askokcancel(t('dialog_confirm_title'), t('dialog_processing_exit')):
                return
        self.root.destroy()


def main():
    try:
        import ctypes
        ctypes.windll.shcore.SetProcessDpiAwareness(1)
    except:
        pass
    root = tk.Tk()
    InsightSlidesApp(root)
    root.mainloop()


if __name__ == "__main__":
    main()
