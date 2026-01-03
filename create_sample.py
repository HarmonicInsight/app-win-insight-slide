"""
テスト用のサンプルPPTXを生成
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor


def create_sample_pptx(output_path: str = "sample.pptx"):
    """サンプルPPTXを作成"""
    prs = Presentation()
    
    # スライド1: タイトル
    slide_layout = prs.slide_layouts[6]  # 空白
    slide1 = prs.slides.add_slide(slide_layout)
    
    title1 = slide1.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
    title1.name = "Title"
    tf = title1.text_frame
    tf.text = "会社概要"
    tf.paragraphs[0].font.size = Pt(44)
    tf.paragraphs[0].font.bold = True
    
    body1 = slide1.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(3))
    body1.name = "Body"
    tf = body1.text_frame
    tf.text = "当社は建設DXを推進する企業です。\n28年の経験と700社以上の導入実績があります。"
    tf.paragraphs[0].font.size = Pt(24)
    
    # スライド2: 事業内容
    slide2 = prs.slides.add_slide(slide_layout)
    
    title2 = slide2.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
    title2.name = "Title"
    tf = title2.text_frame
    tf.text = "事業内容"
    tf.paragraphs[0].font.size = Pt(44)
    tf.paragraphs[0].font.bold = True
    
    body2 = slide2.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(4))
    body2.name = "Body"
    tf = body2.text_frame
    tf.text = "・建設業向けDXコンサルティング\n・業務システム導入支援\n・RPA導入・運用サポート\n・BIM/CIM活用支援"
    tf.paragraphs[0].font.size = Pt(20)
    
    # スライド3: 実績
    slide3 = prs.slides.add_slide(slide_layout)
    
    title3 = slide3.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
    title3.name = "Title"
    tf = title3.text_frame
    tf.text = "導入実績"
    tf.paragraphs[0].font.size = Pt(44)
    tf.paragraphs[0].font.bold = True
    
    body3 = slide3.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(4))
    body3.name = "Body"
    tf = body3.text_frame
    tf.text = "大手ゼネコンから中小建設会社まで\n幅広い規模の企業様に導入いただいています。\n\n新国立競技場プロジェクトにも参画。"
    tf.paragraphs[0].font.size = Pt(20)
    
    prs.save(output_path)
    print(f"サンプルPPTXを作成しました: {output_path}")
    return output_path


if __name__ == "__main__":
    create_sample_pptx()
