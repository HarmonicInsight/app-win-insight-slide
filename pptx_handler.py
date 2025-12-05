"""
PPTX Handler - PowerPointファイルのJSON抽出・反映
"""
import json
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt


def extract_to_json(pptx_path: str) -> dict:
    """PPTXファイルからテキストを抽出してJSON構造で返す"""
    prs = Presentation(pptx_path)
    
    data = {
        "file": str(Path(pptx_path).name),
        "slides": []
    }
    
    for slide_idx, slide in enumerate(prs.slides, start=1):
        slide_data = {
            "slide": slide_idx,
            "shapes": []
        }
        
        for shape in slide.shapes:
            if shape.has_text_frame:
                # テキストフレーム内の全テキストを結合
                text_parts = []
                for paragraph in shape.text_frame.paragraphs:
                    para_text = "".join(run.text for run in paragraph.runs)
                    if para_text:
                        text_parts.append(para_text)
                
                full_text = "\n".join(text_parts)
                
                if full_text.strip():  # 空でないテキストのみ追加
                    shape_data = {
                        "id": shape.shape_id,
                        "name": shape.name,
                        "text": full_text,
                        "original": full_text  # 元テキストを保持
                    }
                    slide_data["shapes"].append(shape_data)
        
        if slide_data["shapes"]:  # テキストがあるスライドのみ追加
            data["slides"].append(slide_data)
    
    return data


def apply_from_json(pptx_path: str, json_data: dict, output_path: str = None) -> str:
    """JSON構造からPPTXファイルにテキストを反映"""
    prs = Presentation(pptx_path)
    
    # JSONデータをshape_idでインデックス化
    text_map = {}
    for slide_data in json_data.get("slides", []):
        slide_idx = slide_data["slide"]
        for shape_data in slide_data.get("shapes", []):
            key = (slide_idx, shape_data["id"])
            text_map[key] = shape_data["text"]
    
    # PPTXに反映
    for slide_idx, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if shape.has_text_frame:
                key = (slide_idx, shape.shape_id)
                if key in text_map:
                    new_text = text_map[key]
                    lines = new_text.split("\n")
                    
                    tf = shape.text_frame
                    
                    # 最初のパラグラフのフォーマットを保存
                    first_para = tf.paragraphs[0]
                    first_font = None
                    if first_para.runs:
                        first_run = first_para.runs[0]
                        first_font = {
                            'size': first_run.font.size,
                            'bold': first_run.font.bold,
                            'italic': first_run.font.italic,
                            'name': first_run.font.name,
                        }
                    
                    # 全パラグラフをクリア（最初のパラグラフは残す）
                    # python-pptxでは最初のパラグラフは削除できないので、テキストをクリア
                    for para in tf.paragraphs:
                        for run in para.runs:
                            run.text = ""
                    
                    # 最初の行を最初のパラグラフに設定
                    if lines:
                        if first_para.runs:
                            first_para.runs[0].text = lines[0]
                        else:
                            run = first_para.add_run()
                            run.text = lines[0]
                            if first_font and first_font['size']:
                                run.font.size = first_font['size']
                    
                    # 2行目以降を追加
                    for i, line in enumerate(lines[1:], start=1):
                        if i < len(tf.paragraphs):
                            # 既存のパラグラフがあれば使用
                            para = tf.paragraphs[i]
                            if para.runs:
                                para.runs[0].text = line
                            else:
                                run = para.add_run()
                                run.text = line
                        else:
                            # 新しいパラグラフを追加
                            para = tf.add_paragraph()
                            run = para.add_run()
                            run.text = line
                            # フォーマットをコピー
                            if first_font:
                                if first_font['size']:
                                    run.font.size = first_font['size']
                                if first_font['bold']:
                                    run.font.bold = first_font['bold']
                                if first_font['name']:
                                    run.font.name = first_font['name']
    
    # 出力パス設定
    if output_path is None:
        p = Path(pptx_path)
        output_path = str(p.parent / f"{p.stem}_edited{p.suffix}")
    
    prs.save(output_path)
    return output_path


def save_json(data: dict, json_path: str):
    """JSONをファイルに保存"""
    with open(json_path, "w", encoding="utf-8") as f:
        json.dump(data, f, ensure_ascii=False, indent=2)


def load_json(json_path: str) -> dict:
    """JSONファイルを読み込み"""
    with open(json_path, "r", encoding="utf-8") as f:
        return json.load(f)


# テスト用
if __name__ == "__main__":
    import sys
    if len(sys.argv) > 1:
        pptx_file = sys.argv[1]
        data = extract_to_json(pptx_file)
        print(json.dumps(data, ensure_ascii=False, indent=2))
