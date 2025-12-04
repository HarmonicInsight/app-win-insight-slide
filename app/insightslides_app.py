"""
Insight Slides - PowerPoint Text Extract & Update Tool v2.2
Freemium Model: Free / Standard / Professional
Bilingual Support: English / Japanese

by Harmonic Insight

Feature Matrix:
┌─────────────────┬───────┬──────────┬──────────────┐
│ Feature         │ Free  │ Standard │ Professional │
├─────────────────┼───────┼──────────┼──────────────┤
│ Extract         │Single │ Unlimited│ Unlimited    │
│ Update          │3 slides│Unlimited│ Unlimited+Batch│
│ Format Quality  │ Basic │ Advanced │ Advanced     │
│ Diff Preview    │  ×   │   ×     │  ○          │
│ Auto Backup     │  ×   │   ×     │  ○          │
└─────────────────┴───────┴──────────┴──────────────┘
"""
import tkinter as tk
from tkinter import ttk, filedialog, messagebox, scrolledtext, simpledialog
import pptx
import openpyxl
from openpyxl.styles import Font as XLFont, PatternFill
import os
import re
import csv
import json
import hashlib
import base64
import webbrowser
import traceback
from datetime import datetime, timedelta
import threading
from pathlib import Path
from typing import Dict, Tuple, List, Optional
import shutil

# ライセンス共通モジュールからインポート
from licensing.core import (
    LICENSE_SECRET,
    LICENSE_TIERS,
    validate_key,
    get_tier_from_key,
    get_expiration_from_key,
    generate_checksum,
)

# ============== App Info ==============
APP_VERSION = "1.0.0"
APP_NAME = "Insight Slides"

# ============== Config Paths ==============
CONFIG_DIR = Path.home() / ".insightslides"
CONFIG_FILE = CONFIG_DIR / "config.json"
LICENSE_FILE = CONFIG_DIR / "license.key"
ERROR_LOG_FILE = CONFIG_DIR / "error_log.txt"

# ============== Support Links ==============
SUPPORT_LINKS = {
    "faq": "https://example.com/insightslides/faq",
    "tutorial": "https://example.com/insightslides/tutorial",
    "purchase": "https://example.com/insightslides/purchase",
    "contact": "mailto:support@example.com",
}

# ============== Internationalization (i18n) ==============
LANGUAGES = {
    'en': {
        # App
        'app_subtitle': 'Extract → AI Edit → Batch Update PowerPoint Text',
        'welcome_title': 'Welcome to Insight Slides!',
        
        # Modes
        'mode_extract': 'Extract Mode',
        'mode_update': 'Update Mode',
        'mode_extract_short': 'Extract',
        'mode_update_short': 'Update',
        
        # Panels
        'panel_mode': 'Mode Selection',
        'panel_file': 'File Operations',
        'panel_settings': 'Settings',
        'panel_status': 'Status',
        'panel_output': 'Output',
        'panel_extract_options': 'Extract Options',
        'panel_update_options': 'Update Options',
        'panel_extract_run': 'Run Extract',
        'panel_update_run': 'Run Update',
        'panel_analysis_tools': 'Analysis Tools',
        'panel_pro_features': 'Pro Features',
        
        # Buttons
        'btn_single_file': 'Select File',
        'btn_batch_folder': 'Folder Batch',
        'btn_from_excel': 'From Excel',
        'btn_from_json': 'From JSON',
        'btn_batch_update': 'Folder Batch Update',
        'btn_diff_preview': 'Diff Preview',
        'diff_preview_desc': 'Check changes before updating',
        'btn_cancel': 'Stop',
        'btn_clear': 'Clear Log',
        'btn_copy': 'Copy Log',
        'btn_license': 'License',
        'btn_activate': 'Activate',
        'btn_deactivate': 'Deactivate',
        'btn_purchase': 'Purchase',
        'btn_close': 'Close',
        'btn_start': 'Get Started',
        'btn_continue': 'Continue',
        
        # Settings
        'setting_output_format': 'Output Format:',
        'setting_include_meta': 'Include file name & date',
        'setting_auto_backup': 'Auto backup before update',
        'setting_auto_backup_pro': 'Auto backup (Pro)',
        'format_tab': 'Tab-separated',
        'format_csv': 'CSV',
        'format_excel': 'Excel',
        
        # Status
        'status_waiting': 'Waiting...',
        'status_processing': 'Processing...',
        'status_complete': 'Complete',
        'status_cancelled': 'Cancelled',
        'status_error': 'Error',
        
        # Messages
        'msg_extract_desc': 'Extract text from PowerPoint files.',
        'msg_update_desc': 'Apply edited text back to PowerPoint.',
        'msg_update_limit': 'Update: First {0} slides only\nUpgrade to Standard for unlimited!',
        'msg_batch_standard': '※ Batch extract available in Standard+',
        'msg_batch_pro': '※ Batch update, diff preview, auto backup in Pro',
        'msg_processing_file': 'Processing: {0}',
        'msg_saved': 'Saved: {0}',
        'msg_extracted': 'Extracted: {0} items from {1} slides',
        'msg_updated': 'Updated: {0} items, Skipped: {1}',
        'msg_no_pptx': 'No PPTX files found',
        'msg_no_data': 'No update data found',
        'msg_cancel_requested': 'Cancellation requested...',
        'msg_backup_created': 'Backup created: {0}',
        'msg_backup_failed': 'Backup failed: {0}',
        'msg_free_limit': 'Free version: Only first {0} slides will be updated.\n\nContinue?\n\nUpgrade to Standard (¥2,980/year) for unlimited!',
        'msg_copied': 'Copied to clipboard',
        'msg_support_copied': 'Support log copied to clipboard.\nPaste when contacting support.',
        
        # License
        'license_title': 'License Management',
        'license_current': 'Current License',
        'license_compare': 'Plan Comparison:',
        'license_enter_key': 'Enter License Key:',
        'license_activated': '{0} has been activated',
        'license_deactivated': 'License deactivated',
        'license_invalid': 'Invalid license key',
        'license_enter_prompt': 'Please enter a license key',
        'license_confirm_deactivate': 'Deactivate license?',
        'license_free_features': ['Extract: Single file only', 'Update: Up to 3 slides'],
        'license_standard_features': ['Extract: Unlimited + Batch', 'Update: Unlimited'],
        'license_pro_features': ['All Standard features', 'Batch update', 'Diff preview', 'Auto backup'],
        'license_key_label': 'License Key',
        'license_expires_label': 'Expires',
        'license_expired': 'Expired',
        'license_days_left': '{0} days left',
        'license_permanent': 'Permanent (No expiration)',
        
        # Upgrade
        'upgrade_title': 'Upgrade',
        'upgrade_to_standard': 'This feature requires Standard or higher.\n\n📘 Standard (¥2,980/year)\n  - Unlimited updates\n  - Batch extraction\n\n⭐ Pro (¥5,980/year)\n  - All Standard features\n  - Batch update\n  - Diff preview\n  - Auto backup\n\nPurchase license?',
        'upgrade_to_pro': 'This feature requires Pro.\n\n⭐ Pro (¥5,980/year)\n  - Batch update\n  - Diff preview\n  - Auto backup\n\nUpgrade?',
        
        # Welcome dialog
        'welcome_step1_title': '1. Extract',
        'welcome_step1_desc': 'Extract text from PowerPoint to Excel/TXT',
        'welcome_step2_title': '2. Edit',
        'welcome_step2_desc': 'Translate or edit in Excel/TXT',
        'welcome_step3_title': '3. Update',
        'welcome_step3_desc': 'Apply changes back to PowerPoint (with formatting)',
        'welcome_tip_title': '💡 Tips',
        'welcome_tip1': '• Access guides and FAQ from the Help menu',
        'welcome_tip2': '• Use "Copy Support Log" when reporting issues',
        
        # Language
        'lang_select_title': 'Select Language',
        'lang_restart_msg': 'Language changed.\nPlease restart the app to fully apply.',
        'lang_changed': 'Language changed.',
        'lang_menu': 'Language',
        'font_size_menu': 'Font Size',
        'font_size_small': 'Small',
        'font_size_medium': 'Medium',
        'font_size_large': 'Large',
        'font_size_changed': 'Font size changed.',
        'font_size_restart_msg': 'Font size changed.\nPlease restart the app to fully apply.',
        
        # Menu
        'menu_help': 'Help',
        'menu_guide': 'User Guide',
        'menu_faq': 'FAQ',
        'menu_support_log': 'Copy Support Log',
        'menu_license': 'License Management',
        'menu_purchase': 'Purchase License',
        'menu_about': 'About',
        
        # About
        'about_title': 'About',
        'about_text': '{0} v{1}\n\nLicense: {2}\n\nA tool to efficiently extract and update\nPowerPoint text content.\n\nby Harmonic Insight\n© 2025',
        'info_workflow_title': 'Basic Workflow',
        'info_workflow_1': '1. Extract: PowerPoint → Excel / Tab text / JSON',
        'info_workflow_2': '2. Edit: Manual or AI processing (typo fixes, translation, terminology, etc.)',
        'info_workflow_3': '3. Update: Excel / JSON → PowerPoint',
        'info_use_cases': 'Use Cases: Translation, Batch Editing, AI Integration, Font Check',
        'info_format_note': 'Formatting (font, color, size) is automatically preserved',
        'info_caution_title': 'Notes',
        'info_caution_1': 'Multiple fonts in one text box will be unified to the primary font',
        'info_caution_2': 'Text in images or chart labels cannot be extracted',
        'info_caution_3': 'Do not modify Slide Number or Object ID when editing',
        
        # Dialogs
        'dialog_confirm': 'Confirm',
        'dialog_error': 'Error',
        'dialog_complete': 'Complete',
        'dialog_processing': 'Processing',
        'dialog_select_pptx': 'Select PowerPoint File',
        'dialog_select_folder': 'Select Folder',
        'dialog_select_data': 'Select Edited File',
        'dialog_save_pptx': 'Save Updated PowerPoint',
        
        # Headers
        'header_slide': 'Slide',
        'header_id': 'Object ID',
        'header_type': 'Type',
        'header_text': 'Text Content',
        'header_filename': 'Filename',
        'header_datetime': 'Extracted At',
        
        # Diff preview
        'diff_title': 'Diff Preview',
        'diff_changes': 'Changes: {0} items',
        'diff_no_changes': 'No changes detected',
        'diff_old': 'Old',
        'diff_new': 'New',
        
        # Pro features
        'pro_section_label': 'Pro Features',
        'btn_font_analysis': 'Font Analysis',
        'btn_conditional_extract': 'Filter Extract',
        'chk_include_notes': 'Include Speaker Notes',
        'type_notes': 'Notes',
        'font_analysis_title': 'Font Analysis Report',
        'font_summary_title': 'Font Usage Summary',
        'font_warning_mixed': '⚠️ Warning: Multiple fonts detected',
        'header_font': 'Font',
        'header_size': 'Size',
        'header_bold': 'Bold',
        'conditional_title': 'Filter Extract',
        'conditional_prompt': 'Enter filter keyword or regex:',
        'conditional_no_match': 'No matching text found',
    },
    'ja': {
        # App
        'app_subtitle': 'PowerPointテキストを抽出 → AI編集 → 一括反映',
        'welcome_title': 'Insight Slides へようこそ！',
        
        # Modes
        'mode_extract': '抽出モード',
        'mode_update': '更新モード',
        'mode_extract_short': '抽出',
        'mode_update_short': '更新',
        
        # Panels
        'panel_mode': 'モード選択',
        'panel_file': 'ファイル操作',
        'panel_settings': '処理設定',
        'panel_status': '処理状況',
        'panel_output': '処理結果',
        'panel_extract_options': '抽出オプション',
        'panel_update_options': '更新オプション',
        'panel_extract_run': '抽出実行',
        'panel_update_run': '更新実行',
        'panel_analysis_tools': '分析ツール',
        'panel_pro_features': '拡張機能',
        
        # Buttons
        'btn_single_file': 'ファイル選択',
        'btn_batch_folder': 'フォルダー一括',
        'btn_from_excel': 'Excelから更新',
        'btn_from_json': 'JSONから更新',
        'btn_batch_update': 'フォルダー一括更新',
        'btn_diff_preview': '差分プレビュー',
        'diff_preview_desc': '更新前に変更箇所を確認',
        'btn_cancel': '中止',
        'btn_clear': 'ログクリア',
        'btn_copy': 'ログコピー',
        'btn_license': 'ライセンス',
        'btn_activate': 'アクティベート',
        'btn_deactivate': 'ライセンス解除',
        'btn_purchase': '購入ページ',
        'btn_close': '閉じる',
        'btn_start': '始める',
        'btn_continue': '続ける',
        
        # Settings
        'setting_output_format': '出力形式:',
        'setting_include_meta': 'ファイル名・日時を含める',
        'setting_auto_backup': '更新前に自動バックアップ',
        'setting_auto_backup_pro': '自動バックアップ (Pro)',
        'format_tab': 'タブ区切り',
        'format_csv': 'CSV形式',
        'format_excel': 'Excel形式',
        
        # Status
        'status_waiting': '処理待機中...',
        'status_processing': '処理中...',
        'status_complete': '完了',
        'status_cancelled': 'キャンセルされました',
        'status_error': 'エラー',
        
        # Messages
        'msg_extract_desc': 'PowerPointからテキストを抽出します。',
        'msg_update_desc': '編集したファイルの変更をPowerPointに反映します。',
        'msg_update_limit': '更新機能: 最初の{0}スライドのみ\nStandard版で無制限に！',
        'msg_batch_standard': '※ Standard以上でフォルダ一括処理が可能',
        'msg_batch_pro': '※ Pro版でバッチ更新・差分プレビュー・自動バックアップ',
        'msg_processing_file': '処理中: {0}',
        'msg_saved': '保存完了: {0}',
        'msg_extracted': '抽出: {0}件 / スライド: {1}枚',
        'msg_updated': '更新: {0}件 / スキップ: {1}件',
        'msg_no_pptx': 'PPTXファイルが見つかりません',
        'msg_no_data': '更新データがありません',
        'msg_cancel_requested': 'キャンセルをリクエスト...',
        'msg_backup_created': 'バックアップ作成: {0}',
        'msg_backup_failed': 'バックアップ失敗: {0}',
        'msg_free_limit': 'Free版では最初の{0}スライドのみ更新されます。\n\n続行しますか？\n\nStandard版(¥2,980/年)で無制限更新！',
        'msg_copied': 'クリップボードにコピーしました',
        'msg_support_copied': 'サポート用ログをクリップボードにコピーしました。\nサポートへの問い合わせ時に貼り付けてください。',
        
        # License
        'license_title': 'ライセンス管理',
        'license_current': '現在のライセンス',
        'license_compare': 'プラン比較:',
        'license_enter_key': 'ライセンスキー:',
        'license_activated': '{0}版がアクティベートされました',
        'license_deactivated': 'ライセンスを解除しました',
        'license_invalid': '無効なライセンスキーです',
        'license_enter_prompt': 'ライセンスキーを入力してください',
        'license_confirm_deactivate': 'ライセンスを解除しますか？',
        'license_free_features': ['抽出: 単体ファイルのみ', '更新: 3スライドまで'],
        'license_standard_features': ['抽出: 無制限+バッチ', '更新: 無制限'],
        'license_pro_features': ['Standard全機能', 'バッチ更新', '差分プレビュー', '自動バックアップ'],
        'license_key_label': 'ライセンスキー',
        'license_expires_label': '有効期限',
        'license_expired': '期限切れ',
        'license_days_left': '残り{0}日',
        'license_permanent': '永続（期限なし）',
        
        # Upgrade
        'upgrade_title': 'アップグレード',
        'upgrade_to_standard': 'この機能はStandard版以上で利用できます。\n\n📘 Standard (¥2,980/年)\n  - 無制限更新\n  - フォルダ一括抽出\n\n⭐ Pro (¥5,980/年)\n  - 上記全て\n  - バッチ更新\n  - 差分プレビュー\n  - 自動バックアップ\n\nライセンスを購入しますか？',
        'upgrade_to_pro': 'この機能はPro版で利用できます。\n\n⭐ Pro (¥5,980/年)\n  - バッチ更新\n  - 差分プレビュー\n  - 自動バックアップ\n\nアップグレードしますか？',
        
        # Welcome dialog
        'welcome_step1_title': '1. 抽出',
        'welcome_step1_desc': 'PowerPointからテキストをExcel/TXTに抽出',
        'welcome_step2_title': '2. 編集',
        'welcome_step2_desc': 'Excel/TXTで翻訳・修正作業',
        'welcome_step3_title': '3. 更新',
        'welcome_step3_desc': '編集内容をPowerPointに反映（書式維持）',
        'welcome_tip_title': '💡 ヒント',
        'welcome_tip1': '• メニュー「ヘルプ」から使い方ガイドやFAQを参照できます',
        'welcome_tip2': '• 問題発生時は「サポート用ログをコピー」をご利用ください',
        
        # Language
        'lang_select_title': '言語選択',
        'lang_restart_msg': '言語を変更しました。\n完全に反映するにはアプリを再起動してください。',
        'lang_changed': '言語を変更しました。',
        'lang_menu': '言語 / Language',
        'font_size_menu': '文字サイズ',
        'font_size_small': '小',
        'font_size_medium': '中',
        'font_size_large': '大',
        'font_size_changed': '文字サイズを変更しました。',
        'font_size_restart_msg': '文字サイズを変更しました。\n完全に反映するにはアプリを再起動してください。',
        
        # Menu
        'menu_help': 'ヘルプ',
        'menu_guide': '使い方ガイド',
        'menu_faq': 'よくある質問 (FAQ)',
        'menu_support_log': 'サポート用ログをコピー',
        'menu_license': 'ライセンス管理',
        'menu_purchase': 'ライセンス購入',
        'menu_about': 'バージョン情報',
        
        # About
        'about_title': 'バージョン情報',
        'about_text': '{0} v{1}\n\nライセンス: {2}\n\nPowerPointテキストの抽出・更新を\n効率化するツールです。\n\nby Harmonic Insight\n© 2025',
        'info_workflow_title': '基本の流れ',
        'info_workflow_1': '1. 抽出: PowerPoint → Excel / Tab text / JSON',
        'info_workflow_2': '2. 編集: 手動修正 or AI処理（誤字脱字修正、翻訳、表記統一 等）',
        'info_workflow_3': '3. 更新: Excel / JSON → PowerPoint に反映',
        'info_use_cases': '活用シーン: 翻訳、一括校正、AI連携、フォント統一チェック',
        'info_format_note': '書式（フォント・色・サイズ）は自動で維持されます',
        'info_caution_title': '注意事項',
        'info_caution_1': '1テキスト内に複数フォントがある場合、代表フォントで統一されます',
        'info_caution_2': '画像内の文字やグラフラベルは抽出できません',
        'info_caution_3': '編集時、スライド番号・IDは変更しないでください',
        
        # Dialogs
        'dialog_confirm': '確認',
        'dialog_error': 'エラー',
        'dialog_complete': '完了',
        'dialog_processing': '処理中',
        'dialog_select_pptx': 'PowerPointファイルを選択',
        'dialog_select_folder': 'フォルダを選択',
        'dialog_select_data': '編集済みファイルを選択',
        'dialog_save_pptx': '更新したPowerPointファイルを保存',
        
        # Headers
        'header_slide': 'スライド番号',
        'header_id': 'オブジェクトID',
        'header_type': 'タイプ',
        'header_text': 'テキスト内容',
        'header_filename': 'ファイル名',
        'header_datetime': '抽出日時',
        
        # Diff preview
        'diff_title': '差分プレビュー',
        'diff_changes': '変更箇所: {0}件',
        'diff_no_changes': '変更箇所なし',
        'diff_old': '旧',
        'diff_new': '新',
        
        # Pro features
        'pro_section_label': 'Pro版限定機能',
        'btn_font_analysis': 'フォント診断',
        'btn_conditional_extract': '絞り込み抽出',
        'chk_include_notes': 'スピーカーノート含む',
        'type_notes': 'ノート',
        'font_analysis_title': 'フォント診断レポート',
        'font_summary_title': 'フォント使用状況',
        'font_warning_mixed': '⚠️ 警告: 複数フォントが混在しています',
        'header_font': 'フォント',
        'header_size': 'サイズ',
        'header_bold': '太字',
        'conditional_title': '絞り込み抽出',
        'conditional_prompt': 'フィルター条件（キーワードまたは正規表現）:',
        'conditional_no_match': '条件に一致するテキストがありません',
    },
}

# Current language (default: English)
_current_lang = 'en'

def t(key: str, *args) -> str:
    """Translate key to current language"""
    text = LANGUAGES.get(_current_lang, LANGUAGES['en']).get(key, key)
    if args:
        return text.format(*args)
    return text

def set_language(lang: str):
    """Set current language"""
    global _current_lang
    if lang in LANGUAGES:
        _current_lang = lang

def get_language() -> str:
    """Get current language"""
    return _current_lang

# ============== ライセンス設定 ==============
# LICENSE_SECRET と LICENSE_TIERS は licensing.core からインポート済み

# ============== Error Messages (Bilingual) ==============
ERROR_MESSAGES = {
    'en': {
        "No such file or directory": "File not found. Please check the path.",
        "Permission denied": "Access denied. Please close the file and try again.",
        "Package not found": "Required library is not installed.",
        "is not a valid OPC package": "Not a valid PowerPoint file.",
        "Workbook": "Excel file format error. Please save as .xlsx.",
        "codec can't decode": "Character encoding error. Please save as UTF-8.",
        "list index out of range": "Data format error. Please check header row.",
        "invalid literal for int": "Number format error. Please check slide numbers.",
    },
    'ja': {
        "No such file or directory": "ファイルが見つかりません。パスを確認してください。",
        "Permission denied": "アクセス権限がありません。ファイルを閉じてから再試行してください。",
        "Package not found": "必要なライブラリがインストールされていません。",
        "is not a valid OPC package": "有効なPowerPointファイルではありません。",
        "Workbook": "Excelファイル形式に問題があります。xlsx形式で保存し直してください。",
        "codec can't decode": "文字コードに問題があります。UTF-8で保存し直してください。",
        "list index out of range": "データ形式が不正です。ヘッダー行を確認してください。",
        "invalid literal for int": "数値形式が不正です。スライド番号を確認してください。",
    },
}

def translate_error(error_msg: str) -> str:
    """Translate error message to current language"""
    error_str = str(error_msg)
    msgs = ERROR_MESSAGES.get(_current_lang, ERROR_MESSAGES['en'])
    for eng, translated in msgs.items():
        if eng in error_str:
            return translated
    return error_str

# ============== デザインシステム ==============
COLOR_PALETTE = {
    # Backgrounds
    "bg_primary": "#FAFBFC", "bg_secondary": "#F3F4F6", "bg_elevated": "#FFFFFF",
    "bg_input": "#FFFFFF", "bg_code": "#1E1E1E", "bg_log": "#F9FAFB",
    "bg_standard_section": "#EEF2FF",  # Light indigo for Standard section
    "bg_pro_section": "#F5F3FF",  # Light purple for Pro section
    # Text
    "text_primary": "#1A202C", "text_secondary": "#4A5568", "text_muted": "#718096",
    # Mode colors - Premium navy (Extract=Navy, Update=Teal)
    "brand_primary": "#1E40AF", "brand_primary_hover": "#1E3A8A",  # Deep navy blue
    "brand_update": "#0D9488", "brand_update_hover": "#0F766E",  # Update mode teal
    # Tier colors (Free=Gray, Standard=Indigo, Pro=Deep Navy)
    "brand_standard": "#3730A3", "brand_standard_hover": "#312E81",  # Deep Indigo
    "brand_pro": "#1E3A8A", "brand_pro_hover": "#1E3A8A",  # Deep Navy (Premium)
    # Status
    "success": "#059669", "warning": "#D97706", "error": "#DC2626", "info": "#2563EB",
    # Borders
    "border_light": "#E5E7EB", "border_medium": "#D1D5DB", "border_focus": "#1E40AF",
    "border_standard": "#C7D2FE",  # Light indigo border
    "border_pro": "#DBEAFE",  # Light blue border
    "header_border": "#E5E7EB",  # Header separator line
    # Interactive
    "hover": "#F9FAFB", "active": "#F3F4F6", "selection": "#DBEAFE",
    # Badges
    "free_badge": "#6B7280", "standard_badge": "#3730A3", "pro_badge": "#1E3A8A",
    # Section labels
    "standard_link": "#3730A3", "pro_link": "#1E3A8A",
}

# Professional font stack (Yu Gothic UI for Japanese support)
FONT_FAMILY = "Yu Gothic UI"

# フォントサイズ設定（小/中/大）
FONT_SIZE_PRESETS = {
    'small': {'base': 10, 'label_key': 'font_size_small'},
    'medium': {'base': 11, 'label_key': 'font_size_medium'},
    'large': {'base': 13, 'label_key': 'font_size_large'},
}

def get_fonts(size_preset: str = 'medium') -> dict:
    """フォントサイズプリセットに基づいてフォント辞書を生成"""
    base = FONT_SIZE_PRESETS.get(size_preset, FONT_SIZE_PRESETS['medium'])['base']
    return {
        "display": (FONT_FAMILY, base + 9, "bold"), 
        "heading": (FONT_FAMILY, base + 3, "bold"),
        "subheading": (FONT_FAMILY, base + 1, "bold"), 
        "body": (FONT_FAMILY, base, "normal"),
        "body_medium": (FONT_FAMILY, base, "bold"), 
        "caption": (FONT_FAMILY, base, "normal"),
        "code": (FONT_FAMILY, base, "normal"),  # 処理結果エリア用（Yu Gothic UIで統一）
        "button": (FONT_FAMILY, base, "bold"),  # ボタンもbase統一
        "badge": (FONT_FAMILY, base - 1, "bold"),
    }

# デフォルトフォント（後で設定から上書き）
FONTS = get_fonts('medium')

SPACING = {"xs": 2, "sm": 4, "md": 8, "lg": 12, "xl": 16, "2xl": 20, "3xl": 24, "4xl": 32}

# Professional Unicode symbols (no emoji - cleaner look)
ICONS = {
    "extract": "↑", "update": "↓", "folder": "▤", "file": "▢", "excel": "▦",
    "settings": "⚙", "help": "?", "save": "↓", "clear": "×", "copy": "⎘",
    "powerpoint": "◈", "success": "✓", "processing": "◌", "batch": "▣",
    "switch": "⇄", "info": "i", "pro": "★", "standard": "◆", "lock": "◉",
    "unlock": "○", "key": "⚿", "preview": "◎", "backup": "▤", "cancel": "×",
    "warning": "△", "support": "?", "link": "↗",
}

PLACEHOLDER_TYPES = {
    'en': {1: "Title", 2: "Body", 3: "Chart", 4: "Date", 5: "Slide Number", 6: "Header", 7: "Footer"},
    'ja': {1: "タイトル", 2: "本文", 3: "図表", 4: "日付", 5: "スライド番号", 6: "ヘッダー", 7: "フッター"},
}

def get_placeholder_type(type_id: int) -> str:
    types = PLACEHOLDER_TYPES.get(_current_lang, PLACEHOLDER_TYPES['en'])
    return types.get(type_id, f"Placeholder({type_id})")


def save_error_log(error: Exception, context: str = ""):
    try:
        CONFIG_DIR.mkdir(parents=True, exist_ok=True)
        timestamp = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
        with open(ERROR_LOG_FILE, 'a', encoding='utf-8') as f:
            f.write(f"\n{'='*60}\nTime: {timestamp}\nContext: {context}\nError: {error}\nDetails:\n{traceback.format_exc()}\n")
    except:
        pass


class LicenseManager:
    """
    ライセンス管理クラス
    検証ロジックは licensing.core モジュールを使用
    """
    def __init__(self):
        self._ensure_config_dir()
        self.license_info = self._load_license()

    def _ensure_config_dir(self):
        CONFIG_DIR.mkdir(parents=True, exist_ok=True)

    def _load_license(self) -> Dict:
        if LICENSE_FILE.exists():
            try:
                with open(LICENSE_FILE, 'r', encoding='utf-8') as f:
                    data = json.load(f)
                    if data.get('key'):
                        # Check if license is still valid (共通モジュールを使用)
                        is_valid, _ = validate_key(data['key'])
                        if is_valid:
                            return data
                        else:
                            # License expired, reset to free
                            return {'type': 'free', 'key': '', 'expires': None}
            except:
                pass
        return {'type': 'free', 'key': '', 'expires': None}

    def _save_license(self, data: Dict):
        with open(LICENSE_FILE, 'w', encoding='utf-8') as f:
            json.dump(data, f, ensure_ascii=False, indent=2)

    def _validate_key(self, key: str) -> bool:
        """キーの検証（共通モジュールを使用）"""
        is_valid, _ = validate_key(key)
        return is_valid

    def _get_license_type_from_key(self, key: str) -> str:
        """キーからライセンスタイプを取得（共通モジュールを使用）"""
        return get_tier_from_key(key)

    def generate_license_key(self, user_id: str, license_type: str = 'standard', days_valid: int = 365) -> str:
        """レガシーフォーマットのキー生成（チェックサム付き）"""
        type_code = 'PRO' if license_type == 'professional' else 'STD'
        user_hash = hashlib.md5(user_id.encode()).hexdigest()[:4].upper()
        expire_code = base64.b64encode((datetime.now() + timedelta(days=days_valid)).strftime('%Y%m%d').encode()).decode()[:4].upper()
        import random
        random_part = ''.join(random.choices('ABCDEFGHIJKLMNOPQRSTUVWXYZ0123456789', k=4))
        key_body = f"{type_code}-{user_hash}-{expire_code}-{random_part}"
        return f"{key_body}-{generate_checksum(key_body)}"

    def activate_license(self, key: str) -> Tuple[bool, str]:
        if not key:
            return False, t('license_enter_prompt')

        # 共通モジュールを使用してキーを検証
        is_valid, license_type = validate_key(key)
        if not is_valid:
            return False, t('license_invalid')

        # 共通モジュールを使用して有効期限を取得
        expires = get_expiration_from_key(key)

        self.license_info = {
            'type': license_type,
            'key': key.strip().upper(),
            'activated': datetime.now().isoformat(),
            'expires': expires
        }
        self._save_license(self.license_info)
        return True, t('license_activated', LICENSE_TIERS[license_type]['name'])
    
    def deactivate_license(self):
        self.license_info = {'type': 'free', 'key': '', 'expires': None}
        if LICENSE_FILE.exists():
            LICENSE_FILE.unlink()
    
    def get_license_type(self) -> str:
        return self.license_info.get('type', 'free')
    
    def get_tier_info(self) -> Dict:
        return LICENSE_TIERS.get(self.get_license_type(), LICENSE_TIERS['free'])
    
    def can_use_feature(self, feature: str) -> bool:
        return self.get_tier_info().get(feature, False)
    
    def get_update_limit(self) -> Optional[int]:
        return self.get_tier_info().get('update_slide_limit')
    
    def get_format_quality(self) -> str:
        return self.get_tier_info().get('format_quality', 'basic')


class ConfigManager:
    DEFAULT_CONFIG = {
        'language': 'en',  # Default: English
        'output_format': 'tab',  # 'tab', 'csv', 'excel'
        'include_metadata': True,
        'auto_backup': True,
        'last_directory': '',
        'window_geometry': '1200x900',
        'first_launch': True,
        'font_size': 'medium',
    }
    
    def __init__(self):
        global FONTS
        CONFIG_DIR.mkdir(parents=True, exist_ok=True)
        self.config = self._load_config()
        # Apply saved language
        set_language(self.config.get('language', 'en'))
        # Apply saved font size
        FONTS = get_fonts(self.config.get('font_size', 'medium'))
    
    def _load_config(self) -> Dict:
        if CONFIG_FILE.exists():
            try:
                with open(CONFIG_FILE, 'r', encoding='utf-8') as f:
                    return {**self.DEFAULT_CONFIG, **json.load(f)}
            except:
                pass
        return self.DEFAULT_CONFIG.copy()
    
    def save(self):
        with open(CONFIG_FILE, 'w', encoding='utf-8') as f:
            json.dump(self.config, f, ensure_ascii=False, indent=2)
    
    def get(self, key: str, default=None):
        return self.config.get(key, default)
    
    def set(self, key: str, value):
        self.config[key] = value
        self.save()


class ModernButton(ttk.Button):
    def __init__(self, parent, text="", command=None, style_type="primary", icon=None, **kwargs):
        display_text = f"{icon} {text}" if icon else text
        super().__init__(parent, text=display_text, command=command, **kwargs)
        self.configure(style=f"{style_type.title()}.Modern.TButton")


class ModernCard(ttk.Frame):
    def __init__(self, parent, title=None, badge=None, **kwargs):
        super().__init__(parent, **kwargs)
        self.configure(style='Card.TFrame')
        self.grid_columnconfigure(0, weight=1)
        if title:
            header = ttk.Frame(self, style='CardHeader.TFrame')
            header.grid(row=0, column=0, sticky='ew', padx=SPACING["lg"], pady=(SPACING["lg"], SPACING["sm"]))
            header.grid_columnconfigure(0, weight=1)
            ttk.Label(header, text=title, style='CardTitle.TLabel').grid(row=0, column=0, sticky='w')
            if badge:
                ttk.Label(header, text=badge, style='Badge.TLabel').grid(row=0, column=1, sticky='e', padx=(SPACING["sm"], 0))
            self.content_row = 1
        else:
            self.content_row = 0
    
    def add_content(self, widget):
        widget.grid(row=self.content_row, column=0, sticky='nsew', padx=SPACING["lg"], pady=(0, SPACING["lg"]))


# ============== メインアプリケーション ==============
class InsightSlidesApp:
    def __init__(self, root):
        self.root = root
        self.license_manager = LicenseManager()
        self.config_manager = ConfigManager()
        self.current_mode = "extract"
        self.processing = False
        self.cancel_requested = False
        self.presentation = None
        self.log_buffer = []
        
        self._setup_window()
        self._apply_styles()
        self._create_menu()
        self._create_layout()
        self.root.protocol("WM_DELETE_WINDOW", self._on_closing)
        
        if self.config_manager.get('first_launch', True):
            self.root.after(500, self._show_language_dialog)
    
    def _setup_window(self):
        tier = self.license_manager.get_tier_info()
        self.root.title(f"{APP_NAME} {tier['name']} - {t('app_subtitle')} v{APP_VERSION}")
        self.root.geometry(self.config_manager.get('window_geometry', '1200x900'))
        self.root.minsize(1000, 700)
        self.root.configure(bg=COLOR_PALETTE["bg_primary"])
    
    def _apply_styles(self):
        self.style = ttk.Style()
        self.style.theme_use('clam')
        self.style.configure('Main.TFrame', background=COLOR_PALETTE["bg_primary"])
        self.style.configure('Content.TFrame', background=COLOR_PALETTE["bg_primary"])
        self.style.configure('Card.TFrame', background=COLOR_PALETTE["bg_elevated"], relief='flat', borderwidth=1)
        self.style.configure('CardHeader.TFrame', background=COLOR_PALETTE["bg_elevated"])
        self.style.configure('CardTitle.TLabel', background=COLOR_PALETTE["bg_elevated"], foreground=COLOR_PALETTE["text_primary"], font=FONTS["subheading"])
        self.style.configure('Header.TFrame', background=COLOR_PALETTE["bg_elevated"])
        self.style.configure('AppTitle.TLabel', background=COLOR_PALETTE["bg_elevated"], foreground=COLOR_PALETTE["brand_primary"], font=FONTS["display"])
        
        for name, bg in [('ProBadge', 'pro_badge'), ('StandardBadge', 'standard_badge'), ('FreeBadge', 'free_badge')]:
            self.style.configure(f'{name}.TLabel', background=COLOR_PALETTE[bg], foreground='white', font=FONTS["badge"], padding=(6, 2))
        
        btn_styles = {
            'Primary': (COLOR_PALETTE["brand_primary"], COLOR_PALETTE["brand_primary_hover"]),  # Blue (Extract)
            'Update': (COLOR_PALETTE["brand_update"], COLOR_PALETTE["brand_update_hover"]),  # Teal (Update mode)
            'Standard': (COLOR_PALETTE["brand_standard"], COLOR_PALETTE["brand_standard_hover"]),  # Indigo
            'Pro': (COLOR_PALETTE["brand_pro"], COLOR_PALETTE["brand_pro_hover"]),  # Purple
            'Secondary': (COLOR_PALETTE["brand_standard"], COLOR_PALETTE["brand_standard_hover"]),  # Indigo (same as Standard)
            'Outline': (COLOR_PALETTE["bg_elevated"], COLOR_PALETTE["hover"]),
            'Disabled': (COLOR_PALETTE["border_light"], COLOR_PALETTE["border_light"]),
        }
        for name, (bg, hover) in btn_styles.items():
            fg = 'white' if name not in ['Outline', 'Disabled'] else COLOR_PALETTE["text_primary"]
            self.style.configure(f'{name}.Modern.TButton', background=bg, foreground=fg, font=FONTS["button"], borderwidth=1, relief='flat', padding=(SPACING["md"], SPACING["sm"]))
            self.style.map(f'{name}.Modern.TButton', background=[('active', hover), ('disabled', COLOR_PALETTE["border_light"])])
    
    def _create_menu(self):
        menubar = tk.Menu(self.root)
        self.root.config(menu=menubar)
        
        help_menu = tk.Menu(menubar, tearoff=0)
        menubar.add_cascade(label=t('menu_help'), menu=help_menu)
        help_menu.add_command(label=f"{ICONS['help']} {t('menu_guide')}", command=lambda: webbrowser.open(SUPPORT_LINKS["tutorial"]))
        help_menu.add_command(label=f"{ICONS['info']} {t('menu_faq')}", command=lambda: webbrowser.open(SUPPORT_LINKS["faq"]))
        help_menu.add_separator()
        help_menu.add_command(label=f"{ICONS['support']} {t('menu_support_log')}", command=self._copy_support_log)
        help_menu.add_separator()
        help_menu.add_command(label=f"{ICONS['key']} {t('menu_license')}", command=self._show_license_dialog)
        help_menu.add_command(label=f"{ICONS['link']} {t('menu_purchase')}", command=lambda: webbrowser.open(SUPPORT_LINKS["purchase"]))
        help_menu.add_separator()
        
        # Language submenu
        lang_menu = tk.Menu(help_menu, tearoff=0)
        help_menu.add_cascade(label=f"🌐 {t('lang_menu')}", menu=lang_menu)
        lang_menu.add_command(label="English", command=lambda: self._change_language('en'))
        lang_menu.add_command(label="日本語", command=lambda: self._change_language('ja'))
        
        # Font size submenu
        font_menu = tk.Menu(help_menu, tearoff=0)
        help_menu.add_cascade(label=f"🔤 {t('font_size_menu')}", menu=font_menu)
        current_size = self.config_manager.get('font_size', 'medium')
        for size_key, size_info in FONT_SIZE_PRESETS.items():
            check = "● " if size_key == current_size else "○ "
            font_menu.add_command(label=f"{check}{t(size_info['label_key'])}", command=lambda s=size_key: self._change_font_size(s))
        
        help_menu.add_separator()
        help_menu.add_command(label=t('menu_about'), command=self._show_about)
    
    def _create_layout(self):
        # 既存のメインコンテナがあれば削除
        if hasattr(self, 'main_container') and self.main_container:
            self.main_container.destroy()
        
        self.main_container = ttk.Frame(self.root, style='Main.TFrame')
        self.main_container.pack(fill='both', expand=True, padx=SPACING["xl"], pady=SPACING["xl"])
        self.main_container.grid_columnconfigure(0, weight=1)
        self.main_container.grid_rowconfigure(1, weight=1)
        
        self._create_header(self.main_container)
        
        content = ttk.Frame(self.main_container, style='Content.TFrame')
        content.grid(row=1, column=0, sticky='nsew')
        content.grid_columnconfigure(1, weight=1)
        content.grid_rowconfigure(0, weight=1)
        
        self._create_controls(content)
        self._create_output(content)
    
    def _create_header(self, parent):
        # ヘッダーコンテナ（下線付き）
        header_container = tk.Frame(parent, bg=COLOR_PALETTE["bg_elevated"])
        header_container.grid(row=0, column=0, sticky='ew', pady=(0, SPACING["2xl"]))
        
        header = tk.Frame(header_container, bg=COLOR_PALETTE["bg_elevated"], padx=SPACING["xl"], pady=SPACING["lg"])
        header.pack(fill='x')
        
        # 薄いグレーの下線（高級感）
        separator = tk.Frame(header_container, height=1, bg=COLOR_PALETTE["header_border"])
        separator.pack(fill='x')
        
        # 左側：タイトル + サブタイトル
        left_frame = tk.Frame(header, bg=COLOR_PALETTE["bg_elevated"])
        left_frame.pack(side='left', anchor='w')
        
        # タイトル行
        title_frame = tk.Frame(left_frame, bg=COLOR_PALETTE["bg_elevated"])
        title_frame.pack(anchor='w')
        
        # ロゴアイコン（プロフェッショナルなシンボル）
        tk.Label(title_frame, text="◈", font=FONTS["display"], fg=COLOR_PALETTE["brand_primary"], bg=COLOR_PALETTE["bg_elevated"]).pack(side='left', padx=(0, 8))
        
        # タイトル（Insight Slides）
        tk.Label(title_frame, text="Insight Slides", font=FONTS["display"], fg=COLOR_PALETTE["text_primary"], bg=COLOR_PALETTE["bg_elevated"]).pack(side='left')
        
        # ライセンスラベル（動的に更新可能）- タイトルより少し控えめ
        self.license_title_label = tk.Label(title_frame, text="", font=FONTS["heading"], bg=COLOR_PALETTE["bg_elevated"])
        self.license_title_label.pack(side='left', padx=(4, 0))
        self._update_license_title_label()
        
        # ライセンスバッジ参照保持（互換性のため）
        self.license_badge = self.license_title_label
        
        # サブタイトル
        tk.Label(left_frame, text=t('app_subtitle'), font=FONTS["caption"], fg=COLOR_PALETTE["text_muted"], bg=COLOR_PALETTE["bg_elevated"]).pack(anchor='w', pady=(2, 0))
        
        # 右側：モード + ライセンス + バージョン
        right_frame = tk.Frame(header, bg=COLOR_PALETTE["bg_elevated"])
        right_frame.pack(side='right', anchor='ne')
        
        # モードインジケーター
        self.mode_frame = tk.Frame(right_frame, bg=COLOR_PALETTE["bg_elevated"])
        self.mode_frame.pack(side='left', padx=(0, SPACING["lg"]))
        
        self.mode_dot = tk.Label(self.mode_frame, text="●", font=FONTS["caption"], fg=COLOR_PALETTE["brand_primary"], bg=COLOR_PALETTE["bg_elevated"])
        self.mode_dot.pack(side='left', padx=(0, 4))
        self.mode_label = tk.Label(self.mode_frame, text=t('mode_extract'), font=FONTS["caption"], fg=COLOR_PALETTE["text_muted"], bg=COLOR_PALETTE["bg_elevated"])
        self.mode_label.pack(side='left')
        
        # 区切り
        tk.Label(right_frame, text="·", font=FONTS["caption"], fg=COLOR_PALETTE["border_medium"], bg=COLOR_PALETTE["bg_elevated"]).pack(side='left', padx=(0, SPACING["lg"]))
        
        # ライセンスリンク
        license_btn = tk.Label(right_frame, text=t('btn_license'), font=FONTS["caption"], fg=COLOR_PALETTE["text_muted"], bg=COLOR_PALETTE["bg_elevated"], cursor="hand2")
        license_btn.pack(side='left', padx=(0, SPACING["lg"]))
        license_btn.bind("<Button-1>", lambda e: self._show_license_dialog())
        license_btn.bind("<Enter>", lambda e: license_btn.config(fg=COLOR_PALETTE["brand_primary"]))
        license_btn.bind("<Leave>", lambda e: license_btn.config(fg=COLOR_PALETTE["text_muted"]))
        
        # 区切り
        tk.Label(right_frame, text="·", font=FONTS["caption"], fg=COLOR_PALETTE["border_medium"], bg=COLOR_PALETTE["bg_elevated"]).pack(side='left', padx=(0, SPACING["lg"]))
        
        # バージョン
        tk.Label(right_frame, text=f"v{APP_VERSION}", font=FONTS["caption"], fg=COLOR_PALETTE["text_muted"], bg=COLOR_PALETTE["bg_elevated"]).pack(side='left')
    
    def _update_license_title_label(self):
        """ヘッダーのライセンス表示を更新"""
        tier = self.license_manager.get_tier_info()
        if tier['name'] == 'Pro':
            self.license_title_label.configure(text=" Professional", fg=COLOR_PALETTE["brand_primary"])
        elif tier['name'] == 'Standard':
            self.license_title_label.configure(text=" Standard", fg=COLOR_PALETTE["brand_standard"])
        else:
            self.license_title_label.configure(text="", fg=COLOR_PALETTE["text_muted"])
    
    def _create_controls(self, parent):
        frame = ttk.Frame(parent, style='Content.TFrame')
        frame.grid(row=0, column=0, sticky='nsew', padx=(0, SPACING["lg"]))
        frame.grid_rowconfigure(3, weight=1)
        
        # Mode card
        mode_card = ModernCard(frame, title=t('panel_mode'))
        mode_card.grid(row=0, column=0, sticky='ew', pady=(0, SPACING["lg"]))
        mode_content = ttk.Frame(mode_card, style='CardHeader.TFrame')
        mode_card.add_content(mode_content)
        mode_content.grid_columnconfigure(0, weight=1)
        mode_content.grid_columnconfigure(1, weight=1)
        self.extract_btn = ModernButton(mode_content, t('mode_extract_short'), self._switch_extract, style_type="primary", icon=ICONS["extract"])
        self.extract_btn.grid(row=0, column=0, sticky='ew', padx=(0, SPACING["sm"]))
        self.update_btn = ModernButton(mode_content, t('mode_update_short'), self._switch_update, style_type="outline", icon=ICONS["update"])
        self.update_btn.grid(row=0, column=1, sticky='ew')
        
        # File card (contains mode-specific panels)
        self.file_card = ModernCard(frame, title=t('panel_file'))
        self.file_card.grid(row=1, column=0, sticky='ew', pady=(0, SPACING["lg"]))
        self.file_content = ttk.Frame(self.file_card, style='CardHeader.TFrame')
        self.file_card.add_content(self.file_content)
        self.file_content.grid_columnconfigure(0, weight=1)
        self._create_extract_panel()
        self._create_update_panel()
        
        # Status bar (compact, at bottom of left panel)
        status_frame = ttk.Frame(frame, style='Content.TFrame')
        status_frame.grid(row=2, column=0, sticky='ew')
        status_frame.grid_columnconfigure(0, weight=1)
        
        self.status_label = ttk.Label(status_frame, text=t('status_waiting'), font=FONTS["caption"], foreground=COLOR_PALETTE["text_muted"], background=COLOR_PALETTE["bg_primary"])
        self.status_label.grid(row=0, column=0, sticky='w', pady=(SPACING["sm"], SPACING["xs"]))
        
        self.progress = ttk.Progressbar(status_frame, mode='indeterminate')
        self.progress.grid(row=1, column=0, sticky='ew', pady=(0, SPACING["sm"]))
        
        # Action buttons row
        btns = ttk.Frame(status_frame, style='Content.TFrame')
        btns.grid(row=2, column=0, sticky='ew')
        self.cancel_btn = ModernButton(btns, t('btn_cancel'), self._cancel, style_type="outline", icon=ICONS["cancel"])
        self.cancel_btn.pack(side='left', padx=(0, SPACING["sm"]))
        self.cancel_btn.configure(state='disabled')
        ModernButton(btns, t('btn_clear'), self._clear_output, style_type="outline", icon=ICONS["clear"]).pack(side='left', padx=(0, SPACING["sm"]))
        ModernButton(btns, t('btn_copy'), self._copy_output, style_type="outline", icon=ICONS["copy"]).pack(side='left')
        
        self._switch_extract()
    
    def _create_extract_panel(self):
        self.extract_frame = ttk.Frame(self.file_content, style='CardHeader.TFrame')
        self.extract_frame.grid_columnconfigure(0, weight=1)
        
        row_idx = 0
        
        # === 抽出オプション セクション ===
        options_section = tk.Frame(self.extract_frame, bg=COLOR_PALETTE["bg_secondary"], padx=SPACING["md"], pady=SPACING["md"])
        options_section.grid(row=row_idx, column=0, sticky='ew', pady=(0, SPACING["md"]))
        options_section.grid_columnconfigure(1, weight=1)
        row_idx += 1
        
        # セクションヘッダー
        tk.Label(options_section, text=f"⚙ {t('panel_extract_options')}", font=FONTS["caption"], 
                 fg=COLOR_PALETTE["text_secondary"], bg=COLOR_PALETTE["bg_secondary"]).grid(row=0, column=0, columnspan=2, sticky='w', pady=(0, SPACING["sm"]))
        
        # 出力形式
        tk.Label(options_section, text=t('setting_output_format'), font=FONTS["body"], 
                 bg=COLOR_PALETTE["bg_secondary"]).grid(row=1, column=0, sticky='w', pady=SPACING["xs"])
        self.output_format_var = tk.StringVar(value=self.config_manager.get('output_format', 'excel'))
        self.format_combo = ttk.Combobox(options_section, textvariable=self.output_format_var, 
                                         values=['excel', 'tab', 'json'], state="readonly", width=15)
        self.format_combo.grid(row=1, column=1, sticky='w', padx=(SPACING["sm"], 0), pady=SPACING["xs"])
        self.format_combo.bind('<<ComboboxSelected>>', lambda e: self.config_manager.set('output_format', self.output_format_var.get()))
        
        # ファイル名・日時を含める
        self.include_metadata_var = tk.BooleanVar(value=self.config_manager.get('include_metadata'))
        meta_check = tk.Checkbutton(options_section, text=t('setting_include_meta'), variable=self.include_metadata_var,
                                    font=FONTS["body"], bg=COLOR_PALETTE["bg_secondary"],
                                    activebackground=COLOR_PALETTE["bg_secondary"],
                                    command=lambda: self.config_manager.set('include_metadata', self.include_metadata_var.get()))
        meta_check.grid(row=2, column=0, columnspan=2, sticky='w', pady=SPACING["xs"])
        
        # スピーカーノート含む (Pro)
        can_notes = self.license_manager.can_use_feature('speaker_notes')
        self.include_notes_var = tk.BooleanVar(value=False)
        notes_frame = tk.Frame(options_section, bg=COLOR_PALETTE["bg_secondary"])
        notes_frame.grid(row=3, column=0, columnspan=2, sticky='w', pady=SPACING["xs"])
        notes_color = COLOR_PALETTE["text_primary"] if can_notes else COLOR_PALETTE["text_muted"]
        self.notes_check = tk.Checkbutton(notes_frame, text=t('chk_include_notes'), variable=self.include_notes_var,
                                          font=FONTS["body"], bg=COLOR_PALETTE["bg_secondary"], fg=notes_color,
                                          activebackground=COLOR_PALETTE["bg_secondary"],
                                          state='normal' if can_notes else 'disabled')
        self.notes_check.pack(side='left')
        if not can_notes:
            tk.Label(notes_frame, text=f" {ICONS['pro']}Pro", font=FONTS["caption"], 
                     fg=COLOR_PALETTE["pro_link"], bg=COLOR_PALETTE["bg_secondary"]).pack(side='left')
        
        # === 抽出実行 セクション ===
        run_section = tk.Frame(self.extract_frame, bg=COLOR_PALETTE["bg_elevated"])
        run_section.grid(row=row_idx, column=0, sticky='ew', pady=(0, SPACING["md"]))
        run_section.grid_columnconfigure(0, weight=1)
        row_idx += 1
        
        # セクションヘッダー
        tk.Label(run_section, text=f"📤 {t('panel_extract_run')}", font=FONTS["caption"], 
                 fg=COLOR_PALETTE["text_secondary"], bg=COLOR_PALETTE["bg_elevated"]).grid(row=0, column=0, sticky='w', pady=(0, SPACING["sm"]))
        
        # ファイル選択ボタン
        ModernButton(run_section, t('btn_single_file'), self._extract_single, 
                     style_type="primary", icon=ICONS["file"]).grid(row=1, column=0, sticky='ew', pady=(0, SPACING["sm"]))
        
        # フォルダー一括 (Standard+)
        can_batch = self.license_manager.can_use_feature('batch_extract')
        batch_frame = tk.Frame(run_section, bg=COLOR_PALETTE["bg_elevated"])
        batch_frame.grid(row=2, column=0, sticky='ew')
        batch_frame.grid_columnconfigure(0, weight=1)
        
        batch_style = "standard" if can_batch else "disabled"
        batch_icon = ICONS["batch"] if can_batch else ICONS["lock"]
        self.batch_extract_btn = ModernButton(batch_frame, t('btn_batch_folder'), 
                                              self._extract_batch if can_batch else self._show_upgrade_prompt, 
                                              style_type=batch_style, icon=batch_icon)
        self.batch_extract_btn.grid(row=0, column=0, sticky='ew')
        if not can_batch:
            tk.Label(batch_frame, text=f" {ICONS['standard']}Standard", font=FONTS["caption"], 
                     fg=COLOR_PALETTE["standard_link"], bg=COLOR_PALETTE["bg_elevated"]).grid(row=0, column=1, padx=(SPACING["sm"], 0))
        
        # === 分析ツール セクション (Pro) ===
        can_font = self.license_manager.can_use_feature('font_analysis')
        can_conditional = self.license_manager.can_use_feature('conditional_extract')
        
        pro_section = tk.Frame(self.extract_frame, bg=COLOR_PALETTE["bg_pro_section"], padx=SPACING["md"], pady=SPACING["md"])
        pro_section.grid(row=row_idx, column=0, sticky='ew', pady=(SPACING["sm"], 0))
        pro_section.grid_columnconfigure(0, weight=1)
        
        # Pro section header
        tk.Label(pro_section, text=f"{ICONS['pro']} {t('panel_analysis_tools')} (Pro)", font=FONTS["caption"], 
                 fg=COLOR_PALETTE["pro_link"], bg=COLOR_PALETTE["bg_pro_section"]).grid(row=0, column=0, sticky='w', pady=(0, SPACING["sm"]))
        
        # Conditional extract button (Pro)
        cond_style = "pro" if can_conditional else "disabled"
        cond_icon = "🔍" if can_conditional else ICONS["lock"]
        self.conditional_btn = ModernButton(pro_section, t('btn_conditional_extract'), 
                                            self._conditional_extract if can_conditional else self._show_upgrade_prompt,
                                            style_type=cond_style, icon=cond_icon)
        self.conditional_btn.grid(row=1, column=0, sticky='ew', pady=(0, SPACING["xs"]))
        
        # Font analysis button (Pro)
        font_style = "pro" if can_font else "disabled"
        font_icon = "📊" if can_font else ICONS["lock"]
        self.font_btn = ModernButton(pro_section, t('btn_font_analysis'),
                                     self._font_analysis if can_font else self._show_upgrade_prompt,
                                     style_type=font_style, icon=font_icon)
        self.font_btn.grid(row=2, column=0, sticky='ew')
    
    def _create_update_panel(self):
        self.update_frame = ttk.Frame(self.file_content, style='CardHeader.TFrame')
        self.update_frame.grid_columnconfigure(0, weight=1)
        
        row_idx = 0
        
        # === 更新実行 セクション ===
        tier = self.license_manager.get_tier_info()
        limit = tier['update_slide_limit']
        can_backup = self.license_manager.can_use_feature('auto_backup')
        
        run_section = tk.Frame(self.update_frame, bg=COLOR_PALETTE["bg_elevated"])
        run_section.grid(row=row_idx, column=0, sticky='ew', pady=(0, SPACING["md"]))
        run_section.grid_columnconfigure(0, weight=1)
        row_idx += 1
        
        # セクションヘッダー
        tk.Label(run_section, text=f"📥 {t('panel_update_run')}", font=FONTS["caption"], 
                 fg=COLOR_PALETTE["text_secondary"], bg=COLOR_PALETTE["bg_elevated"]).grid(row=0, column=0, sticky='w', pady=(0, SPACING["sm"]))
        
        # 制限警告（Free版のみ）
        if limit:
            tk.Label(run_section, text=t('msg_update_limit', limit), font=FONTS["caption"], 
                     fg=COLOR_PALETTE["warning"], bg=COLOR_PALETTE["bg_elevated"], wraplength=280).grid(row=1, column=0, sticky='w', pady=(0, SPACING["sm"]))
        
        # Excelから更新
        ModernButton(run_section, t('btn_from_excel'), self._update_excel, 
                     style_type="update", icon=ICONS["excel"]).grid(row=2, column=0, sticky='ew', pady=(0, SPACING["sm"]))
        
        # JSONから更新
        ModernButton(run_section, t('btn_from_json'), self._update_json, 
                     style_type="update", icon=ICONS["file"]).grid(row=3, column=0, sticky='ew', pady=(0, SPACING["sm"]))
        
        # 自動バックアップ (Pro) - 更新実行内に配置
        self.auto_backup_var = tk.BooleanVar(value=self.config_manager.get('auto_backup') if can_backup else False)
        backup_frame = tk.Frame(run_section, bg=COLOR_PALETTE["bg_elevated"])
        backup_frame.grid(row=4, column=0, sticky='w', pady=(SPACING["sm"], 0))
        backup_color = COLOR_PALETTE["text_primary"] if can_backup else COLOR_PALETTE["text_muted"]
        self.backup_check = tk.Checkbutton(backup_frame, text=t('setting_auto_backup'), variable=self.auto_backup_var,
                                           font=FONTS["body"], bg=COLOR_PALETTE["bg_elevated"], fg=backup_color,
                                           activebackground=COLOR_PALETTE["bg_elevated"],
                                           state='normal' if can_backup else 'disabled',
                                           command=lambda: self.config_manager.set('auto_backup', self.auto_backup_var.get()))
        self.backup_check.pack(side='left')
        if not can_backup:
            tk.Label(backup_frame, text=f" {ICONS['pro']}Pro", font=FONTS["caption"], 
                     fg=COLOR_PALETTE["pro_link"], bg=COLOR_PALETTE["bg_elevated"]).pack(side='left')
        
        # === 拡張機能 セクション (Pro) ===
        can_preview = self.license_manager.can_use_feature('diff_preview')
        can_batch = self.license_manager.can_use_feature('batch_update')
        
        pro_section = tk.Frame(self.update_frame, bg=COLOR_PALETTE["bg_pro_section"], padx=SPACING["md"], pady=SPACING["md"])
        pro_section.grid(row=row_idx, column=0, sticky='ew', pady=(SPACING["sm"], 0))
        pro_section.grid_columnconfigure(0, weight=1)
        
        # Pro section header
        tk.Label(pro_section, text=f"{ICONS['pro']} {t('panel_pro_features')} (Pro)", font=FONTS["caption"], 
                 fg=COLOR_PALETTE["pro_link"], bg=COLOR_PALETTE["bg_pro_section"]).grid(row=0, column=0, sticky='w', pady=(0, SPACING["sm"]))
        
        # Diff preview (Pro) - 説明付き
        preview_style = "pro" if can_preview else "disabled"
        preview_icon = ICONS["preview"] if can_preview else ICONS["lock"]
        self.preview_btn = ModernButton(pro_section, t('btn_diff_preview'), 
                                        self._run_preview if can_preview else self._show_upgrade_prompt, 
                                        style_type=preview_style, icon=preview_icon)
        self.preview_btn.grid(row=1, column=0, sticky='ew', pady=(0, SPACING["xs"]))
        
        # 差分プレビュー説明
        tk.Label(pro_section, text=t('diff_preview_desc'), font=FONTS["caption"], 
                 fg=COLOR_PALETTE["text_muted"], bg=COLOR_PALETTE["bg_pro_section"]).grid(row=2, column=0, sticky='w', pady=(0, SPACING["sm"]))
        
        # Batch update (Pro)
        batch_style = "pro" if can_batch else "disabled"
        batch_icon = ICONS["batch"] if can_batch else ICONS["lock"]
        self.batch_update_btn = ModernButton(pro_section, t('btn_batch_update'), 
                                             self._update_batch if can_batch else self._show_upgrade_prompt, 
                                             style_type=batch_style, icon=batch_icon)
        self.batch_update_btn.grid(row=3, column=0, sticky='ew')
    
    def _create_output(self, parent):
        card = ModernCard(parent, title=t('panel_output'))
        card.grid(row=0, column=1, sticky='nsew')
        content = ttk.Frame(card, style='CardHeader.TFrame')
        card.add_content(content)
        content.grid_columnconfigure(0, weight=1)
        content.grid_rowconfigure(0, weight=1)
        content.grid_rowconfigure(1, weight=0)  # 説明エリア用
        
        self.output_text = scrolledtext.ScrolledText(content, wrap=tk.WORD, state=tk.DISABLED, font=FONTS["code"], bg=COLOR_PALETTE["bg_log"], fg=COLOR_PALETTE["text_primary"], relief='flat', padx=SPACING["md"], pady=SPACING["md"])
        self.output_text.grid(row=0, column=0, sticky='nsew')
        
        # 説明エリア（固定表示）
        self.info_frame = tk.Frame(content, bg=COLOR_PALETTE["bg_elevated"], padx=SPACING["lg"], pady=SPACING["md"])
        self.info_frame.grid(row=1, column=0, sticky='ew')
        
        info_text = (
            f"📋 {t('info_workflow_title')}\n"
            f"  {t('info_workflow_1')}\n"
            f"  {t('info_workflow_2')}\n"
            f"  {t('info_workflow_3')}\n"
            "\n"
            f"🎯 {t('info_use_cases')}\n"
            f"💡 {t('info_format_note')}\n"
            "\n"
            f"⚠️ {t('info_caution_title')}\n"
            f"  • {t('info_caution_1')}\n"
            f"  • {t('info_caution_2')}\n"
            f"  • {t('info_caution_3')}"
        )
        
        info_label = tk.Label(self.info_frame, text=info_text, font=FONTS["body"], 
                              bg=COLOR_PALETTE["bg_elevated"], fg=COLOR_PALETTE["text_secondary"],
                              justify='left', anchor='w')
        info_label.pack(fill='x')
        
        self._show_welcome()
    
    def _show_welcome(self):
        tier = self.license_manager.get_tier_info()
        self._update_output(f"{t('welcome_title')}\n{APP_NAME} v{APP_VERSION} ({tier['name']})\n\n", clear=True)
        self._update_output(f"📤 {t('mode_extract')}: {t('msg_extract_desc')}\n")
        self._update_output(f"📥 {t('mode_update')}: {t('msg_update_desc')}\n\n")
        
        limit = tier['update_slide_limit']
        if limit:
            self._update_output(f"⚠️ Free: {t('msg_update_limit', limit)}\n\n")
    
    # === Output helpers ===
    def _update_output(self, text, clear=False):
        self.output_text.configure(state=tk.NORMAL)
        if clear:
            self.output_text.delete('1.0', tk.END)
            self.log_buffer = []
        self.output_text.insert(tk.END, text)
        self.output_text.see(tk.END)
        self.output_text.configure(state=tk.DISABLED)
        self.log_buffer.append(text)
    
    def _update_output_safe(self, text, clear=False):
        self.root.after(0, lambda: self._update_output(text, clear))
    
    def _update_status(self, text, color=None):
        self.status_label.configure(text=text)
        if color:
            self.status_label.configure(foreground=color)
    
    def _update_status_safe(self, text, color=None):
        self.root.after(0, lambda: self._update_status(text, color))
    
    def _log(self, msg, level="info"):
        timestamp = datetime.now().strftime("%H:%M:%S")
        prefix = {"error": "❌ ", "warning": "⚠️ ", "success": "✅ "}.get(level, "")
        self._update_output_safe(f"[{timestamp}] {prefix}{msg}\n")
    
    def _show_result_summary(self, title, stats):
        """Show prominent result summary at top of log"""
        separator = "═" * 40
        summary = f"\n{separator}\n"
        summary += f"  {title}\n"
        summary += f"{separator}\n"
        for key, value in stats.items():
            summary += f"  {key}: {value}\n"
        summary += f"{separator}\n\n"
        self._update_output_safe(summary)
    
    def _start_progress(self):
        self.progress.start(10)
        self.processing = True
        self.cancel_requested = False
        self.root.after(0, lambda: self.cancel_btn.configure(state='normal'))
    
    def _stop_progress(self):
        self.progress.stop()
        self.processing = False
        self.root.after(0, lambda: self.cancel_btn.configure(state='disabled'))
    
    def _cancel(self):
        if self.processing:
            self.cancel_requested = True
            self._log(t('msg_cancel_requested'), "warning")
    
    def _clear_output(self):
        self.output_text.configure(state=tk.NORMAL)
        self.output_text.delete('1.0', tk.END)
        self.output_text.configure(state=tk.DISABLED)
        self.log_buffer = []
        self._update_status(t('status_waiting'), COLOR_PALETTE["text_muted"])
    
    def _copy_output(self):
        content = self.output_text.get('1.0', tk.END).strip()
        if content:
            self.root.clipboard_clear()
            self.root.clipboard_append(content)
            messagebox.showinfo(t('dialog_complete'), t('msg_copied'))
    
    def _copy_support_log(self):
        """Copy support log to clipboard"""
        tier = self.license_manager.get_tier_info()
        log_content = f"=== {APP_NAME} Support Info ===\n"
        log_content += f"Version: {APP_VERSION}\n"
        log_content += f"License: {tier['name']}\n"
        log_content += f"Language: {get_language()}\n"
        log_content += f"Time: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}\n"
        log_content += f"OS: {os.name}\n\n"
        log_content += f"=== Log ===\n"
        log_content += ''.join(self.log_buffer[-50:])
        
        self.root.clipboard_clear()
        self.root.clipboard_append(log_content)
        messagebox.showinfo(t('dialog_complete'), t('msg_support_copied'))
    
    # === Mode switching ===
    def _switch_extract(self):
        self.current_mode = "extract"
        self.mode_dot.configure(fg=COLOR_PALETTE["brand_primary"])
        self.mode_label.configure(text=t('mode_extract'))
        self.extract_btn.configure(style='Primary.Modern.TButton')
        self.update_btn.configure(style='Outline.Modern.TButton')
        self.update_frame.grid_remove()
        self.extract_frame.grid(row=0, column=0, sticky='nsew')
        self._update_status(f"{t('mode_extract')} - {t('status_waiting')}", COLOR_PALETTE["text_muted"])
    
    def _switch_update(self):
        self.current_mode = "update"
        self.mode_dot.configure(fg=COLOR_PALETTE["brand_update"])
        self.mode_label.configure(text=t('mode_update'))
        self.extract_btn.configure(style='Outline.Modern.TButton')
        self.update_btn.configure(style='Update.Modern.TButton')
        self.extract_frame.grid_remove()
        self.update_frame.grid(row=0, column=0, sticky='nsew')
        self._update_status(f"{t('mode_update')} - {t('status_waiting')}", COLOR_PALETTE["text_muted"])
    
    # === Language ===
    def _change_language(self, lang: str):
        """Change language and refresh UI"""
        if lang != get_language():
            self.config_manager.set('language', lang)
            set_language(lang)
            # UIを再構築
            self._refresh_ui()
            messagebox.showinfo(t('dialog_complete'), t('lang_restart_msg'))
    
    def _change_font_size(self, size_preset: str):
        """Change font size and refresh UI"""
        global FONTS
        current = self.config_manager.get('font_size', 'medium')
        if size_preset != current:
            self.config_manager.set('font_size', size_preset)
            # グローバルFONTSを更新
            new_fonts = get_fonts(size_preset)
            FONTS.clear()
            FONTS.update(new_fonts)
            # スタイルを再適用してUIを再構築
            self._setup_styles()
            self._refresh_ui()
            messagebox.showinfo(t('dialog_complete'), t('font_size_restart_msg'))
    
    def _center_dialog(self, dialog, width, height):
        """ダイアログをアプリウィンドウの中央に配置"""
        dialog.update_idletasks()
        # 親ウィンドウの位置とサイズを取得
        parent_x = self.root.winfo_x()
        parent_y = self.root.winfo_y()
        parent_width = self.root.winfo_width()
        parent_height = self.root.winfo_height()
        # 中央位置を計算
        x = parent_x + (parent_width - width) // 2
        y = parent_y + (parent_height - height) // 2
        dialog.geometry(f"{width}x{height}+{x}+{y}")
    
    def _show_language_dialog(self):
        """Show language selection dialog on first launch"""
        dialog = tk.Toplevel(self.root)
        dialog.title(t('lang_select_title'))
        self._center_dialog(dialog, 400, 250)
        dialog.resizable(False, False)
        dialog.transient(self.root)
        dialog.grab_set()
        
        frame = ttk.Frame(dialog, padding=SPACING["xl"])
        frame.pack(fill='both', expand=True)
        
        ttk.Label(frame, text="🌐 Select Language / 言語を選択", font=FONTS["heading"]).pack(pady=(0, SPACING["lg"]))
        
        lang_var = tk.StringVar(value='en')
        
        ttk.Radiobutton(frame, text="English", variable=lang_var, value='en').pack(anchor='w', pady=SPACING["sm"])
        ttk.Radiobutton(frame, text="日本語 (Japanese)", variable=lang_var, value='ja').pack(anchor='w', pady=SPACING["sm"])
        
        def apply_lang():
            selected = lang_var.get()
            set_language(selected)
            self.config_manager.set('language', selected)
            self.config_manager.set('first_launch', False)
            dialog.destroy()
            self._refresh_ui()
            self._show_welcome_dialog()
        
        ttk.Button(frame, text="Continue / 続ける", command=apply_lang).pack(pady=(SPACING["xl"], 0))
    
    # === Utility ===
    def clean_text(self, text):
        """Clean and normalize text for extraction"""
        if text is None:
            return ""
        # Remove control characters (keep newlines)
        text = re.sub(r'[\x00-\x08\x0B-\x0C\x0E-\x1F]', '', text)
        # Normalize line endings (PowerPoint uses \v for line breaks within shapes)
        text = text.replace('\r\n', '\n').replace('\r', '\n').replace('\v', '\n')
        # Normalize whitespace
        text = text.replace('\u00A0', ' ')  # Non-breaking space → space
        text = text.replace('\u3000', ' ')  # Full-width space → space
        # Remove trailing whitespace from each line
        lines = [line.rstrip() for line in text.split('\n')]
        text = '\n'.join(lines)
        # Escape Excel formula characters
        if text.startswith(('=', '+', '-', '@')):
            text = "'" + text
        return text
    
    def _normalize_for_compare(self, text):
        """Normalize text for comparison to avoid false positives"""
        if text is None:
            return ""
        # Remove leading quote added by clean_text
        if text.startswith("'") and len(text) > 1 and text[1] in ('=', '+', '-', '@'):
            text = text[1:]
        # Normalize line endings (including vertical tab used by PowerPoint)
        text = text.replace('\r\n', '\n').replace('\r', '\n').replace('\v', '\n')
        # Normalize whitespace
        text = text.replace('\u00A0', ' ')  # Non-breaking space → space
        text = text.replace('\u3000', ' ')  # Full-width space → space
        text = text.replace('\t', ' ')      # Tab → space
        # Remove trailing/leading whitespace from each line
        lines = [line.strip() for line in text.split('\n')]
        text = '\n'.join(lines)
        # Remove trailing whitespace
        text = text.rstrip()
        return text
    
    def _texts_are_equal(self, old_text, new_text):
        """Compare texts with normalization"""
        return self._normalize_for_compare(old_text) == self._normalize_for_compare(new_text)
    
    def _normalize_excel_text(self, text):
        """Normalize text loaded from Excel/text files for update"""
        if text is None:
            return ""
        # Remove leading quote added during extraction
        if text.startswith("'") and len(text) > 1 and text[1] in ('=', '+', '-', '@'):
            text = text[1:]
        # Normalize line endings (Excel uses \r\n, some editors use \r)
        text = text.replace('\r\n', '\n').replace('\r', '\n')
        # Excel internal line break (Alt+Enter) is char(10), already \n
        # Normalize whitespace
        text = text.replace('\u00A0', ' ')  # Non-breaking space → space
        text = text.replace('\u3000', ' ')  # Full-width space → space
        # Remove trailing whitespace from each line
        lines = [line.rstrip() for line in text.split('\n')]
        text = '\n'.join(lines)
        return text
    
    def _safe_excel_value(self, value):
        """Make value safe for Excel output (prevent formula injection)"""
        if value is None:
            return ""
        text = str(value)
        # Remove/replace problematic characters
        text = text.replace('\t', ' ')  # Tab → space
        text = text.replace('\v', ' ')  # Vertical tab → space
        # Escape formula-triggering characters
        if text.startswith(('=', '+', '-', '@')):
            text = "'" + text
        return text
    
    def get_shape_type(self, shape):
        try:
            if shape.is_placeholder:
                return get_placeholder_type(shape.placeholder_format.type)
            elif hasattr(shape, "has_table") and shape.has_table:
                return "Table" if get_language() == 'en' else "表"
            elif shape.shape_type == 1:
                return "TextBox" if get_language() == 'en' else "テキストボックス"
            return "Other" if get_language() == 'en' else "その他"
        except:
            return "Unknown" if get_language() == 'en' else "不明"
    
    def _create_backup(self, path: str) -> Optional[str]:
        if not self.license_manager.can_use_feature('auto_backup') or not self.auto_backup_var.get():
            return None
        try:
            backup_dir = Path(path).parent / "backup"
            backup_dir.mkdir(exist_ok=True)
            backup_name = f"{Path(path).stem}_{datetime.now().strftime('%Y%m%d_%H%M%S')}{Path(path).suffix}"
            backup_path = backup_dir / backup_name
            shutil.copy2(path, backup_path)
            self._log(t('msg_backup_created', backup_name))
            return str(backup_path)
        except Exception as e:
            self._log(t('msg_backup_failed', translate_error(e)), "warning")
            return None
    
    # === Extract ===
    def extract_from_ppt(self, path: str, include_notes: bool = False) -> Tuple[List, Dict]:
        try:
            prs = pptx.Presentation(path)
            data = []
            meta = {'file_name': os.path.basename(path), 'slide_count': len(prs.slides), 'time': datetime.now().strftime("%Y-%m-%d %H:%M:%S")}
            
            for slide_num, slide in enumerate(prs.slides, 1):
                if self.cancel_requested:
                    break
                for shape in slide.shapes:
                    try:
                        sid = str(shape.shape_id)
                        stype = self.get_shape_type(shape)
                        
                        if hasattr(shape, "text") and shape.text.strip():
                            row = [slide_num, sid, stype, self.clean_text(shape.text)]
                            if self.include_metadata_var.get():
                                row.extend([meta['file_name'], meta['time']])
                            data.append(row)
                        
                        if hasattr(shape, "has_table") and shape.has_table:
                            for r, row in enumerate(shape.table.rows):
                                for c, cell in enumerate(row.cells):
                                    if cell.text.strip():
                                        cell_row = [slide_num, f"{sid}_t{r}_{c}", f"表({r+1},{c+1})", self.clean_text(cell.text)]
                                        if self.include_metadata_var.get():
                                            cell_row.extend([meta['file_name'], meta['time']])
                                        data.append(cell_row)
                    except Exception as e:
                        self._log(f"スライド{slide_num}でエラー: {translate_error(e)}", "warning")
                
                # Extract speaker notes (Pro feature)
                if include_notes:
                    try:
                        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                            notes_text = slide.notes_slide.notes_text_frame.text.strip()
                            if notes_text:
                                notes_row = [slide_num, "notes", t('type_notes'), self.clean_text(notes_text)]
                                if self.include_metadata_var.get():
                                    notes_row.extend([meta['file_name'], meta['time']])
                                data.append(notes_row)
                    except Exception as e:
                        pass  # Skip if notes extraction fails
            
            return data, meta
        except Exception as e:
            save_error_log(e, f"extract_from_ppt: {path}")
            self._log(f"読み込みエラー: {translate_error(e)}", "error")
            return [], {}
    
    def save_to_file(self, data: List, path: str, fmt: str = "txt", meta: Dict = None) -> bool:
        try:
            if fmt == "excel":
                return self._save_excel(data, path)
            elif fmt == "json":
                return self._save_json(data, path, meta)
            # tab-separated text
            with open(path, 'w', encoding='utf-8', newline='') as f:
                w = csv.writer(f, delimiter='\t')
                headers = [t('header_slide'), t('header_id'), t('header_type'), t('header_text')]
                if self.include_metadata_var.get():
                    headers.extend([t('header_filename'), t('header_datetime')])
                w.writerow(headers)
                w.writerows(data)
            return True
        except Exception as e:
            save_error_log(e, f"save_to_file: {path}")
            self._log(f"保存エラー: {translate_error(e)}", "error")
            return False
    
    def _save_json(self, data: List, path: str, meta: Dict = None) -> bool:
        """Save as JSON with hierarchical structure for AI consumption"""
        try:
            # Group by slide number
            slides_dict = {}
            for row in data:
                slide_num = row[0]
                if slide_num not in slides_dict:
                    slides_dict[slide_num] = []
                shape_data = {
                    "id": row[1],
                    "type": row[2],
                    "text": row[3]
                }
                slides_dict[slide_num].append(shape_data)
            
            # Build structured output
            output = {
                "file": meta.get('file_name', '') if meta else os.path.basename(path),
                "extracted_at": meta.get('time', datetime.now().strftime("%Y-%m-%d %H:%M:%S")) if meta else datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
                "slide_count": len(slides_dict),
                "total_items": len(data),
                "slides": [
                    {
                        "number": slide_num,
                        "shapes": shapes
                    }
                    for slide_num, shapes in sorted(slides_dict.items())
                ]
            }
            
            with open(path, 'w', encoding='utf-8') as f:
                json.dump(output, f, ensure_ascii=False, indent=2)
            return True
        except Exception as e:
            save_error_log(e, f"_save_json: {path}")
            self._log(f"JSON保存エラー: {translate_error(e)}", "error")
            return False
    
    def _save_excel(self, data: List, path: str) -> bool:
        try:
            wb = openpyxl.Workbook()
            ws = wb.active
            ws.title = "Data"
            headers = [t('header_slide'), t('header_id'), t('header_type'), t('header_text')]
            if self.include_metadata_var.get():
                headers.extend([t('header_filename'), t('header_datetime')])
            for col, h in enumerate(headers, 1):
                ws.cell(1, col, h).font = XLFont(bold=True)
            for r, row in enumerate(data, 2):
                for c, val in enumerate(row, 1):
                    ws.cell(r, c, val)
            for col in ws.columns:
                max_len = max((len(str(cell.value or "")) for cell in col), default=10)
                ws.column_dimensions[col[0].column_letter].width = min(max_len + 2, 50)
            wb.save(path)
            return True
        except Exception as e:
            save_error_log(e, f"_save_excel: {path}")
            self._log(f"Excel保存エラー: {translate_error(e)}", "error")
            return False
    
    def _extract_single(self):
        if self.processing:
            return messagebox.showwarning("処理中", "現在処理中です")
        path = filedialog.askopenfilename(title="PowerPointファイルを選択", filetypes=[("PowerPoint", "*.pptx")], initialdir=self.config_manager.get('last_directory'))
        if not path:
            return
        self.config_manager.set('last_directory', os.path.dirname(path))
        
        # Check if notes should be included (Pro feature)
        include_notes = self.include_notes_var.get() if hasattr(self, 'include_notes_var') else False
        
        def run():
            try:
                self._start_progress()
                self._update_status_safe(f"{ICONS['processing']} 処理中...", COLOR_PALETTE["info"])
                self._update_output_safe(f"\n{ICONS['file']} 処理開始: {os.path.basename(path)}\n", clear=True)
                
                data, meta = self.extract_from_ppt(path, include_notes=include_notes)
                if self.cancel_requested:
                    return self._log("キャンセルされました", "warning")
                
                if data:
                    fmt = self.output_format_var.get()
                    ext = ".xlsx" if fmt == "excel" else (".json" if fmt == "json" else ".txt")
                    out = os.path.splitext(path)[0] + "_抽出" + ext
                    if self.save_to_file(data, out, fmt, meta):
                        notes_msg = " (ノート含む)" if include_notes else ""
                        self._show_result_summary(f"{ICONS['success']} 抽出完了{notes_msg}", {
                            "抽出件数": f"{len(data)} 件",
                            "スライド": f"{meta.get('slide_count', 0)} 枚",
                            "出力ファイル": out
                        })
                        self._update_status_safe(f"{ICONS['success']} 完了: {len(data)}件", COLOR_PALETTE["success"])
                else:
                    self._log("テキストが見つかりませんでした", "warning")
            except Exception as e:
                save_error_log(e, "_extract_single")
                self._log(f"エラー: {translate_error(e)}", "error")
            finally:
                self._stop_progress()
        
        threading.Thread(target=run, daemon=True).start()
    
    def _extract_batch(self):
        if self.processing:
            return messagebox.showwarning("処理中", "現在処理中です")
        folder = filedialog.askdirectory(title="フォルダを選択", initialdir=self.config_manager.get('last_directory'))
        if not folder:
            return
        self.config_manager.set('last_directory', folder)
        
        # Check if notes should be included (Pro feature)
        include_notes = self.include_notes_var.get() if hasattr(self, 'include_notes_var') else False
        
        def run():
            try:
                self._start_progress()
                self._update_status_safe(f"{ICONS['processing']} バッチ処理中...", COLOR_PALETTE["info"])
                self._update_output_safe(f"\n{ICONS['folder']} フォルダ処理: {folder}\n", clear=True)
                
                # Skip temp files (~$xxx.pptx)
                files = [f for f in Path(folder).glob("*.pptx") if not f.name.startswith("~$")]
                if not files:
                    return self._log("PPTXファイルが見つかりません", "warning")
                
                self._log(f"発見: {len(files)}件")
                total, success = 0, 0
                created_files = []
                
                for i, f in enumerate(files, 1):
                    if self.cancel_requested:
                        break
                    self._log(f"[{i}/{len(files)}] {f.name}")
                    data, meta = self.extract_from_ppt(str(f), include_notes=include_notes)
                    if data:
                        fmt = self.output_format_var.get()
                        ext = ".xlsx" if fmt == "excel" else (".json" if fmt == "json" else ".txt")
                        out = str(f.with_suffix('')) + "_抽出" + ext
                        if self.save_to_file(data, out, fmt, meta):
                            total += len(data)
                            success += 1
                            created_files.append(out)
                            self._log(f"  → 作成: {os.path.basename(out)}")
                
                notes_msg = " (ノート含む)" if include_notes else ""
                self._show_result_summary(f"{ICONS['success']} バッチ抽出完了{notes_msg}", {
                    "処理ファイル": f"{success}/{len(files)} ファイル",
                    "抽出件数": f"{total} 件",
                })
                self._log("作成されたファイル:")
                for cf in created_files:
                    self._log(f"  {cf}")
                self._update_status_safe(f"{ICONS['success']} 完了", COLOR_PALETTE["success"])
            except Exception as e:
                save_error_log(e, "_extract_batch")
                self._log(f"エラー: {translate_error(e)}", "error")
            finally:
                self._stop_progress()
        
        threading.Thread(target=run, daemon=True).start()
    
    # === Pro Features: Font Analysis ===
    def _font_analysis(self):
        """Analyze fonts used in PowerPoint file (Pro feature)"""
        if self.processing:
            return messagebox.showwarning("処理中", "現在処理中です")
        path = filedialog.askopenfilename(title="PowerPointファイルを選択", filetypes=[("PowerPoint", "*.pptx")], initialdir=self.config_manager.get('last_directory'))
        if not path:
            return
        self.config_manager.set('last_directory', os.path.dirname(path))
        
        def run():
            try:
                self._start_progress()
                self._update_status_safe(f"{ICONS['processing']} フォント分析中...", COLOR_PALETTE["info"])
                self._update_output_safe(f"\n📊 {t('font_analysis_title')}: {os.path.basename(path)}\n\n", clear=True)
                
                prs = pptx.Presentation(path)
                data = []
                font_stats = {}
                
                for slide_num, slide in enumerate(prs.slides, 1):
                    if self.cancel_requested:
                        break
                    for shape in slide.shapes:
                        try:
                            if hasattr(shape, "text_frame"):
                                for para in shape.text_frame.paragraphs:
                                    for run in para.runs:
                                        if run.text.strip():
                                            font_name = run.font.name or "(default)"
                                            font_size = f"{run.font.size.pt:.1f}pt" if run.font.size else "(default)"
                                            is_bold = "Yes" if run.font.bold else "No"
                                            
                                            # Track font usage
                                            font_stats[font_name] = font_stats.get(font_name, 0) + 1
                                            
                                            row = [slide_num, str(shape.shape_id), self.get_shape_type(shape),
                                                   self._safe_excel_value(run.text[:30]), 
                                                   self._safe_excel_value(font_name), 
                                                   font_size, is_bold]
                                            data.append(row)
                            
                            if hasattr(shape, "has_table") and shape.has_table:
                                for r, trow in enumerate(shape.table.rows):
                                    for c, cell in enumerate(trow.cells):
                                        for para in cell.text_frame.paragraphs:
                                            for run in para.runs:
                                                if run.text.strip():
                                                    font_name = run.font.name or "(default)"
                                                    font_size = f"{run.font.size.pt:.1f}pt" if run.font.size else "(default)"
                                                    is_bold = "Yes" if run.font.bold else "No"
                                                    font_stats[font_name] = font_stats.get(font_name, 0) + 1
                                                    row = [slide_num, f"{shape.shape_id}_t{r}_{c}", f"表({r+1},{c+1})",
                                                           self._safe_excel_value(run.text[:30]), 
                                                           self._safe_excel_value(font_name), 
                                                           font_size, is_bold]
                                                    data.append(row)
                        except Exception as e:
                            pass
                
                if self.cancel_requested:
                    return self._log("キャンセルされました", "warning")
                
                # Show font usage summary in log
                self._log(f"\n{'═'*40}")
                self._log(f"  {t('font_summary_title')}")
                self._log(f"{'═'*40}")
                
                # フォント×サイズの組み合わせで集計
                font_size_stats = {}
                for row in data:
                    key = (row[4], row[5])  # (font_name, font_size)
                    font_size_stats[key] = font_size_stats.get(key, 0) + 1
                
                total = sum(font_size_stats.values())
                summary_data = []
                for (font, size), count in sorted(font_size_stats.items(), key=lambda x: (-x[1], x[0][0])):
                    pct = count * 100 / total if total else 0
                    self._log(f"  {font} / {size}: {count}件 ({pct:.1f}%)")
                    summary_data.append([font, size, count, f"{pct:.1f}%"])
                
                self._log(f"{'═'*40}\n")
                
                # Save to Excel with 2 sheets (Details + Summary)
                if data:
                    out = os.path.splitext(path)[0] + "_フォント診断.xlsx"
                    wb = openpyxl.Workbook()
                    
                    # Sheet 1: Details
                    ws1 = wb.active
                    ws1.title = "詳細"
                    headers = [t('header_slide'), t('header_id'), t('header_type'), 
                               t('header_text'), t('header_font'), t('header_size'), t('header_bold')]
                    ws1.append(headers)
                    for row in data:
                        ws1.append(row)
                    # Style header
                    for cell in ws1[1]:
                        cell.font = XLFont(bold=True)
                        cell.fill = PatternFill(start_color="E0E0E0", end_color="E0E0E0", fill_type="solid")
                    # Adjust column widths
                    ws1.column_dimensions['A'].width = 10
                    ws1.column_dimensions['B'].width = 12
                    ws1.column_dimensions['C'].width = 15
                    ws1.column_dimensions['D'].width = 35
                    ws1.column_dimensions['E'].width = 20
                    ws1.column_dimensions['F'].width = 10
                    ws1.column_dimensions['G'].width = 8
                    
                    # Sheet 2: Summary (Font × Size)
                    ws2 = wb.create_sheet("サマリー")
                    ws2.append([t('header_font'), t('header_size'), "件数", "割合"])
                    for row in summary_data:
                        ws2.append(row)
                    # Style header
                    for cell in ws2[1]:
                        cell.font = XLFont(bold=True)
                        cell.fill = PatternFill(start_color="E0E0E0", end_color="E0E0E0", fill_type="solid")
                    # Adjust column widths
                    ws2.column_dimensions['A'].width = 25
                    ws2.column_dimensions['B'].width = 12
                    ws2.column_dimensions['C'].width = 10
                    ws2.column_dimensions['D'].width = 10
                    
                    wb.save(out)
                    self._log(f"出力ファイル: {out}")
                    self._update_status_safe(f"{ICONS['success']} 完了: {len(data)}件", COLOR_PALETTE["success"])
                else:
                    self._log("フォント情報が見つかりませんでした", "warning")
            except Exception as e:
                save_error_log(e, "_font_analysis")
                self._log(f"エラー: {translate_error(e)}", "error")
            finally:
                self._stop_progress()
        
        threading.Thread(target=run, daemon=True).start()
    
    # === Pro Features: Conditional Extract ===
    def _conditional_extract(self):
        """Extract only text matching a condition (Pro feature)"""
        if self.processing:
            return messagebox.showwarning("処理中", "現在処理中です")
        
        # Get filter condition from user
        condition = simpledialog.askstring(t('conditional_title'), t('conditional_prompt'))
        if not condition:
            return
        
        path = filedialog.askopenfilename(title="PowerPointファイルを選択", filetypes=[("PowerPoint", "*.pptx")], initialdir=self.config_manager.get('last_directory'))
        if not path:
            return
        self.config_manager.set('last_directory', os.path.dirname(path))
        
        include_notes = self.include_notes_var.get() if hasattr(self, 'include_notes_var') else False
        
        def run():
            try:
                self._start_progress()
                self._update_status_safe(f"{ICONS['processing']} 条件抽出中...", COLOR_PALETTE["info"])
                self._update_output_safe(f"\n🔍 {t('conditional_title')}: {os.path.basename(path)}\n", clear=True)
                self._log(f"条件: {condition}")
                
                # Compile regex pattern
                try:
                    pattern = re.compile(condition, re.IGNORECASE)
                except re.error:
                    # If not valid regex, use as simple substring
                    pattern = re.compile(re.escape(condition), re.IGNORECASE)
                
                # Extract all data first
                data, meta = self.extract_from_ppt(path, include_notes=include_notes)
                
                if self.cancel_requested:
                    return self._log("キャンセルされました", "warning")
                
                # Filter by condition
                filtered = [row for row in data if pattern.search(str(row[3]))]
                
                if filtered:
                    fmt = self.output_format_var.get()
                    ext = ".xlsx" if fmt == "excel" else (".json" if fmt == "json" else ".txt")
                    out = os.path.splitext(path)[0] + "_条件抽出" + ext
                    if self.save_to_file(filtered, out, fmt, meta):
                        self._show_result_summary(f"{ICONS['success']} 条件抽出完了", {
                            "条件": condition,
                            "マッチ件数": f"{len(filtered)} 件 (全{len(data)}件中)",
                            "出力ファイル": out
                        })
                        self._update_status_safe(f"{ICONS['success']} 完了: {len(filtered)}件", COLOR_PALETTE["success"])
                else:
                    self._log(t('conditional_no_match'), "warning")
            except Exception as e:
                save_error_log(e, "_conditional_extract")
                self._log(f"エラー: {translate_error(e)}", "error")
            finally:
                self._stop_progress()
        
        threading.Thread(target=run, daemon=True).start()
    
    # === Update ===
    def _load_excel_updates(self, path: str) -> Dict:
        updates = {}
        try:
            wb = openpyxl.load_workbook(path)
            ws = wb.active
            headers = [c.value for c in ws[1]]
            try:
                si, oi, ti = headers.index("スライド番号"), headers.index("オブジェクトID"), headers.index("テキスト内容")
            except ValueError:
                self._log("ヘッダー形式が不正です", "error")
                return {}
            for row in list(ws.rows)[1:]:
                try:
                    sn = int(row[si].value) if row[si].value else None
                    oid = str(row[oi].value) if row[oi].value else None
                    txt = str(row[ti].value) if row[ti].value else ""
                    if txt == "None":
                        txt = ""
                    # Normalize text from Excel (handles Excel line breaks, etc.)
                    txt = self._normalize_excel_text(txt)
                    if sn and oid:
                        updates[(sn, oid)] = txt
                except:
                    pass
            return updates
        except Exception as e:
            save_error_log(e, f"_load_excel_updates: {path}")
            self._log(f"Excel読み込みエラー: {translate_error(e)}", "error")
            return {}
    
    def _load_txt_updates(self, path: str) -> Dict:
        updates = {}
        try:
            with open(path, 'r', encoding='utf-8', newline='') as f:
                first = f.readline()
                f.seek(0)
                delim = '\t' if '\t' in first else ','
                reader = csv.reader(f, delimiter=delim)
                headers = next(reader)
                try:
                    si, oi, ti = headers.index("スライド番号"), headers.index("オブジェクトID"), headers.index("テキスト内容")
                except ValueError:
                    self._log("ヘッダー形式が不正です", "error")
                    return {}
                for row in reader:
                    try:
                        if len(row) > max(si, oi, ti):
                            sn = int(row[si]) if row[si] else None
                            oid = row[oi] if row[oi] else None
                            txt = row[ti] if len(row) > ti else ""
                            # Normalize text from file
                            txt = self._normalize_excel_text(txt)
                            if sn and oid:
                                updates[(sn, oid)] = txt
                    except:
                        pass
            return updates
        except Exception as e:
            save_error_log(e, f"_load_txt_updates: {path}")
            self._log(f"テキスト読み込みエラー: {translate_error(e)}", "error")
            return {}
    
    def _load_json_updates(self, path: str) -> Dict:
        """Load updates from JSON file"""
        updates = {}
        try:
            with open(path, 'r', encoding='utf-8') as f:
                data = json.load(f)
            
            if not isinstance(data, list):
                self._log("JSON形式が不正です（配列が必要）", "error")
                return {}
            
            for item in data:
                try:
                    # Support both Japanese and English keys
                    sn = item.get('スライド番号') or item.get('slide')
                    oid = item.get('オブジェクトID') or item.get('id')
                    txt = item.get('テキスト内容') or item.get('text', '')
                    
                    if sn and oid:
                        txt = self._normalize_excel_text(str(txt))
                        updates[(int(sn), str(oid))] = txt
                except:
                    pass
            
            return updates
        except json.JSONDecodeError as e:
            self._log(f"JSON解析エラー: {e}", "error")
            return {}
        except Exception as e:
            save_error_log(e, f"_load_json_updates: {path}")
            self._log(f"JSON読み込みエラー: {translate_error(e)}", "error")
            return {}
    
    def _update_ppt(self, ppt_path: str, updates: Dict, preview: bool = False) -> Tuple[int, int, List]:
        limit = self.license_manager.get_update_limit()
        quality = self.license_manager.get_format_quality()
        preserve = True  # 書式維持は常に有効（品質はライセンスで分岐）
        
        self.presentation = pptx.Presentation(ppt_path)
        updated, skipped = 0, 0
        changes = []
        
        for slide_idx, slide in enumerate(self.presentation.slides, 1):
            if self.cancel_requested:
                break
            if limit and slide_idx > limit:
                skipped += len([k for k in updates if k[0] == slide_idx])
                continue
            
            for shape in slide.shapes:
                try:
                    sid = str(shape.shape_id)
                    key = (slide_idx, sid)
                    
                    if key in updates and hasattr(shape, "text"):
                        new_txt = updates[key]
                        old_txt = shape.text
                        # Use normalized comparison to avoid false positives
                        if not self._texts_are_equal(old_txt, new_txt):
                            changes.append({'slide': slide_idx, 'id': sid, 'old': old_txt[:50], 'new': new_txt[:50]})
                            if not preview:
                                # Remove leading quote before updating
                                update_txt = self._normalize_for_compare(new_txt)
                                if preserve and quality == 'advanced':
                                    self._update_shape_advanced(shape, update_txt)
                                elif preserve:
                                    self._update_shape_basic(shape, update_txt)
                                else:
                                    shape.text = update_txt
                                updated += 1
                    
                    if hasattr(shape, "has_table") and shape.has_table:
                        for r, row in enumerate(shape.table.rows):
                            for c, cell in enumerate(row.cells):
                                cid = f"{sid}_t{r}_{c}"
                                ckey = (slide_idx, cid)
                                if ckey in updates:
                                    new_txt = updates[ckey]
                                    old_txt = cell.text
                                    # Use normalized comparison
                                    if not self._texts_are_equal(old_txt, new_txt):
                                        changes.append({'slide': slide_idx, 'id': cid, 'old': old_txt[:30], 'new': new_txt[:30]})
                                        if not preview:
                                            update_txt = self._normalize_for_compare(new_txt)
                                            if preserve and quality == 'advanced':
                                                self._update_cell_advanced(cell, update_txt)
                                            elif preserve:
                                                self._update_cell_basic(cell, update_txt)
                                            else:
                                                cell.text = update_txt
                                            updated += 1
                except Exception as e:
                    self._log(f"スライド{slide_idx}更新エラー: {translate_error(e)}", "warning")
                    skipped += 1
            
            # Update speaker notes (Pro feature)
            notes_key = (slide_idx, "notes")
            if notes_key in updates:
                try:
                    new_txt = updates[notes_key]
                    # Get or create notes slide
                    if slide.has_notes_slide:
                        notes_slide = slide.notes_slide
                    else:
                        # Create notes slide if it doesn't exist
                        notes_slide = slide.notes_slide  # This auto-creates notes slide
                    
                    notes_frame = notes_slide.notes_text_frame
                    old_txt = notes_frame.text if notes_frame.text else ""
                    if not self._texts_are_equal(old_txt, new_txt):
                        changes.append({'slide': slide_idx, 'id': 'notes', 'old': old_txt[:50], 'new': new_txt[:50]})
                        if not preview:
                            update_txt = self._normalize_for_compare(new_txt)
                            notes_frame.text = update_txt
                            updated += 1
                except Exception as e:
                    self._log(f"スライド{slide_idx}ノート更新エラー: {translate_error(e)}", "warning")
        
        return updated, skipped, changes
    
    def _update_shape_basic(self, shape, new_text: str):
        """基本書式維持（Free版）"""
        try:
            tf = shape.text_frame
            if tf.paragraphs:
                p = tf.paragraphs[0]
                font_info = {}
                if p.runs:
                    run = p.runs[0]
                    font_info = {'name': run.font.name, 'size': run.font.size, 'bold': run.font.bold}
                tf.text = new_text
                if font_info and tf.paragraphs and tf.paragraphs[0].runs:
                    run = tf.paragraphs[0].runs[0]
                    if font_info.get('name'):
                        run.font.name = font_info['name']
                    if font_info.get('size'):
                        run.font.size = font_info['size']
                    run.font.bold = font_info.get('bold')
            else:
                shape.text = new_text
        except:
            shape.text = new_text
    
    def _update_shape_advanced(self, shape, new_text: str):
        """高度書式維持（Standard/Pro版）"""
        try:
            tf = shape.text_frame
            styles = []
            for p in tf.paragraphs:
                ps = {'alignment': p.alignment, 'level': p.level, 'line_spacing': p.line_spacing, 'space_before': p.space_before, 'space_after': p.space_after, 'runs': []}
                for run in p.runs:
                    try:
                        color_data, color_type = None, None
                        if hasattr(run.font, 'color') and run.font.color and hasattr(run.font.color, 'type') and run.font.color.type == 1:
                            color_data = run.font.color.rgb
                            color_type = 1
                        ps['runs'].append({'name': run.font.name, 'size': run.font.size, 'bold': run.font.bold, 'italic': run.font.italic, 'underline': run.font.underline, 'color': color_data, 'color_type': color_type, 'lang': getattr(run.font, 'language_id', None)})
                    except:
                        pass
                styles.append(ps)
            
            tf.text = new_text
            
            for i, p in enumerate(tf.paragraphs):
                if i < len(styles):
                    s = styles[i]
                    try:
                        p.alignment, p.level, p.line_spacing, p.space_before, p.space_after = s['alignment'], s['level'], s['line_spacing'], s['space_before'], s['space_after']
                    except:
                        pass
                    if p.runs and s['runs']:
                        run, rs = p.runs[0], s['runs'][0]
                        try:
                            if rs.get('name'):
                                run.font.name = rs['name']
                            if rs.get('size'):
                                run.font.size = rs['size']
                            run.font.bold, run.font.italic, run.font.underline = rs.get('bold'), rs.get('italic'), rs.get('underline')
                            if rs.get('color') and rs.get('color_type') == 1:
                                run.font.color.rgb = rs['color']
                            if rs.get('lang'):
                                run.font.language_id = rs['lang']
                        except:
                            pass
        except:
            shape.text = new_text
    
    def _update_cell_basic(self, cell, new_text: str):
        try:
            tf = cell.text_frame
            font_info = {}
            if tf.paragraphs and tf.paragraphs[0].runs:
                run = tf.paragraphs[0].runs[0]
                font_info = {'name': run.font.name, 'size': run.font.size, 'bold': run.font.bold}
            cell.text = new_text
            if font_info and tf.paragraphs and tf.paragraphs[0].runs:
                run = tf.paragraphs[0].runs[0]
                if font_info.get('name'):
                    run.font.name = font_info['name']
                if font_info.get('size'):
                    run.font.size = font_info['size']
                run.font.bold = font_info.get('bold')
        except:
            cell.text = new_text
    
    def _update_cell_advanced(self, cell, new_text: str):
        try:
            tf = cell.text_frame
            styles = []
            for p in tf.paragraphs:
                ps = {'alignment': p.alignment, 'runs': []}
                for run in p.runs:
                    try:
                        color_data = None
                        if hasattr(run.font, 'color') and run.font.color and hasattr(run.font.color, 'type') and run.font.color.type == 1:
                            color_data = run.font.color.rgb
                        ps['runs'].append({'name': run.font.name, 'size': run.font.size, 'bold': run.font.bold, 'italic': run.font.italic, 'color': color_data})
                    except:
                        pass
                styles.append(ps)
            
            cell.text = new_text
            
            for i, p in enumerate(tf.paragraphs):
                if i < len(styles):
                    s = styles[i]
                    try:
                        p.alignment = s['alignment']
                    except:
                        pass
                    if p.runs and s['runs']:
                        run, rs = p.runs[0], s['runs'][0]
                        try:
                            if rs.get('name'):
                                run.font.name = rs['name']
                            if rs.get('size'):
                                run.font.size = rs['size']
                            run.font.bold, run.font.italic = rs.get('bold'), rs.get('italic')
                            if rs.get('color'):
                                run.font.color.rgb = rs['color']
                        except:
                            pass
        except:
            cell.text = new_text
    
    def _run_preview(self):
        if self.processing:
            return messagebox.showwarning("処理中", "現在処理中です")
        data_path = filedialog.askopenfilename(title="編集済みファイルを選択", filetypes=[("Excel/TXT", "*.xlsx *.txt *.csv")], initialdir=self.config_manager.get('last_directory'))
        if not data_path:
            return
        ppt_path = filedialog.askopenfilename(title="PowerPointファイルを選択", filetypes=[("PowerPoint", "*.pptx")], initialdir=os.path.dirname(data_path))
        if not ppt_path:
            return
        
        def run():
            try:
                self._start_progress()
                self._update_output_safe(f"\n{ICONS['preview']} 差分プレビュー\n", clear=True)
                updates = self._load_excel_updates(data_path) if data_path.endswith('.xlsx') else self._load_txt_updates(data_path)
                if not updates:
                    return self._log("更新データなし", "warning")
                self._log(f"読み込み: {len(updates)}件")
                _, _, changes = self._update_ppt(ppt_path, updates, preview=True)
                if changes:
                    self._log(f"\n変更箇所: {len(changes)}件\n{'='*40}")
                    for i, c in enumerate(changes[:20], 1):
                        old_display = c['old'].replace('\n', '↵').replace('\r', '⏎').replace('\v', '⇓')
                        new_display = c['new'].replace('\n', '↵').replace('\r', '⏎').replace('\v', '⇓')
                        self._update_output_safe(f"[{i}] スライド{c['slide']} ID:{c['id']}\n  旧: {old_display}\n  新: {new_display}\n\n")
                    if len(changes) > 20:
                        self._log(f"... 他 {len(changes)-20}件")
                else:
                    self._log("変更箇所なし")
                self._update_status_safe(f"{ICONS['success']} プレビュー完了", COLOR_PALETTE["success"])
            except Exception as e:
                save_error_log(e, "_run_preview")
                self._log(f"エラー: {translate_error(e)}", "error")
            finally:
                self._stop_progress()
        
        threading.Thread(target=run, daemon=True).start()
    
    def _update_excel(self):
        self._run_update("excel")
    
    def _update_json(self):
        self._run_update("json")
    
    def _run_update(self, source: str):
        if self.processing:
            return messagebox.showwarning("処理中", "現在処理中です")
        
        limit = self.license_manager.get_update_limit()
        if limit:
            if not messagebox.askyesno("Free版制限", f"Free版では最初の{limit}スライドのみ更新されます。\n\n続行しますか？\n\nStandard版(¥2,980/年)で無制限更新！"):
                return self._show_license_dialog()
        
        if source == "excel":
            ftypes = [("Excel", "*.xlsx")]
        elif source == "json":
            ftypes = [("JSON", "*.json")]
        else:
            ftypes = [("テキスト", "*.txt *.csv")]
        
        data_path = filedialog.askopenfilename(title="編集済みファイルを選択", filetypes=ftypes, initialdir=self.config_manager.get('last_directory'))
        if not data_path:
            return
        ppt_path = filedialog.askopenfilename(title="PowerPointファイルを選択", filetypes=[("PowerPoint", "*.pptx")], initialdir=os.path.dirname(data_path))
        if not ppt_path:
            return
        self.config_manager.set('last_directory', os.path.dirname(data_path))
        
        def run():
            try:
                self._start_progress()
                self._update_output_safe(f"\n{ICONS['update']} 更新処理開始\n", clear=True)
                self._log(f"データ: {os.path.basename(data_path)}")
                self._log(f"PPT: {os.path.basename(ppt_path)}")
                
                self._create_backup(ppt_path)
                
                if source == "excel":
                    updates = self._load_excel_updates(data_path)
                elif source == "json":
                    updates = self._load_json_updates(data_path)
                else:
                    updates = self._load_txt_updates(data_path)
                
                if not updates:
                    return self._log("更新データなし", "warning")
                self._log(f"読み込み: {len(updates)}件")
                
                if limit:
                    self._log(f"⚠️ Free版: スライド1〜{limit}のみ更新", "warning")
                
                updated, skipped, _ = self._update_ppt(ppt_path, updates)
                
                if self.cancel_requested:
                    return self._log("キャンセルされました", "warning")
                
                def save():
                    out = filedialog.asksaveasfilename(title="保存先を選択", defaultextension=".pptx", filetypes=[("PowerPoint", "*.pptx")], initialdir=os.path.dirname(ppt_path), initialfile=os.path.splitext(os.path.basename(ppt_path))[0] + "_更新済み.pptx")
                    if out:
                        self.presentation.save(out)
                        self._log(f"保存完了: {os.path.basename(out)}", "success")
                        self._log(f"更新: {updated}件 / スキップ: {skipped}件")
                        self._update_status_safe(f"{ICONS['success']} 完了: {updated}件更新", COLOR_PALETTE["success"])
                        messagebox.showinfo("完了", f"更新: {updated}件\nスキップ: {skipped}件\n\n保存: {os.path.basename(out)}")
                    else:
                        self._log("保存キャンセル", "warning")
                
                self.root.after(0, save)
            except Exception as e:
                save_error_log(e, "_run_update")
                self._log(f"エラー: {translate_error(e)}", "error")
            finally:
                self._stop_progress()
        
        threading.Thread(target=run, daemon=True).start()
    
    def _update_batch(self):
        if self.processing:
            return messagebox.showwarning("処理中", "現在処理中です")
        folder = filedialog.askdirectory(title="フォルダを選択", initialdir=self.config_manager.get('last_directory'))
        if not folder:
            return
        self.config_manager.set('last_directory', folder)
        
        def run():
            try:
                self._start_progress()
                self._update_output_safe(f"\n{ICONS['batch']} バッチ更新開始: {folder}\n", clear=True)
                
                # Skip temp files (~$xxx.pptx)
                ppt_files = [f for f in Path(folder).glob("*.pptx") if not f.name.startswith("~$")]
                if not ppt_files:
                    return self._log("PPTXファイルなし", "warning")
                
                self._log(f"発見: {len(ppt_files)}件")
                total_updated, total_files = 0, 0
                
                for ppt in ppt_files:
                    if self.cancel_requested:
                        break
                    
                    # 対応する更新ファイルを探す
                    base = str(ppt.with_suffix(''))
                    data_path = None
                    for ext in ['_抽出.xlsx', '.xlsx', '_抽出.txt', '.txt', '_抽出.csv', '.csv']:
                        candidate = base + ext
                        if os.path.exists(candidate):
                            data_path = candidate
                            break
                    
                    if not data_path:
                        self._log(f"スキップ: {ppt.name} (更新ファイルなし)")
                        continue
                    
                    self._log(f"処理: {ppt.name}")
                    self._create_backup(str(ppt))
                    
                    updates = self._load_excel_updates(data_path) if data_path.endswith('.xlsx') else self._load_txt_updates(data_path)
                    if updates:
                        updated, _, _ = self._update_ppt(str(ppt), updates)
                        out = base + "_更新済み.pptx"
                        self.presentation.save(out)
                        total_updated += updated
                        total_files += 1
                        self._log(f"  → {updated}件更新")
                
                self._log(f"完了: {total_files}ファイル, {total_updated}件更新", "success")
                self._update_status_safe(f"{ICONS['success']} バッチ完了", COLOR_PALETTE["success"])
            except Exception as e:
                save_error_log(e, "_update_batch")
                self._log(f"エラー: {translate_error(e)}", "error")
            finally:
                self._stop_progress()
        
        threading.Thread(target=run, daemon=True).start()
    
    # === Dialogs ===
    def _show_upgrade_prompt(self):
        tier = self.license_manager.get_tier_info()
        if tier['name'] == 'Free':
            msg = t('upgrade_to_standard')
        else:
            msg = t('upgrade_to_pro')
        
        if messagebox.askyesno(t('upgrade_title'), msg):
            webbrowser.open(SUPPORT_LINKS["purchase"])
    
    def _show_license_dialog(self):
        dialog = tk.Toplevel(self.root)
        dialog.title(t('license_title'))
        self._center_dialog(dialog, 520, 700)
        dialog.resizable(False, False)
        dialog.transient(self.root)
        dialog.grab_set()
        
        frame = ttk.Frame(dialog, padding=SPACING["xl"])
        frame.pack(fill='both', expand=True)
        
        tier = self.license_manager.get_tier_info()
        license_type = self.license_manager.get_license_type()
        license_info = self.license_manager.license_info
        
        ttk.Label(frame, text=t('license_current'), font=FONTS["heading"]).pack(anchor='w')
        status_color = {"free": COLOR_PALETTE["free_badge"], "standard": COLOR_PALETTE["standard_badge"], "professional": COLOR_PALETTE["pro_badge"]}.get(license_type, COLOR_PALETTE["free_badge"])
        ttk.Label(frame, text=f"{tier['badge_text']} ({tier['name']})", font=FONTS["body_medium"], foreground=status_color).pack(anchor='w', pady=(SPACING["sm"], 0))
        
        # 登録キーと有効期限を表示
        if license_type != 'free' and license_info.get('key'):
            key_frame = ttk.Frame(frame)
            key_frame.pack(anchor='w', fill='x', pady=(SPACING["sm"], 0))
            
            # キー（マスク表示）
            key = license_info.get('key', '')
            masked_key = key[:8] + '****' + key[-4:] if len(key) > 12 else key
            ttk.Label(key_frame, text=f"{t('license_key_label')}: ", font=FONTS["caption"]).pack(side='left')
            ttk.Label(key_frame, text=masked_key, font=FONTS["code"], foreground=COLOR_PALETTE["text_muted"]).pack(side='left')
            
            # 有効期限
            expires = license_info.get('expires')
            if expires:
                exp_frame = ttk.Frame(frame)
                exp_frame.pack(anchor='w', fill='x', pady=(2, 0))
                ttk.Label(exp_frame, text=f"{t('license_expires_label')}: ", font=FONTS["caption"]).pack(side='left')
                
                # 残り日数を計算
                try:
                    exp_date = datetime.strptime(expires, '%Y-%m-%d')
                    days_left = (exp_date - datetime.now()).days
                    if days_left < 0:
                        exp_text = f"{expires} ({t('license_expired')})"
                        exp_color = COLOR_PALETTE["error"]
                    elif days_left <= 30:
                        exp_text = f"{expires} ({t('license_days_left', days_left)})"
                        exp_color = COLOR_PALETTE["warning"]
                    else:
                        exp_text = f"{expires} ({t('license_days_left', days_left)})"
                        exp_color = COLOR_PALETTE["success"]
                except:
                    exp_text = expires
                    exp_color = COLOR_PALETTE["text_muted"]
                
                ttk.Label(exp_frame, text=exp_text, font=FONTS["code"], foreground=exp_color).pack(side='left')
            else:
                # 永続ライセンス
                exp_frame = ttk.Frame(frame)
                exp_frame.pack(anchor='w', fill='x', pady=(2, 0))
                ttk.Label(exp_frame, text=f"{t('license_expires_label')}: ", font=FONTS["caption"]).pack(side='left')
                ttk.Label(exp_frame, text=t('license_permanent'), font=FONTS["code"], foreground=COLOR_PALETTE["success"]).pack(side='left')
        
        ttk.Separator(frame, orient='horizontal').pack(fill='x', pady=SPACING["md"])
        
        ttk.Label(frame, text=t('license_compare'), font=FONTS["body_medium"]).pack(anchor='w')
        
        plans = [
            ("Free (¥0)", t('license_free_features')),
            ("Standard (¥2,980/year)", t('license_standard_features')),
            ("Pro (¥5,980/year)", t('license_pro_features')),
        ]
        for name, features in plans:
            ttk.Label(frame, text=name, font=FONTS["body_medium"]).pack(anchor='w', pady=(SPACING["sm"], 0))
            for f in features:
                ttk.Label(frame, text=f"  • {f}", font=FONTS["caption"]).pack(anchor='w')
        
        ttk.Separator(frame, orient='horizontal').pack(fill='x', pady=SPACING["md"])
        
        ttk.Label(frame, text=t('license_enter_key'), font=FONTS["body"]).pack(anchor='w', pady=(SPACING["sm"], SPACING["sm"]))
        key_var = tk.StringVar()
        ttk.Entry(frame, textvariable=key_var, width=45, font=FONTS["code"]).pack(fill='x', pady=(0, SPACING["md"]))
        
        def activate():
            ok, msg = self.license_manager.activate_license(key_var.get().strip())
            if ok:
                messagebox.showinfo(t('dialog_complete'), msg)
                dialog.destroy()
                self._refresh_ui()
            else:
                messagebox.showerror(t('dialog_error'), msg)
        
        def deactivate():
            if messagebox.askyesno(t('dialog_confirm'), t('license_confirm_deactivate')):
                self.license_manager.deactivate_license()
                messagebox.showinfo(t('dialog_complete'), t('license_deactivated'))
                dialog.destroy()
                self._refresh_ui()
        
        btns = ttk.Frame(frame)
        btns.pack(fill='x')
        if license_type != 'free':
            ttk.Button(btns, text=t('btn_deactivate'), command=deactivate).pack(side='left')
        ttk.Button(btns, text=t('btn_activate'), command=activate).pack(side='left', padx=(SPACING["sm"], 0))
        ttk.Button(btns, text=t('btn_purchase'), command=lambda: webbrowser.open(SUPPORT_LINKS["purchase"])).pack(side='left', padx=(SPACING["sm"], 0))
        ttk.Button(btns, text=t('btn_close'), command=dialog.destroy).pack(side='right')
    
    def _show_welcome_dialog(self):
        dialog = tk.Toplevel(self.root)
        dialog.title(t('welcome_title'))
        self._center_dialog(dialog, 500, 400)
        dialog.resizable(False, False)
        dialog.transient(self.root)
        dialog.grab_set()
        
        frame = ttk.Frame(dialog, padding=SPACING["xl"])
        frame.pack(fill='both', expand=True)
        
        ttk.Label(frame, text=f"{ICONS['powerpoint']} {t('welcome_title')}", font=FONTS["heading"]).pack(pady=(0, SPACING["lg"]))
        
        steps = [
            (ICONS["extract"], t('welcome_step1_title'), t('welcome_step1_desc')),
            (ICONS["excel"], t('welcome_step2_title'), t('welcome_step2_desc')),
            (ICONS["update"], t('welcome_step3_title'), t('welcome_step3_desc')),
        ]
        for icon, title, desc in steps:
            ttk.Label(frame, text=f"{icon} {title}", font=FONTS["body_medium"]).pack(anchor='w')
            ttk.Label(frame, text=f"    {desc}", font=FONTS["caption"], foreground=COLOR_PALETTE["text_secondary"]).pack(anchor='w', pady=(0, SPACING["sm"]))
        
        ttk.Label(frame, text=f"\n{t('welcome_tip_title')}", font=FONTS["body_medium"]).pack(anchor='w')
        ttk.Label(frame, text=t('welcome_tip1'), font=FONTS["caption"]).pack(anchor='w')
        ttk.Label(frame, text=t('welcome_tip2'), font=FONTS["caption"]).pack(anchor='w')
        
        ttk.Button(frame, text=t('btn_start'), command=dialog.destroy).pack(pady=(SPACING["xl"], 0))
    
    def _show_about(self):
        tier = self.license_manager.get_tier_info()
        messagebox.showinfo(t('about_title'), t('about_text', APP_NAME, APP_VERSION, tier['name']))
    
    def _refresh_ui(self):
        tier = self.license_manager.get_tier_info()
        self.root.title(f"{APP_NAME} {tier['name']} - {t('app_subtitle')} v{APP_VERSION}")
        
        # メニューを再構築
        self._create_menu()
        
        # スタイルを再適用
        self._setup_styles()
        
        # レイアウト全体を再構築
        self._create_layout()
        
        if self.current_mode == "extract":
            self._switch_extract()
        else:
            self._switch_update()
        
        self._show_welcome()
    
    def _on_closing(self):
        if self.processing and not messagebox.askokcancel(t('dialog_confirm'), t('dialog_processing')):
            return
        self.config_manager.set('window_geometry', self.root.geometry())
        self.root.destroy()


def main():
    try:
        import ctypes
        ctypes.windll.shcore.SetProcessDpiAwareness(1)
    except:
        pass
    root = tk.Tk()
    InsightSlidesApp(root)
    root.mainloop()


if __name__ == "__main__":
    main()
